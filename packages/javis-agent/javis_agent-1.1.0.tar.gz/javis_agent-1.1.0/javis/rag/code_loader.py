"""
Code loader for RAG system
Load and chunk codebase for embedding
"""
import os
import time
from typing import List, Dict, Optional
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from javis.modules.config import setup_logging

logger = setup_logging()

# Optional imports for document parsing
try:
    import fitz  # PyMuPDF - much faster than pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyMuPDF not installed. PDF files will be skipped. Install with: pip install PyMuPDF")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed. DOCX files will be skipped. Install with: pip install python-docx")

try:
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("pandas/openpyxl not installed. Excel files will be skipped. Install with: pip install pandas openpyxl")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PowerPoint files will be skipped. Install with: pip install python-pptx")

class CodeLoader:
    """Load and prepare code documents for RAG"""
    
    def __init__(self, 
                 chunk_size: int = 1000, 
                 chunk_overlap: int = 200,
                 extensions: tuple = ('.cpp', '.h', '.hpp', '.py', '.java', '.js', '.ts', '.md', '.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.pptm', '.txt'),
                 adaptive_chunking: bool = True):
        """
        Initialize code loader
        
        Args:
            chunk_size: Default size of each chunk
            chunk_overlap: Overlap between chunks
            extensions: File extensions to load
            adaptive_chunking: Auto-adjust chunk size based on file size
        """
        self.default_chunk_size = chunk_size
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.extensions = extensions
        self.adaptive_chunking = adaptive_chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""],
            length_function=len,
        )
    
    def _calculate_optimal_chunk_size(self, file_size: int) -> int:
        """
        Calculate optimal chunk size based on file size (tuned for higher search granularity)
        
        Args:
            file_size: Size of file in characters
        Returns:
            Optimal chunk size
        """
        # Small files (< 10KB): 500 chars
        # if file_size < 10_000:
        #     return 500
        # # Medium files (10KB - 100KB): 1000 chars
        # elif file_size < 100_000:
        #     return 1000
        # # Large files (100KB - 1MB): 2000 chars
        # elif file_size < 1_000_000:
        #     return 2000
        # # Very large files (1MB - 5MB): 3000 chars
        # elif file_size < 5_000_000:
        #     return 3000
        # # Huge files (> 5MB): 5000 chars
        # else:
        #     return 5000
        return 1000
    
    def _detect_section_title(self, chunk: Document) -> None:
        """
        Detect section title from chunk content (for markdown/docs)
        Modifies chunk.metadata in-place
        
        Args:
            chunk: Document chunk to analyze
        """
        lines = chunk.page_content.split('\n')
        
        # Look for markdown headers (# Header)
        for line in lines[:5]:  # Check first 5 lines only
            line = line.strip()
            if line.startswith('#'):
                # Extract header text
                header_text = line.lstrip('#').strip()
                if header_text:
                    chunk.metadata['section_title'] = header_text
                    return
            
            # Also check for underlined headers (===, ---)
            if len(lines) > 1:
                next_line = lines[1].strip()
                if next_line and all(c in '=-' for c in next_line):
                    chunk.metadata['section_title'] = line
                    return
    
    def _extract_text_from_pdf(self, file_path: str) -> tuple:
        """
        Extract text from PDF file using PyMuPDF (5-10x faster than pdfplumber)
        Returns (full_text, page_metadata_list)
        """
        if not PDF_AVAILABLE:
            logger.warning(f"Skipping PDF {file_path}: PyMuPDF not installed")
            return "", []
        
        try:
            text_parts = []
            page_metadata = []
            # PyMuPDF is much faster than pdfplumber
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()  # Fast plain text extraction
                if text and text.strip():
                    text_parts.append(text)
                    page_metadata.append({
                        'page_number': page_num + 1,
                        'page_text_length': len(text),
                        'page_start_pos': sum(len(t) for t in text_parts[:-1]) + len(text_parts) - 1
                    })
            
            doc.close()
            full_text = "\n\n".join(text_parts)
            return full_text, page_metadata
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return "", []
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file (optimized)"""
        if not DOCX_AVAILABLE:
            logger.warning(f"Skipping DOCX {file_path}: python-docx not installed")
            return ""
        
        try:
            doc = DocxDocument(file_path)
            # Fast extraction - just paragraphs, no complex formatting
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file (optimized with row limits)"""
        if not EXCEL_AVAILABLE:
            logger.warning(f"Skipping Excel {file_path}: pandas/openpyxl not installed")
            return ""
        
        try:
            text_parts = []
            # Read with optimizations
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                # Limit rows to prevent memory issues with huge sheets
                # Read only first 10000 rows for performance
                df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=10000)
                
                if df.empty:
                    continue
                
                text_parts.append(f"## Sheet: {sheet_name}\n")
                # Convert to CSV-like format (faster than to_string)
                text_parts.append(df.to_csv(index=False, sep='\t'))
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {e}")
            return ""
    
    def _extract_text_from_pptx(self, file_path: str) -> tuple:
        """
        Extract text from PowerPoint file (optimized)
        Returns (full_text, slide_metadata_list)
        """
        if not PPTX_AVAILABLE:
            logger.warning(f"Skipping PowerPoint {file_path}: python-pptx not installed")
            return "", []
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            slide_metadata = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # Fast extraction - just text content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                    
                    # Extract from tables quickly
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_texts.append(" | ".join(row_text))
                
                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    text_parts.append(f"## Slide {slide_num}\n" + slide_content)
                    slide_metadata.append({
                        'slide_number': slide_num,
                        'slide_text_length': len(slide_content),
                        'slide_start_pos': sum(len(t) for t in text_parts[:-1]) + len(text_parts) - 1
                    })
            
            full_text = "\n\n".join(text_parts)
            return full_text, slide_metadata
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return "", []
    
    def _process_single_file(self, file_path: str, directory: str) -> Optional[Document]:
        """
        Process a single file and return Document object with enhanced metadata
        
        Args:
            file_path: Path to file
            directory: Base directory for relative path calculation
            
        Returns:
            Document object or None if failed
        """
        file_ext = Path(file_path).suffix.lower()
        file_name = os.path.basename(file_path)
        
        try:
            start_time = time.time()
            
            # Extract content based on file type
            # For PDF and PowerPoint, also extract page/slide metadata
            page_metadata = []
            if file_ext == '.pdf':
                content, page_metadata = self._extract_text_from_pdf(file_path)
            elif file_ext == '.docx':
                content = self._extract_text_from_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                content = self._extract_text_from_excel(file_path)
            elif file_ext in ['.pptx', '.pptm']:
                content, page_metadata = self._extract_text_from_pptx(file_path)
            else:
                # Plain text files (code, markdown, txt, etc.)
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            # Skip empty content
            if not content or not content.strip():
                logger.warning(f"Skipping empty file: {file_path}")
                return None
            
            # Create metadata with enhanced citation information
            relative_path = os.path.relpath(file_path, directory)
            metadata = {
                'source': file_path,
                'relative_path': relative_path,
                'file_name': file_name,
                'file_type': file_ext,
                'file_size': len(content),
                'page_metadata': page_metadata  # Page/slide info for PDF/PowerPoint
            }
            
            elapsed = time.time() - start_time
            logger.debug(f"Processed {file_name} ({len(content)} chars) in {elapsed:.2f}s")
            
            # Create document
            return Document(page_content=content, metadata=metadata)
            
        except Exception as e:
            logger.warning(f"Error reading {file_path}: {e}")
            return None
        
    def load_codebase(self, directory: str, max_workers: int = 4) -> List[Document]:
        """
        Load all code files from directory (parallel processing for better performance)
        
        Args:
            directory: Root directory of codebase
            max_workers: Number of parallel workers (default: 4)
            
        Returns:
            List of Document objects
        """
    def load_codebase(self, directory: str, max_workers: int = 4) -> List[Document]:
        """
        Load all code files from directory (parallel processing for better performance)
        
        Args:
            directory: Root directory of codebase
            max_workers: Number of parallel workers (default: 4)
            
        Returns:
            List of Document objects
        """
        documents = []
        directory_path = Path(directory)
        
        if not directory_path.exists():
            logger.error(f"Directory does not exist: {directory}")
            return documents
        
        logger.info(f"Loading codebase from: {directory}")
        
        # Collect all file paths first
        file_paths = []
        for root, dirs, files in os.walk(directory):
            # Skip common ignore directories
            dirs[:] = [d for d in dirs if d not in ['.git', '__pycache__', 'node_modules', 'build', 'dist', '.venv', 'venv','uts','tests','sldd','sample-apps','doc']]
            
            for file in files:
                if file.endswith(self.extensions):
                    file_path = os.path.join(root, file)
                    file_paths.append(file_path)
        
        if not file_paths:
            logger.warning(f"No files found with extensions {self.extensions}")
            return documents
        
        logger.info(f"Found {len(file_paths)} files to process")
        start_time = time.time()
        
        # Process files in parallel using ThreadPoolExecutor
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all tasks
            future_to_file = {
                executor.submit(self._process_single_file, file_path, directory): file_path
                for file_path in file_paths
            }
            
            # Collect results as they complete
            processed = 0
            for future in as_completed(future_to_file):
                processed += 1
                doc = future.result()
                if doc is not None:
                    documents.append(doc)
                
                # Progress logging every 10 files or at the end
                if processed % 10 == 0 or processed == len(file_paths):
                    logger.info(f"Progress: {processed}/{len(file_paths)} files processed ({len(documents)} valid)")
        
        elapsed = time.time() - start_time
        logger.info(f"✅ Loaded {len(documents)} documents from {len(file_paths)} files in {elapsed:.2f}s ({len(file_paths)/elapsed:.1f} files/sec)")
        
        return documents
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """
        Split documents into chunks with adaptive chunk size
        
        Args:
            documents: List of documents to chunk
            
        Returns:
            List of chunked documents
        """
        logger.info(f"Chunking {len(documents)} documents...")
        
        all_chunks = []
        
        for doc in documents:
            # Determine chunk size based on document size
            if self.adaptive_chunking:
                file_size = doc.metadata.get('file_size', len(doc.page_content))
                optimal_chunk_size = self._calculate_optimal_chunk_size(file_size)
                
                # Log if chunk size is different from default
                if optimal_chunk_size != self.default_chunk_size:
                    logger.info(
                        f"File '{doc.metadata.get('file_name', 'unknown')}' "
                        f"({file_size:,} chars) -> chunk_size={optimal_chunk_size}"
                    )
                
                # Create custom splitter for this document
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=optimal_chunk_size,
                    chunk_overlap=int(optimal_chunk_size * 0.2),  # 20% overlap
                    separators=["\n\n", "\n", " ", ""],
                    length_function=len,
                )
            else:
                splitter = self.text_splitter
            
            # Split document
            doc_chunks = splitter.split_documents([doc])
            
            # Add chunk index and position metadata for citation
            page_metadata = doc.metadata.get('page_metadata', [])
            for i, chunk in enumerate(doc_chunks):
                chunk.metadata['chunk_index'] = i
                chunk.metadata['total_chunks'] = len(doc_chunks)
                
                # Find which page/slide this chunk belongs to (for PDF/PowerPoint)
                if page_metadata:
                    chunk_start = sum(len(c.page_content) for c in doc_chunks[:i])
                    for page_info in page_metadata:
                        if chunk_start >= page_info.get('page_start_pos', 0) or chunk_start >= page_info.get('slide_start_pos', 0):
                            if 'page_number' in page_info:
                                chunk.metadata['page_number'] = page_info['page_number']
                            if 'slide_number' in page_info:
                                chunk.metadata['slide_number'] = page_info['slide_number']
                
                # Add section title if detectable (for markdown/code files)
                self._detect_section_title(chunk)
            
            all_chunks.extend(doc_chunks)
        
        logger.info(f"Created {len(all_chunks)} chunks")
        return all_chunks
    
    def load_and_chunk(self, directory: str) -> List[Document]:
        """
        Load codebase and chunk in one step
        
        Args:
            directory: Root directory of codebase
            
        Returns:
            List of chunked documents
        """
        documents = self.load_codebase(directory)
        if not documents:
            return []
        return self.chunk_documents(documents)
