"""File processor module for converting multimodal files to markdown using markitdown.

This module provides unified file processing for various file types including:
- Office documents (DOCX, XLSX, PPTX)
- Documents (PDF, TXT, HTML)
- Images (JPG, PNG with OCR)
- Data files (CSV, JSON)

Files are converted to markdown format for optimal LLM consumption.
"""

import contextlib
import hashlib
import json
import tempfile
import threading
import time
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING, Any, TypeVar

import requests

if TYPE_CHECKING:
    from holodeck.models.config import ExecutionConfig

from holodeck.lib.logging_config import get_logger
from holodeck.lib.logging_utils import log_exception, log_retry
from holodeck.models.test_case import FileInput
from holodeck.models.test_result import ProcessedFileInput

logger = get_logger(__name__)

T = TypeVar("T")


@dataclass
class SourceFile:
    """Source file to be ingested into vector store.

    Represents a file during the ingestion process with metadata and content.

    Attributes:
        path: Absolute file path
        content: File content converted to markdown (populated by FileProcessor)
        mtime: File modification time (Unix timestamp)
        size_bytes: File size in bytes
        file_type: File extension (.txt, .md, .pdf, .csv, .json, etc.)
        chunks: Text chunks after splitting (populated by TextChunker)
    """

    path: Path
    content: str = ""
    mtime: float = 0.0
    size_bytes: int = 0
    file_type: str = ""
    chunks: list[str] = field(default_factory=list)


class FileProcessor:
    """Process files with markitdown for multimodal test inputs."""

    def __init__(
        self,
        cache_dir: str | None = None,
        download_timeout_ms: int = 30000,
        max_retries: int = 3,
        processing_timeout_ms: int = 30000,
    ) -> None:
        """Initialize file processor.

        Args:
            cache_dir: Directory for caching remote files. Defaults to .holodeck/cache/
            download_timeout_ms: Timeout for file downloads in milliseconds
            max_retries: Maximum number of retry attempts for downloads
            processing_timeout_ms: Timeout for file processing in milliseconds.
                Defaults to 30000ms.
        """
        self.cache_dir = Path(cache_dir or ".holodeck/cache")
        self.cache_dir.mkdir(parents=True, exist_ok=True)

        self.download_timeout_ms = download_timeout_ms
        self.max_retries = max_retries
        self.processing_timeout_ms = processing_timeout_ms
        self.md: Any = None  # Initialize lazily

        logger.debug(
            f"FileProcessor initialized: cache_dir={self.cache_dir}, "
            f"download_timeout={download_timeout_ms}ms, "
            f"processing_timeout={processing_timeout_ms}ms, max_retries={max_retries}"
        )

        try:
            from markitdown import MarkItDown  # noqa: F401
        except ImportError as e:
            logger.error("markitdown package not found", exc_info=True)
            raise ImportError(
                "markitdown is required for file processing. "
                "Install with: pip install 'markitdown[all]'"
            ) from e

    @classmethod
    def from_execution_config(
        cls,
        config: "ExecutionConfig",
        cache_dir: str | None = None,
        max_retries: int = 3,
    ) -> "FileProcessor":
        """Create FileProcessor from ExecutionConfig.

        Factory method that handles conversion from ExecutionConfig's
        seconds-based timeouts to FileProcessor's milliseconds-based timeouts.

        Args:
            config: ExecutionConfig with timeout settings in seconds
            cache_dir: Override cache directory (defaults to config.cache_dir)
            max_retries: Maximum retry attempts for downloads

        Returns:
            Configured FileProcessor instance
        """
        # Convert seconds to milliseconds, using defaults if not specified
        download_timeout_ms = (config.download_timeout or 30) * 1000
        processing_timeout_ms = (config.file_timeout or 30) * 1000

        return cls(
            cache_dir=cache_dir or config.cache_dir or ".holodeck/cache",
            download_timeout_ms=download_timeout_ms,
            processing_timeout_ms=processing_timeout_ms,
            max_retries=max_retries,
        )

    def _get_markitdown(self) -> Any:
        """Get or create MarkItDown instance."""
        if self.md is None:
            from markitdown import MarkItDown

            self.md = MarkItDown()
        return self.md

    def _with_timeout(
        self, func: Callable[..., T], timeout_ms: int, *args: Any, **kwargs: Any
    ) -> T:
        """Execute function with timeout support using threading.

        Provides cross-platform timeout support using threading.Timer.
        The timeout works by setting a timer that can interrupt the operation
        (though this is done via exception raising in the calling thread).

        Args:
            func: Function to execute
            timeout_ms: Timeout in milliseconds
            *args: Positional arguments for func
            **kwargs: Keyword arguments for func

        Returns:
            Result of function execution

        Raises:
            TimeoutError: If function execution exceeds timeout
        """
        timeout_sec = timeout_ms / 1000.0
        result_holder: list[T] = []
        exception_holder: list[BaseException] = []

        def target() -> None:
            try:
                result: T = func(*args, **kwargs)
                result_holder.append(result)
            except BaseException as e:
                exception_holder.append(e)

        thread = threading.Thread(target=target, daemon=True)
        thread.start()
        thread.join(timeout=timeout_sec)

        # Check if thread is still alive (timeout occurred)
        if thread.is_alive():
            raise TimeoutError(f"Processing exceeded {timeout_ms}ms timeout limit")

        # Check if an exception occurred in the thread
        if exception_holder:
            raise exception_holder[0]

        # Return result if available
        if result_holder:
            return result_holder[0]

        raise RuntimeError("No result from processing operation")

    def _get_cache_key(self, url_or_path: str) -> str:
        """Generate cache key using MD5 hash of URL or path.

        Args:
            url_or_path: File URL or path

        Returns:
            MD5 hash of the input string
        """
        return hashlib.md5(url_or_path.encode(), usedforsecurity=False).hexdigest()

    def _load_from_cache(self, cache_key: str) -> dict[str, Any] | None:
        """Load processed file from cache.

        Args:
            cache_key: Cache key (MD5 hash)

        Returns:
            Cached data dict or None if not found
        """
        cache_file = self.cache_dir / f"{cache_key}.json"
        if cache_file.exists():
            try:
                with open(cache_file) as f:
                    data = json.load(f)
                    if isinstance(data, dict):
                        logger.debug(f"Cache hit: {cache_key}")
                        return data
                    return None
            except Exception as e:
                logger.warning(f"Failed to load cache for {cache_key}: {e}")
                return None
        logger.debug(f"Cache miss: {cache_key}")
        return None

    def _save_to_cache(
        self,
        cache_key: str,
        markdown_content: str,
        metadata: dict,
        processing_time_ms: int,
    ) -> None:
        """Save processed file to cache.

        Args:
            cache_key: Cache key (MD5 hash)
            markdown_content: Converted markdown content
            metadata: File metadata
            processing_time_ms: Processing time in milliseconds
        """
        cache_file = self.cache_dir / f"{cache_key}.json"
        data = {
            "markdown_content": markdown_content,
            "metadata": metadata,
            "processing_time_ms": processing_time_ms,
        }
        try:
            with open(cache_file, "w") as f:
                json.dump(data, f)
            logger.debug(f"File cached: {cache_key} ({len(markdown_content)} bytes)")
        except Exception as e:
            logger.warning(f"Failed to save cache for {cache_key}: {e}")

    def _download_file(self, url: str) -> bytes | None:
        """Download file from URL with retry logic.

        Implements exponential backoff: 1s, 2s, 4s for retries.

        Args:
            url: Remote file URL

        Returns:
            File content bytes or None if download fails
        """
        timeout_sec = self.download_timeout_ms / 1000.0

        logger.debug(f"Downloading file from URL: {url}")

        for attempt in range(self.max_retries):
            try:
                response = requests.get(url, timeout=timeout_sec)
                response.raise_for_status()
                content: bytes = response.content
                size_bytes = len(content)
                logger.debug(
                    f"File downloaded successfully: {url} ({size_bytes} bytes)"
                )
                return content
            except Exception as e:
                if attempt == self.max_retries - 1:
                    logger.error(
                        f"Failed to download file after {self.max_retries} "
                        f"attempts: {url}"
                    )
                    return None

                # Exponential backoff: 1s, 2s, 4s
                backoff_sec = 2**attempt
                log_retry(
                    logger,
                    f"Download {url}",
                    attempt=attempt + 1,
                    max_attempts=self.max_retries,
                    delay=backoff_sec,
                    error=e,
                )
                time.sleep(backoff_sec)

        return None

    def process_file(self, file_input: FileInput) -> ProcessedFileInput:
        """Process a single file input to markdown.

        Args:
            file_input: File input configuration with path or URL

        Returns:
            ProcessedFileInput with markdown content and metadata
        """
        start_time = time.time()
        file_location = file_input.url or file_input.path or "unknown"

        logger.debug(
            f"Processing file: {file_location} (type={file_input.type}, "
            f"cache={'enabled' if file_input.cache else 'disabled'})"
        )

        try:
            # Determine if file is local or remote
            if file_input.url:
                return self._process_remote_file(file_input, start_time)
            else:
                return self._process_local_file(file_input, start_time)

        except Exception as e:
            elapsed_ms = int((time.time() - start_time) * 1000)
            log_exception(logger, f"File processing failed: {file_location}", e)
            # Create detailed error message with context
            error_msg = self._format_error_message(e, file_location)
            return ProcessedFileInput(
                original=file_input,
                markdown_content="",
                metadata={},
                processing_time_ms=elapsed_ms,
                cached_path=None,
                error=error_msg,
            )

    def _format_error_message(self, error: Exception, file_location: str) -> str:
        """Format error message with context and file location.

        Args:
            error: The exception that occurred
            file_location: Location of file being processed (path or URL)

        Returns:
            Formatted error message with context
        """
        error_type = type(error).__name__
        error_str = str(error)

        # Build error message with file context
        if file_location and file_location != "unknown":
            return f"{error_type}: {error_str} (while processing: {file_location})"
        return f"{error_type}: {error_str}"

    def _preprocess_pdf_pages(self, file_path: Path, pages: list[int]) -> Path:
        """Extract specific pages from PDF into temporary file.

        Args:
            file_path: Path to original PDF file
            pages: List of page numbers to extract (0-indexed)

        Returns:
            Path to temporary PDF file with extracted pages

        Raises:
            ValueError: If page numbers are invalid or out of range
        """
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError as e:
            raise ImportError(
                "pypdf is required for PDF page extraction. "
                "Install with: pip install 'markitdown[all]'"
            ) from e

        logger.debug(f"Extracting pages {pages} from PDF: {file_path}")

        try:
            reader = PdfReader(str(file_path))
            writer = PdfWriter()
            total_pages = len(reader.pages)

            # Validate page numbers
            for page_num in pages:
                if page_num < 0 or page_num >= total_pages:
                    raise ValueError(
                        f"Page {page_num} out of range (PDF has {total_pages} pages)"
                    )

            # Extract specified pages
            for page_num in pages:
                writer.add_page(reader.pages[page_num])

            # Create temporary file
            tmp = tempfile.NamedTemporaryFile(  # noqa: SIM115
                suffix=".pdf", delete=False
            )
            tmp_path = Path(tmp.name)
            writer.write(tmp)
            tmp.close()

            logger.debug(
                f"Extracted {len(pages)} pages from PDF to temp file: {tmp_path}"
            )
            return tmp_path

        except Exception as e:
            logger.error(f"PDF page extraction failed: {e}")
            raise

    def _preprocess_excel_sheet_range(
        self, file_path: Path, sheet: str | None, range_spec: str | None
    ) -> Path:
        """Extract specific sheet/range from Excel into temporary CSV file.

        Args:
            file_path: Path to original Excel file
            sheet: Sheet name to extract (None = first sheet)
            range_spec: Cell range to extract (e.g., "A1:E100")

        Returns:
            Path to temporary CSV file with extracted data

        Raises:
            ValueError: If sheet name or range is invalid
        """
        try:
            import openpyxl
            import pandas as pd
        except ImportError as e:
            raise ImportError(
                "openpyxl and pandas are required for Excel extraction. "
                "Install with: pip install 'markitdown[all]' pandas"
            ) from e

        logger.debug(
            f"Extracting Excel data: sheet={sheet}, range={range_spec} "
            f"from {file_path}"
        )

        try:
            # Load workbook
            wb = openpyxl.load_workbook(str(file_path), data_only=True)

            # Select sheet
            if sheet:
                if sheet not in wb.sheetnames:
                    raise ValueError(
                        f"Sheet '{sheet}' not found. "
                        f"Available sheets: {', '.join(wb.sheetnames)}"
                    )
                ws = wb[sheet]
            else:
                ws = wb.active

            # Extract data based on range
            if range_spec:
                # Parse range like "A1:E100"
                try:
                    data = ws[range_spec]
                    # Convert to list of lists
                    if isinstance(data, tuple):
                        # Single cell or range
                        if isinstance(data[0], tuple):
                            # Range of cells
                            rows = [[cell.value for cell in row] for row in data]
                        else:
                            # Single row
                            rows = [[cell.value for cell in data]]
                    else:
                        # Single cell
                        rows = [[data.value]]
                except Exception as e:
                    raise ValueError(f"Invalid range '{range_spec}': {e}") from e
            else:
                # Extract all data from sheet
                rows = [[cell.value for cell in row] for row in ws.iter_rows()]

            # Convert to DataFrame and save as CSV
            df = pd.DataFrame(rows)

            # Create temporary CSV file
            tmp = tempfile.NamedTemporaryFile(  # noqa: SIM115
                mode="w", suffix=".csv", delete=False, newline=""
            )
            tmp_path = Path(tmp.name)
            df.to_csv(tmp_path, index=False, header=False)

            logger.debug(
                f"Extracted Excel data ({len(rows)} rows) to temp CSV: {tmp_path}"
            )
            return tmp_path

        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise

    def _preprocess_powerpoint_slides(self, file_path: Path, pages: list[int]) -> Path:
        """Extract specific slides from PowerPoint into temporary file.

        Args:
            file_path: Path to original PowerPoint file
            pages: List of slide numbers to extract (0-indexed)

        Returns:
            Path to temporary PowerPoint file with extracted slides

        Raises:
            ValueError: If slide numbers are invalid or out of range
        """
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError(
                "python-pptx is required for PowerPoint slide extraction. "
                "Install with: pip install 'markitdown[all]'"
            ) from e

        logger.debug(f"Extracting slides {pages} from PowerPoint: {file_path}")

        try:
            # Load presentation
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)

            # Validate slide numbers
            for slide_num in pages:
                if slide_num < 0 or slide_num >= total_slides:
                    raise ValueError(
                        f"Slide {slide_num} out of range "
                        f"(presentation has {total_slides} slides)"
                    )

            # Create new presentation with selected slides
            new_prs = Presentation()
            # Copy slide dimensions (only if they exist)
            if prs.slide_width is not None:
                new_prs.slide_width = prs.slide_width
            if prs.slide_height is not None:
                new_prs.slide_height = prs.slide_height

            # Extract specified slides
            for slide_num in pages:
                slide = prs.slides[slide_num]
                # Copy slide layout and content
                new_slide_layout = new_prs.slide_layouts[0]
                new_slide = new_prs.slides.add_slide(new_slide_layout)

                # Copy all shapes from original slide
                for shape in slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

            # Create temporary file
            tmp = tempfile.NamedTemporaryFile(  # noqa: SIM115
                suffix=".pptx", delete=False
            )
            tmp_path = Path(tmp.name)
            new_prs.save(str(tmp_path))
            tmp.close()

            logger.debug(
                f"Extracted {len(pages)} slides from PowerPoint to temp file: "
                f"{tmp_path}"
            )
            return tmp_path

        except Exception as e:
            logger.error(f"PowerPoint slide extraction failed: {e}")
            raise

    def _preprocess_file(self, file_input: FileInput, file_path: Path) -> Path:
        """Main preprocessing dispatcher that routes to specific handlers.

        Checks if file needs preprocessing (pages/sheet/range extraction) and
        routes to appropriate handler. If no preprocessing needed, returns
        original file path.

        Args:
            file_input: File input configuration
            file_path: Path to file

        Returns:
            Path to preprocessed file (or original if no preprocessing)
        """
        # Check if preprocessing is needed
        has_pages = file_input.pages and len(file_input.pages) > 0
        has_sheet = file_input.sheet is not None
        has_range = file_input.range is not None

        # No preprocessing needed
        if not (has_pages or has_sheet or has_range):
            return file_path

        file_type = file_input.type.lower()

        # Route to appropriate preprocessor
        if has_pages and file_type == "pdf":
            # has_pages check ensures file_input.pages is not None
            pages = file_input.pages
            if pages is not None:
                return self._preprocess_pdf_pages(file_path, pages)
        elif has_pages and file_type == "powerpoint":
            # has_pages check ensures file_input.pages is not None
            pages = file_input.pages
            if pages is not None:
                return self._preprocess_powerpoint_slides(file_path, pages)
        elif (has_sheet or has_range) and file_type == "excel":
            return self._preprocess_excel_sheet_range(
                file_path, file_input.sheet, file_input.range
            )

        # No applicable preprocessing or check above failed
        return file_path

    def _process_local_file(
        self, file_input: FileInput, start_time: float
    ) -> ProcessedFileInput:
        """Process a local file.

        Args:
            file_input: File input with local path
            start_time: Processing start time

        Returns:
            ProcessedFileInput with markdown content
        """
        if not file_input.path:
            raise ValueError("Local file must have path specified")

        path = Path(file_input.path)

        # Validate file exists
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_input.path}")

        # Preprocess file if needed (page/sheet/range extraction)
        processed_path = self._preprocess_file(file_input, path)

        # Get file metadata
        file_size = path.stat().st_size
        metadata: dict = {
            "size_bytes": file_size,
            "path": str(path),
            "type": file_input.type,
        }

        # Add preprocessing metadata if preprocessing occurred
        if processed_path != path:
            metadata["preprocessed"] = True
            if file_input.pages:
                metadata["pages_extracted"] = file_input.pages
            if file_input.sheet:
                metadata["sheet_extracted"] = file_input.sheet
            if file_input.range:
                metadata["range_extracted"] = file_input.range

        # Warn if file is large
        size_mb = file_size / (1024 * 1024)
        if size_mb > 100:
            logger.warning(f"Large file detected: {file_input.path} ({size_mb:.2f}MB)")
            metadata["warning"] = f"Large file detected ({size_mb:.2f}MB)"

        # Convert file (use preprocessed path if available) with timeout
        logger.debug(f"Converting local file to markdown: {file_input.path}")
        md = self._get_markitdown()
        try:
            result = self._with_timeout(
                md.convert, self.processing_timeout_ms, str(processed_path)
            )
            markdown_content = result.text_content
        except TimeoutError:
            logger.error(f"File processing timeout: {file_input.path}")
            raise

        # Clean up temporary preprocessed file if created
        if processed_path != path:
            with contextlib.suppress(Exception):
                processed_path.unlink()

        elapsed_ms = int((time.time() - start_time) * 1000)

        logger.debug(
            f"Local file processed: {file_input.path} "
            f"({len(markdown_content)} bytes in {elapsed_ms}ms)"
        )

        return ProcessedFileInput(
            original=file_input,
            markdown_content=markdown_content,
            metadata=metadata,
            processing_time_ms=elapsed_ms,
            cached_path=None,
            error=None,
        )

    def _process_remote_file(
        self, file_input: FileInput, start_time: float
    ) -> ProcessedFileInput:
        """Process a remote file from URL.

        Args:
            file_input: File input with URL
            start_time: Processing start time

        Returns:
            ProcessedFileInput with markdown content
        """
        if not file_input.url:
            raise ValueError("Remote file must have URL specified")

        url = file_input.url

        # Check cache if enabled
        if file_input.cache is not False:
            cache_key = self._get_cache_key(url)
            cached_data = self._load_from_cache(cache_key)

            if cached_data:
                elapsed_ms = int((time.time() - start_time) * 1000)
                logger.debug(f"Using cached remote file: {url}")
                return ProcessedFileInput(
                    original=file_input,
                    markdown_content=cached_data["markdown_content"],
                    metadata=cached_data["metadata"],
                    cached_path=str(self.cache_dir / f"{cache_key}.json"),
                    processing_time_ms=elapsed_ms,
                    error=None,
                )

        # Download file
        file_content = self._download_file(url)
        if file_content is None:
            elapsed_ms = int((time.time() - start_time) * 1000)
            logger.error(f"Failed to download remote file: {url}")
            return ProcessedFileInput(
                original=file_input,
                markdown_content="",
                metadata={"url": url},
                processing_time_ms=elapsed_ms,
                cached_path=None,
                error="Failed to download file after max retries",
            )

        # Save to temporary file and process
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name

        try:
            # Preprocess file if needed (page/sheet/range extraction)
            path_obj = Path(tmp_path)
            processed_path = self._preprocess_file(file_input, path_obj)

            logger.debug(f"Converting remote file to markdown: {url}")
            md = self._get_markitdown()
            try:
                result = self._with_timeout(
                    md.convert, self.processing_timeout_ms, str(processed_path)
                )
                markdown_content = result.text_content
            except TimeoutError:
                logger.error(f"Remote file processing timeout: {url}")
                raise

            # Clean up temporary preprocessed file if created
            if processed_path != path_obj:
                with contextlib.suppress(Exception):
                    processed_path.unlink()

            # Get metadata
            metadata: dict = {
                "url": url,
                "size_bytes": len(file_content),
                "type": file_input.type,
            }

            # Add preprocessing metadata if preprocessing occurred
            if processed_path != path_obj:
                metadata["preprocessed"] = True
                if file_input.pages:
                    metadata["pages_extracted"] = file_input.pages
                if file_input.sheet:
                    metadata["sheet_extracted"] = file_input.sheet
                if file_input.range:
                    metadata["range_extracted"] = file_input.range

            elapsed_ms = int((time.time() - start_time) * 1000)
            cached_path = None

            # Cache if enabled
            if file_input.cache is not False:
                cache_key = self._get_cache_key(url)
                self._save_to_cache(cache_key, markdown_content, metadata, elapsed_ms)
                cached_path = str(self.cache_dir / f"{cache_key}.json")

            logger.debug(
                f"Remote file processed: {url} "
                f"({len(markdown_content)} bytes in {elapsed_ms}ms, "
                f"cached={cached_path is not None})"
            )

            return ProcessedFileInput(
                original=file_input,
                markdown_content=markdown_content,
                metadata=metadata,
                cached_path=cached_path,
                processing_time_ms=elapsed_ms,
                error=None,
            )

        finally:
            # Clean up temporary file
            with contextlib.suppress(Exception):
                Path(tmp_path).unlink()
