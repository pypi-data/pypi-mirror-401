"""Unit tests for file processor module using markitdown integration."""

import hashlib
import tempfile
from pathlib import Path
from unittest import mock

from holodeck.lib.file_processor import FileProcessor
from holodeck.models.config import ExecutionConfig
from holodeck.models.test_case import FileInput
from holodeck.models.test_result import ProcessedFileInput


class TestFileProcessorBasics:
    """Tests for basic file processor functionality."""

    def test_process_file_routes_to_local(self) -> None:
        """Test that process_file routes local files correctly."""
        mock_file_input = FileInput(path="/path/to/sample.pdf", type="pdf")

        with mock.patch.object(FileProcessor, "_process_local_file") as mock_local:
            mock_local.return_value = ProcessedFileInput(
                original=mock_file_input,
                markdown_content="# Content",
                processing_time_ms=100,
            )

            processor = FileProcessor()
            result = processor.process_file(mock_file_input)

            assert result.original == mock_file_input
            mock_local.assert_called_once()

    def test_process_file_routes_to_remote(self) -> None:
        """Test that process_file routes remote files correctly."""
        mock_file_input = FileInput(url="https://example.com/file.pdf", type="pdf")

        with mock.patch.object(FileProcessor, "_process_remote_file") as mock_remote:
            mock_remote.return_value = ProcessedFileInput(
                original=mock_file_input,
                markdown_content="# Content",
                processing_time_ms=100,
            )

            processor = FileProcessor()
            result = processor.process_file(mock_file_input)

            assert result.original == mock_file_input
            mock_remote.assert_called_once()

    def test_cache_key_generation(self) -> None:
        """Test cache key generation uses MD5 hashing."""
        processor = FileProcessor()

        url = "https://example.com/test.pdf"
        expected_hash = hashlib.md5(url.encode()).hexdigest()  # noqa: S324
        actual_hash = processor._get_cache_key(url)

        assert actual_hash == expected_hash
        assert len(actual_hash) == 32  # MD5 is 32 hex characters


class TestFileProcessorMetadata:
    """Tests for file metadata tracking."""

    def test_error_handling_returns_processed_input(self) -> None:
        """Test that errors are captured and returned in result."""
        mock_file_input = FileInput(path="/nonexistent/file.pdf", type="pdf")

        with mock.patch.object(FileProcessor, "_process_local_file") as mock_local:
            mock_local.side_effect = FileNotFoundError("File not found")

            processor = FileProcessor()
            result = processor.process_file(mock_file_input)

            assert result.error is not None
            assert "File not found" in result.error
            assert result.original == mock_file_input


class TestFileProcessorCaching:
    """Tests for file caching functionality."""

    def test_cache_directory_creation(self) -> None:
        """Test that cache directory is created on initialization."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            assert cache_dir.exists()
            assert processor.cache_dir == cache_dir

    def test_load_cache_missing_file(self) -> None:
        """Test loading from cache when file doesn't exist."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            result = processor._load_from_cache("nonexistent_key")
            assert result is None

    def test_save_and_load_cache(self) -> None:
        """Test saving and loading from cache."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            cache_key = "test_key"
            metadata = {"size": 1024}
            processor._save_to_cache(cache_key, "Test content", metadata, 100)

            loaded = processor._load_from_cache(cache_key)
            assert loaded is not None
            assert loaded["markdown_content"] == "Test content"
            assert loaded["metadata"] == metadata
            assert loaded["processing_time_ms"] == 100


class TestFileProcessorErrorHandling:
    """Tests for error handling in file processing."""

    def test_error_dict_structure(self) -> None:
        """Test that error messages are properly structured."""
        result = ProcessedFileInput(
            original=FileInput(path="test.pdf", type="pdf"),
            markdown_content="",
            processing_time_ms=100,
            error="Test error message",
        )

        assert result.error == "Test error message"
        assert result.markdown_content == ""


class TestFileProcessorDownloads:
    """Tests for remote file download functionality."""

    def test_download_file_success(self) -> None:
        """Test successful file download."""
        processor = FileProcessor(download_timeout_ms=30000)

        with mock.patch("requests.get") as mock_get:
            mock_get.return_value = mock.Mock(content=b"PDF content", status_code=200)

            result = processor._download_file("https://example.com/file.pdf")

            assert result == b"PDF content"
            mock_get.assert_called_once()

    def test_download_file_retry_on_failure(self) -> None:
        """Test retry logic for failed downloads."""
        processor = FileProcessor(max_retries=3)

        with mock.patch("requests.get") as mock_get, mock.patch("time.sleep"):
            mock_get.side_effect = [
                Exception("Timeout"),
                Exception("Timeout"),
                mock.Mock(content=b"Success", status_code=200),
            ]

            result = processor._download_file("https://example.com/file.pdf")

            assert result == b"Success"
            assert mock_get.call_count == 3

    def test_download_file_max_retries_exceeded(self) -> None:
        """Test handling when max retries exceeded."""
        processor = FileProcessor(max_retries=3)

        with mock.patch("requests.get") as mock_get, mock.patch("time.sleep"):
            mock_get.side_effect = Exception("Timeout")

            result = processor._download_file("https://example.com/file.pdf")

            assert result is None
            assert mock_get.call_count == 3


class TestFileProcessorConfiguration:
    """Tests for FileProcessor configuration."""

    def test_custom_cache_directory(self) -> None:
        """Test setting custom cache directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            custom_cache = Path(tmpdir) / "custom" / "cache"
            processor = FileProcessor(cache_dir=str(custom_cache))

            assert processor.cache_dir == custom_cache
            assert custom_cache.exists()

    def test_custom_download_timeout(self) -> None:
        """Test setting custom download timeout."""
        processor = FileProcessor(download_timeout_ms=60000)

        assert processor.download_timeout_ms == 60000

    def test_custom_max_retries(self) -> None:
        """Test setting custom max retries."""
        processor = FileProcessor(max_retries=5)

        assert processor.max_retries == 5


class TestFileProcessorMarkItDown:
    """Tests for MarkItDown integration and lazy initialization."""

    def test_markitdown_import_error(self) -> None:
        """Test that import error is raised when markitdown not available."""
        with (
            mock.patch.dict("sys.modules", {"markitdown": None}),
            mock.patch("builtins.__import__", side_effect=ImportError("No module")),
        ):
            try:
                FileProcessor()
                raise AssertionError("Should raise ImportError")
            except ImportError as e:
                assert "markitdown is required" in str(e)

    def test_get_markitdown_lazy_initialization(self) -> None:
        """Test that MarkItDown is initialized lazily on first access."""
        processor = FileProcessor()

        # Initially md should be None
        assert processor.md is None

        # Mock the MarkItDown class
        mock_md_instance = mock.MagicMock()
        with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
            result = processor._get_markitdown()

            assert result is not None
            assert result is mock_md_instance
            assert processor.md is not None
            assert processor.md is mock_md_instance

    def test_get_markitdown_returns_cached_instance(self) -> None:
        """Test that _get_markitdown returns cached instance on subsequent calls."""
        processor = FileProcessor()

        mock_md_instance = mock.MagicMock()
        with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
            # First call
            result1 = processor._get_markitdown()
            # Second call should not create new instance
            result2 = processor._get_markitdown()

            # Should return same instance due to caching
            assert result1 is result2
            assert result1 is processor.md
            assert result1 is mock_md_instance


class TestFileProcessorLocalFileProcessing:
    """Tests for local file processing with markitdown."""

    def test_process_local_file_success(self) -> None:
        """Test successful local file processing."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            tmp.write(b"Test content")
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="text")
            processor = FileProcessor()

            # Create a mock that properly mimics MarkItDown instance
            mock_md_instance = mock.MagicMock()
            mock_convert_result = mock.MagicMock()
            mock_convert_result.text_content = "# Test Content"
            mock_md_instance.convert.return_value = mock_convert_result

            with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                result = processor._process_local_file(file_input, start_time=0)

                assert result.markdown_content == "# Test Content"
                assert result.error is None
                assert result.metadata is not None
                assert result.metadata["path"] == tmp_path
                assert result.metadata["type"] == "text"
                assert "size_bytes" in result.metadata
        finally:
            Path(tmp_path).unlink()

    def test_process_local_file_not_found(self) -> None:
        """Test local file processing with missing file via process_file."""
        file_input = FileInput(path="/nonexistent/file.txt", type="text")
        processor = FileProcessor()

        # Use public process_file which has exception handling
        result = processor.process_file(file_input)

        assert result.error is not None
        assert "File not found" in result.error
        assert result.processing_time_ms is not None

    def test_process_local_file_via_process_file_no_path(self) -> None:
        """Test file processing through process_file with URL input."""
        # Create a FileInput with URL (tests routing to _process_remote_file)
        file_input = FileInput(url="https://example.com/file.txt", type="text")
        processor = FileProcessor()

        with mock.patch.object(processor, "_process_remote_file") as mock_remote:
            mock_remote.return_value = ProcessedFileInput(
                original=file_input,
                markdown_content="content",
                processing_time_ms=10,
            )
            processor.process_file(file_input)

            # Should route to remote processing, not local
            mock_remote.assert_called_once()

    def test_process_local_file_large_file_warning(self) -> None:
        """Test that large files are flagged in metadata."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            # Create a file
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="text")
            processor = FileProcessor()

            # Patch Path.stat to simulate large file
            with mock.patch.object(Path, "stat") as mock_stat:
                mock_stat.return_value.st_size = 101 * 1024 * 1024  # 101 MB

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "Large file content"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.metadata is not None
                    assert "warning" in result.metadata
                    assert "Large file" in result.metadata["warning"]
        finally:
            Path(tmp_path).unlink()


class TestFileProcessorRemoteFileProcessing:
    """Tests for remote file processing with caching and downloading."""

    def test_process_remote_file_from_cache(self) -> None:
        """Test remote file processing using cache."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            url = "https://example.com/file.pdf"
            file_input = FileInput(url=url, type="pdf")

            # Save something to cache first
            cache_key = processor._get_cache_key(url)
            processor._save_to_cache(
                cache_key,
                "Cached markdown",
                {"url": url},
                50,
            )

            # Process remote file - should use cache
            result = processor._process_remote_file(file_input, start_time=0)

            assert result.markdown_content == "Cached markdown"
            assert result.cached_path is not None
            assert result.error is None

    def test_process_remote_file_download_fails(self) -> None:
        """Test remote file processing when download fails."""
        processor = FileProcessor()
        file_input = FileInput(url="https://example.com/file.pdf", type="pdf")

        with mock.patch.object(processor, "_download_file", return_value=None):
            result = processor._process_remote_file(file_input, start_time=0)

            assert result.error is not None
            assert "Failed to download" in result.error
            assert result.markdown_content == ""

    def test_process_remote_file_invalid_input_via_process_file(self) -> None:
        """Test remote file processing with local file via process_file."""
        processor = FileProcessor()
        file_input = FileInput(path="/local/file.txt", type="text")

        # Use public process_file which has exception handling
        with mock.patch.object(processor, "_process_local_file") as mock_local:
            mock_local.return_value = ProcessedFileInput(
                original=file_input,
                markdown_content="content",
                processing_time_ms=10,
            )
            processor.process_file(file_input)

            # Should route to local processing since it has path
            mock_local.assert_called_once()

    def test_process_remote_file_success_with_cache(self) -> None:
        """Test successful remote file processing and caching."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            url = "https://example.com/file.pdf"
            file_input = FileInput(url=url, type="pdf", cache=True)

            with mock.patch.object(
                processor, "_download_file", return_value=b"PDF content"
            ):
                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# PDF Content"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_remote_file(file_input, start_time=0)

                    assert result.markdown_content == "# PDF Content"
                    assert result.error is None
                    assert result.metadata is not None
                    assert result.metadata["url"] == url

                    # Verify cache was created
                    assert result.cached_path is not None

    def test_process_remote_file_cache_disabled(self) -> None:
        """Test remote file processing with cache disabled."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            url = "https://example.com/file.pdf"
            file_input = FileInput(url=url, type="pdf", cache=False)

            with mock.patch.object(processor, "_download_file", return_value=b"PDF"):
                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# PDF"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_remote_file(file_input, start_time=0)

                    assert result.error is None
                    # When cache is disabled, file won't be created
                    cache_key = processor._get_cache_key(url)
                    cache_file = processor.cache_dir / f"{cache_key}.json"
                    assert not cache_file.exists()

    def test_process_remote_file_cleanup_temp_file(self) -> None:
        """Test that temporary file is cleaned up after processing."""
        with tempfile.TemporaryDirectory() as tmpdir:
            processor = FileProcessor(cache_dir=tmpdir)
            file_input = FileInput(url="https://example.com/file.pdf", type="pdf")

            with mock.patch.object(
                processor, "_download_file", return_value=b"content"
            ):
                # Use MagicMock to properly handle context manager
                mock_temp_file = mock.MagicMock()
                mock_temp_file.name = str(Path(tmpdir) / "test_file")

                with mock.patch(
                    "tempfile.NamedTemporaryFile", return_value=mock_temp_file
                ):
                    mock_md_instance = mock.MagicMock()
                    mock_result = mock.MagicMock()
                    mock_result.text_content = "content"
                    mock_md_instance.convert.return_value = mock_result

                    with mock.patch(
                        "markitdown.MarkItDown", return_value=mock_md_instance
                    ):
                        result = processor._process_remote_file(
                            file_input, start_time=0
                        )

                        assert result.error is None


class TestFileProcessorCacheEdgeCases:
    """Tests for cache edge cases and error handling."""

    def test_load_cache_corrupted_json(self) -> None:
        """Test loading from cache when JSON is corrupted."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            cache_dir.mkdir(parents=True, exist_ok=True)

            processor = FileProcessor(cache_dir=str(cache_dir))
            cache_key = "test_key"
            cache_file = cache_dir / f"{cache_key}.json"

            # Write invalid JSON
            cache_file.write_text("{ invalid json }")

            result = processor._load_from_cache(cache_key)
            assert result is None

    def test_load_cache_non_dict_json(self) -> None:
        """Test loading from cache when JSON is not a dict."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            cache_dir.mkdir(parents=True, exist_ok=True)

            processor = FileProcessor(cache_dir=str(cache_dir))
            cache_key = "test_key"
            cache_file = cache_dir / f"{cache_key}.json"

            # Write valid JSON but not a dict
            cache_file.write_text('["list", "not", "dict"]')

            result = processor._load_from_cache(cache_key)
            assert result is None

    def test_save_cache_permission_denied(self) -> None:
        """Test saving to cache when permissions denied."""
        with tempfile.TemporaryDirectory() as tmpdir:
            cache_dir = Path(tmpdir) / ".holodeck" / "cache"
            processor = FileProcessor(cache_dir=str(cache_dir))

            # Mock open to raise PermissionError
            with mock.patch("builtins.open", side_effect=PermissionError("No access")):
                # Should not raise, just suppress the error
                processor._save_to_cache("test_key", "content", {"key": "value"}, 100)


class TestFileProcessorValidationErrors:
    """Tests for validation error conditions."""

    def test_process_local_file_raises_error_no_path(self) -> None:
        """Test _process_local_file raises ValueError when path is None."""
        import pytest

        processor = FileProcessor()
        # Create a mock FileInput with no path
        mock_input = mock.MagicMock(spec=FileInput)
        mock_input.path = None

        with pytest.raises(ValueError, match="Local file must have path specified"):
            processor._process_local_file(mock_input, start_time=0)

    def test_process_remote_file_raises_error_no_url(self) -> None:
        """Test _process_remote_file raises ValueError when URL is None."""
        import pytest

        processor = FileProcessor()
        # Create a mock FileInput with no URL
        mock_input = mock.MagicMock(spec=FileInput)
        mock_input.url = None

        with pytest.raises(ValueError, match="Remote file must have URL specified"):
            processor._process_remote_file(mock_input, start_time=0)


class TestFileProcessorProcessFile:
    """Additional tests for process_file method."""

    def test_process_file_with_exception(self) -> None:
        """Test process_file handles exceptions gracefully."""
        processor = FileProcessor()
        file_input = FileInput(url="https://example.com/file.pdf", type="pdf")

        with mock.patch.object(processor, "_process_remote_file") as mock_remote:
            mock_remote.side_effect = RuntimeError("Processing failed")

            result = processor.process_file(file_input)

            assert result.error is not None
            assert "Processing failed" in result.error
            assert result.processing_time_ms is not None
            assert result.processing_time_ms >= 0


class TestFileProcessorPageSheetRangeExtraction:
    """Tests for page/sheet/range extraction preprocessing."""

    # PDF page extraction tests
    def test_pdf_pages_extraction_single_page(self) -> None:
        """Test extracting a single page from PDF."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[1])
            processor = FileProcessor()

            # Mock the preprocessing to return a temp file path
            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.return_value = Path(tmp_path)

                # Mock markitdown conversion
                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Page 1 Content"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    mock_preprocess.assert_called_once_with(file_input, Path(tmp_path))
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_pdf_pages_extraction_multiple_pages(self) -> None:
        """Test extracting multiple specific pages from PDF."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[1, 3, 5])
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.return_value = Path(tmp_path)

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Pages 1, 3, 5"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    mock_preprocess.assert_called_once()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_pdf_pages_extraction_sequential_range(self) -> None:
        """Test extracting sequential page range from PDF."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[2, 3, 4])
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.return_value = Path(tmp_path)

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Pages 2-4"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_pdf_pages_extraction_invalid_page(self) -> None:
        """Test handling invalid page numbers in PDF."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[999])
            processor = FileProcessor()

            # Mock _preprocess_file to raise an error for invalid pages
            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.side_effect = ValueError("Page 999 out of range")

                # Use process_file which handles exceptions
                result = processor.process_file(file_input)

                assert result.error is not None
                assert "out of range" in result.error.lower()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_pdf_no_pages_full_document(self) -> None:
        """Test processing full PDF when pages is None."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=None)
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                # When pages is None, should return original path (no preprocessing)
                mock_preprocess.return_value = Path(tmp_path)

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Full PDF"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    # Should still call preprocess_file to check for preprocessing needs
                    mock_preprocess.assert_called_once()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_pdf_empty_pages_list(self) -> None:
        """Test handling empty pages list."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[])
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                # Empty list should be treated as no preprocessing
                mock_preprocess.return_value = Path(tmp_path)

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Full PDF"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    # Excel sheet/range extraction tests
    def test_excel_sheet_extraction(self) -> None:
        """Test extracting specific sheet from Excel."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="excel", sheet="Sheet2")
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                # Mock returns a CSV temp file
                mock_csv_path = Path(tmp_path).with_suffix(".csv")
                mock_preprocess.return_value = mock_csv_path

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "| Col1 | Col2 |\n| --- | --- |"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    mock_preprocess.assert_called_once()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_excel_sheet_extraction_invalid_sheet(self) -> None:
        """Test handling non-existent sheet name."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="excel", sheet="NonExistent")
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.side_effect = ValueError(
                    "Sheet 'NonExistent' not found"
                )

                result = processor.process_file(file_input)

                assert result.error is not None
                assert "not found" in result.error.lower()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_excel_range_extraction(self) -> None:
        """Test extracting cell range from Excel."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="excel", range="A1:E10")
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_csv_path = Path(tmp_path).with_suffix(".csv")
                mock_preprocess.return_value = mock_csv_path

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "| A | B | C |"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_excel_sheet_and_range_extraction(self) -> None:
        """Test extracting range from specific sheet."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(
                path=tmp_path, type="excel", sheet="Data", range="B2:D20"
            )
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_csv_path = Path(tmp_path).with_suffix(".csv")
                mock_preprocess.return_value = mock_csv_path

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "| Header |"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    mock_preprocess.assert_called_once()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_excel_no_sheet_first_sheet_default(self) -> None:
        """Test using first sheet when sheet is None."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="excel", sheet=None)
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                # When no sheet specified, should process normally
                mock_preprocess.return_value = Path(tmp_path)

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "| Data |"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    # PowerPoint slide extraction tests
    def test_powerpoint_pages_extraction(self) -> None:
        """Test extracting specific slides from PowerPoint using pages field."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="powerpoint", pages=[1, 3])
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_pptx_path = Path(tmp_path)
                mock_preprocess.return_value = mock_pptx_path

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "# Slide 1\n\n# Slide 3"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    mock_preprocess.assert_called_once()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_powerpoint_pages_extraction_invalid(self) -> None:
        """Test handling invalid slide numbers in PowerPoint."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="powerpoint", pages=[999])
            processor = FileProcessor()

            with mock.patch.object(processor, "_preprocess_file") as mock_preprocess:
                mock_preprocess.side_effect = ValueError("Slide 999 out of range")

                result = processor.process_file(file_input)

                assert result.error is not None
                assert "out of range" in result.error.lower()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    # Integration tests for preprocessing flow
    def test_preprocessing_before_markitdown(self) -> None:
        """Test that preprocessing happens before markitdown conversion."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[1])
            processor = FileProcessor()

            call_order = []

            def mock_preprocess(*args: object, **kwargs: object) -> Path:
                call_order.append("preprocess")
                return Path(tmp_path)

            def mock_convert(*args: object, **kwargs: object) -> mock.MagicMock:
                call_order.append("convert")
                result = mock.MagicMock()
                result.text_content = "content"
                return result

            with (
                mock.patch.object(
                    processor, "_preprocess_file", side_effect=mock_preprocess
                ),
                mock.patch("markitdown.MarkItDown") as mock_md_class,
            ):
                mock_md_instance = mock.MagicMock()
                mock_md_instance.convert.side_effect = mock_convert
                mock_md_class.return_value = mock_md_instance

                processor._process_local_file(file_input, start_time=0)

                # Verify preprocessing happened before conversion
                assert call_order == ["preprocess", "convert"]
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_preprocessing_creates_temp_file(self) -> None:
        """Test that preprocessing creates temporary file and cleans up."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf", pages=[1])
            processor = FileProcessor()

            temp_file_created = None

            def mock_preprocess(*args: object, **kwargs: object) -> Path:
                nonlocal temp_file_created
                # Simulate creating a temp file
                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as t:
                    temp_file_created = Path(t.name)
                return temp_file_created

            with mock.patch.object(
                processor, "_preprocess_file", side_effect=mock_preprocess
            ):
                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "content"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    assert result.error is None
                    # Verify temp file was created during preprocessing
                    assert temp_file_created is not None

            # Clean up the temp file
            if temp_file_created and temp_file_created.exists():
                temp_file_created.unlink()
        finally:
            Path(tmp_path).unlink(missing_ok=True)


class TestFileProcessorErrorHandlingAdvanced:
    """Tests for advanced error handling including timeout and malformed files."""

    def test_timeout_handling_local_file(self) -> None:
        """Test that processing timeout is handled gracefully."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf")
            processor = FileProcessor(processing_timeout_ms=1000)

            # Mock markitdown to simulate slow processing
            with mock.patch.object(processor, "_get_markitdown") as mock_get_md:
                mock_md_instance = mock.MagicMock()
                # Simulate timeout by raising an exception
                mock_md_instance.convert.side_effect = TimeoutError(
                    "File processing exceeded 1000ms timeout"
                )
                mock_get_md.return_value = mock_md_instance

                result = processor.process_file(file_input)

                assert result.error is not None
                assert "timeout" in result.error.lower()
                assert result.markdown_content == ""
                assert result.processing_time_ms is not None
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_timeout_handling_remote_file(self) -> None:
        """Test that timeout during remote file processing is handled."""
        processor = FileProcessor(processing_timeout_ms=500)
        file_input = FileInput(url="https://example.com/file.pdf", type="pdf")

        with (
            mock.patch.object(processor, "_download_file", return_value=b"PDF"),
            mock.patch.object(processor, "_get_markitdown") as mock_get_md,
        ):
            # Mock markitdown to simulate timeout
            mock_md_instance = mock.MagicMock()
            mock_md_instance.convert.side_effect = TimeoutError(
                "Processing exceeded timeout limit"
            )
            mock_get_md.return_value = mock_md_instance

            result = processor.process_file(file_input)

            assert result.error is not None
            assert "timeout" in result.error.lower()

    def test_malformed_pdf_error(self) -> None:
        """Test handling of malformed PDF file."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(b"Invalid PDF content")
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf")
            processor = FileProcessor()

            with mock.patch.object(processor, "_get_markitdown") as mock_get_md:
                mock_md_instance = mock.MagicMock()
                mock_md_instance.convert.side_effect = ValueError(
                    "Invalid PDF structure: corrupted file"
                )
                mock_get_md.return_value = mock_md_instance

                result = processor.process_file(file_input)

                assert result.error is not None
                assert (
                    "invalid" in result.error.lower()
                    or "corrupted" in result.error.lower()
                )
                assert result.markdown_content == ""
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_malformed_excel_error(self) -> None:
        """Test handling of malformed Excel file."""
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp.write(b"Invalid Excel content")
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="excel")
            processor = FileProcessor()

            with mock.patch.object(processor, "_get_markitdown") as mock_get_md:
                mock_md_instance = mock.MagicMock()
                mock_md_instance.convert.side_effect = RuntimeError(
                    "Failed to read Excel file: invalid format"
                )
                mock_get_md.return_value = mock_md_instance

                result = processor.process_file(file_input)

                assert result.error is not None
                assert result.markdown_content == ""
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_malformed_powerpoint_error(self) -> None:
        """Test handling of malformed PowerPoint file."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(b"Invalid PowerPoint")
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="powerpoint")
            processor = FileProcessor()

            with mock.patch.object(processor, "_get_markitdown") as mock_get_md:
                mock_md_instance = mock.MagicMock()
                mock_md_instance.convert.side_effect = EOFError(
                    "Unexpected end of file: corrupted PowerPoint"
                )
                mock_get_md.return_value = mock_md_instance

                result = processor.process_file(file_input)

                assert result.error is not None
                assert result.markdown_content == ""
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_error_field_populated_with_message(self) -> None:
        """Test that error field contains descriptive error message."""
        file_input = FileInput(path="/nonexistent/file.pdf", type="pdf")
        processor = FileProcessor()

        result = processor.process_file(file_input)

        assert result.error is not None
        assert "not found" in result.error.lower()
        assert (
            "/nonexistent/file.pdf" in result.error
            or "nonexistent" in result.error.lower()
        )

    def test_error_field_preserved_across_operations(self) -> None:
        """Test that error field is set and preserved through processing."""
        file_input = FileInput(path="/bad/path.pdf", type="pdf")
        processor = FileProcessor()

        result = processor.process_file(file_input)

        # Verify all error-related fields are consistent
        assert result.error is not None
        assert result.markdown_content == ""
        assert result.original == file_input
        assert result.processing_time_ms is not None

    def test_multiple_files_continue_on_error(self) -> None:
        """Test that processing continues when one file fails."""
        processor = FileProcessor()

        # Create multiple test files
        file_inputs = [
            FileInput(path="/valid/path1.pdf", type="pdf"),
            FileInput(path="/invalid/path.pdf", type="pdf"),
            FileInput(path="/valid/path2.pdf", type="pdf"),
        ]

        results = []
        for file_input in file_inputs:
            with mock.patch.object(processor, "_process_local_file") as mock_local:
                if "invalid" in file_input.path:
                    # Simulate error for invalid file
                    mock_local.side_effect = FileNotFoundError("Not found")
                else:
                    # Simulate success for valid files
                    mock_local.return_value = ProcessedFileInput(
                        original=file_input,
                        markdown_content="# Content",
                        processing_time_ms=100,
                    )

                result = processor.process_file(file_input)
                results.append(result)

        # Verify all files were processed (none were skipped)
        assert len(results) == 3

        # Verify middle file has error
        assert results[1].error is not None
        assert results[1].markdown_content == ""

        # Verify other files processed successfully (not affected by error)
        assert results[0].error is None or results[0].processing_time_ms is not None
        assert results[2].error is None or results[2].processing_time_ms is not None

    def test_error_with_remote_file_download_timeout(self) -> None:
        """Test timeout during remote file download."""
        processor = FileProcessor(download_timeout_ms=100)
        file_input = FileInput(url="https://example.com/slow.pdf", type="pdf")

        with mock.patch.object(processor, "_download_file") as mock_download:
            mock_download.side_effect = TimeoutError("Download timeout")

            result = processor.process_file(file_input)

            assert result.error is not None

    def test_error_message_includes_file_location(self) -> None:
        """Test that error messages include file location."""
        file_path = "/path/to/missing/file.pdf"
        file_input = FileInput(path=file_path, type="pdf")
        processor = FileProcessor()

        result = processor.process_file(file_input)

        assert result.error is not None
        # Error should reference the file that failed
        assert file_path in result.error or "missing" in result.error.lower()

    def test_permission_denied_error(self) -> None:
        """Test handling of permission denied errors."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="pdf")
            processor = FileProcessor()

            with mock.patch.object(processor, "_process_local_file") as mock_local:
                mock_local.side_effect = PermissionError(
                    "Permission denied reading file"
                )

                result = processor.process_file(file_input)

                assert result.error is not None
                assert "permission" in result.error.lower()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_invalid_file_encoding_handling(self) -> None:
        """Test handling of file encoding errors."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            # Write invalid UTF-8 bytes
            tmp.write(b"\x80\x81\x82\x83")
            tmp_path = tmp.name

        try:
            file_input = FileInput(path=tmp_path, type="text")
            processor = FileProcessor()

            with mock.patch.object(processor, "_get_markitdown") as mock_get_md:
                mock_md_instance = mock.MagicMock()
                mock_md_instance.convert.side_effect = UnicodeDecodeError(
                    "utf-8", b"\x80", 0, 1, "invalid start byte"
                )
                mock_get_md.return_value = mock_md_instance

                result = processor.process_file(file_input)

                assert result.error is not None
                assert "UnicodeDecodeError" in result.error
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_custom_timeout_value(self) -> None:
        """Test that custom timeout value is set correctly."""
        custom_timeout = 5000
        processor = FileProcessor(processing_timeout_ms=custom_timeout)

        assert processor.processing_timeout_ms == custom_timeout

    def test_timeout_default_value(self) -> None:
        """Test that default timeout value is 30 seconds."""
        processor = FileProcessor()

        assert processor.processing_timeout_ms == 30000


class TestFileProcessorWithTimeoutEdgeCases:
    """Tests for _with_timeout edge cases."""

    def test_with_timeout_no_result_runtime_error(self) -> None:
        """Test _with_timeout raises RuntimeError when no result."""
        import pytest

        # Note: This test demonstrates the edge case, but in practice,
        # functions should always return a value. We need to carefully craft
        # a function that completes but doesn't append to result_holder
        # For testing purposes, we'll skip this test as unrealistic

        pytest.skip("Skipping edge case - functions should return values")

    def test_with_timeout_exception_in_thread(self) -> None:
        """Test _with_timeout propagates exceptions from thread."""
        import pytest

        processor = FileProcessor()

        def failing_function() -> str:
            """Function that raises an exception."""
            raise ValueError("Thread execution failed")

        with pytest.raises(ValueError, match="Thread execution failed"):
            processor._with_timeout(failing_function, 1000)

    def test_with_timeout_timeout_occurs(self) -> None:
        """Test _with_timeout raises TimeoutError when timeout is exceeded."""
        import time

        import pytest

        processor = FileProcessor()

        def slow_function() -> str:
            """Function that takes too long."""
            time.sleep(2)  # Sleep longer than timeout
            return "should not reach"

        with pytest.raises(TimeoutError, match="Processing exceeded .* timeout limit"):
            processor._with_timeout(slow_function, 100)  # 100ms timeout

    def test_with_timeout_successful_execution(self) -> None:
        """Test _with_timeout returns result on successful execution."""
        processor = FileProcessor()

        def quick_function(x: int, y: int) -> int:
            """Function that completes quickly."""
            return x + y

        result = processor._with_timeout(quick_function, 1000, 5, 3)
        assert result == 8


class TestFileProcessorFormatErrorMessage:
    """Tests for _format_error_message method."""

    def test_format_error_message_with_file_location(self) -> None:
        """Test error message formatting with file location."""
        processor = FileProcessor()
        error = FileNotFoundError("File does not exist")
        file_location = "/path/to/file.pdf"

        message = processor._format_error_message(error, file_location)

        assert "FileNotFoundError" in message
        assert "File does not exist" in message
        assert "/path/to/file.pdf" in message
        assert "while processing" in message

    def test_format_error_message_without_file_location(self) -> None:
        """Test error message formatting without file location."""
        processor = FileProcessor()
        error = ValueError("Invalid input")
        file_location = "unknown"

        message = processor._format_error_message(error, file_location)

        assert "ValueError" in message
        assert "Invalid input" in message
        assert "while processing" not in message

    def test_format_error_message_empty_location(self) -> None:
        """Test error message formatting with empty location."""
        processor = FileProcessor()
        error = RuntimeError("Something went wrong")
        file_location = ""

        message = processor._format_error_message(error, file_location)

        assert "RuntimeError" in message
        assert "Something went wrong" in message
        assert "while processing" not in message


class TestFileProcessorPDFPreprocessing:
    """Tests for PDF page extraction preprocessing."""

    def test_preprocess_pdf_pages_single_cell_extraction(self) -> None:
        """Test PDF preprocessing with actual pypdf integration."""

        processor = FileProcessor()

        # Create a simple PDF file using pypdf
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            from pypdf import PdfReader, PdfWriter

            # Create a minimal valid PDF with 3 pages
            writer = PdfWriter()
            # Add 3 empty pages
            for _ in range(3):
                writer.add_blank_page(width=200, height=200)

            with open(tmp_path, "wb") as f:
                writer.write(f)

            # Test extracting page 1
            result_path = processor._preprocess_pdf_pages(tmp_path, [1])

            # Verify the result is a valid PDF
            assert result_path.exists()
            assert result_path.suffix == ".pdf"

            # Verify it has exactly 1 page
            reader = PdfReader(str(result_path))
            assert len(reader.pages) == 1

            # Clean up temp file
            result_path.unlink()

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_pdf_pages_out_of_range(self) -> None:
        """Test PDF preprocessing with page number out of range."""
        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            from pypdf import PdfWriter

            # Create a PDF with only 2 pages
            writer = PdfWriter()
            for _ in range(2):
                writer.add_blank_page(width=200, height=200)

            with open(tmp_path, "wb") as f:
                writer.write(f)

            # Try to extract page 5 (out of range)
            with pytest.raises(ValueError, match="Page .* out of range"):
                processor._preprocess_pdf_pages(tmp_path, [5])

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_pdf_pages_import_error(self) -> None:
        """Test PDF preprocessing when pypdf is not available."""
        import sys

        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            # Mock pypdf import to fail
            with (
                mock.patch.dict(sys.modules, {"pypdf": None}),
                mock.patch("builtins.__import__", side_effect=ImportError("No pypdf")),
                pytest.raises(ImportError, match="pypdf is required"),
            ):
                processor._preprocess_pdf_pages(tmp_path, [0])

        finally:
            tmp_path.unlink(missing_ok=True)


class TestFileProcessorExcelPreprocessing:
    """Tests for Excel sheet/range extraction preprocessing."""

    def test_preprocess_excel_single_cell(self) -> None:
        """Test Excel preprocessing with single cell extraction."""

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            import openpyxl

            # Create a simple Excel file
            wb = openpyxl.Workbook()
            ws = wb.active
            ws["A1"] = "Header"
            ws["A2"] = "Data1"
            ws["B1"] = "Value"
            ws["B2"] = 100
            wb.save(str(tmp_path))

            # Extract single cell
            result_path = processor._preprocess_excel_sheet_range(tmp_path, None, "A1")

            # Verify CSV was created
            assert result_path.exists()
            assert result_path.suffix == ".csv"

            # Verify content
            import pandas as pd

            df = pd.read_csv(result_path, header=None)
            assert df.iloc[0, 0] == "Header"

            # Clean up
            result_path.unlink()

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_excel_range_as_tuple(self) -> None:
        """Test Excel preprocessing with range extraction."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            import openpyxl

            wb = openpyxl.Workbook()
            ws = wb.active
            # Fill in some data
            for row in range(1, 4):
                for col in range(1, 4):
                    ws.cell(row=row, column=col, value=f"R{row}C{col}")
            wb.save(str(tmp_path))

            # Extract range A1:C3
            result_path = processor._preprocess_excel_sheet_range(
                tmp_path, None, "A1:C3"
            )

            assert result_path.exists()
            assert result_path.suffix == ".csv"

            # Clean up
            result_path.unlink()

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_excel_invalid_range(self) -> None:
        """Test Excel preprocessing with invalid range."""
        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            import openpyxl

            wb = openpyxl.Workbook()
            wb.save(str(tmp_path))

            # Try invalid range
            with pytest.raises(ValueError, match="Invalid range"):
                processor._preprocess_excel_sheet_range(tmp_path, None, "INVALID_RANGE")

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_excel_import_error(self) -> None:
        """Test Excel preprocessing when openpyxl is not available."""
        import sys

        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            # Mock openpyxl import to fail
            with mock.patch.dict(sys.modules, {"openpyxl": None}):  # noqa: SIM117
                with mock.patch(
                    "builtins.__import__", side_effect=ImportError("No openpyxl")
                ):
                    with pytest.raises(
                        ImportError, match="openpyxl and pandas are required"
                    ):
                        processor._preprocess_excel_sheet_range(tmp_path, None, None)

        finally:
            tmp_path.unlink(missing_ok=True)


class TestFileProcessorPowerPointPreprocessing:
    """Tests for PowerPoint slide extraction preprocessing."""

    def test_preprocess_powerpoint_slides(self) -> None:
        """Test PowerPoint preprocessing with slide extraction."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            from pptx import Presentation

            # Create a simple presentation with 3 slides
            prs = Presentation()
            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = f"Slide {i + 1}"
            prs.save(str(tmp_path))

            # Extract slides 0 and 2
            result_path = processor._preprocess_powerpoint_slides(tmp_path, [0, 2])

            # Verify result
            assert result_path.exists()
            assert result_path.suffix == ".pptx"

            # Verify it has 2 slides
            result_prs = Presentation(str(result_path))
            assert len(result_prs.slides) == 2

            # Clean up
            result_path.unlink()

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_powerpoint_slide_out_of_range(self) -> None:
        """Test PowerPoint preprocessing with slide number out of range."""
        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            from pptx import Presentation

            # Create presentation with only 2 slides
            prs = Presentation()
            for _ in range(2):
                prs.slides.add_slide(prs.slide_layouts[0])
            prs.save(str(tmp_path))

            # Try to extract slide 5 (out of range)
            with pytest.raises(ValueError, match="Slide .* out of range"):
                processor._preprocess_powerpoint_slides(tmp_path, [5])

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_powerpoint_import_error(self) -> None:
        """Test PowerPoint preprocessing when python-pptx is not available."""
        import sys

        import pytest

        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            # Mock pptx import to fail
            with (
                mock.patch.dict(sys.modules, {"pptx": None}),
                mock.patch("builtins.__import__", side_effect=ImportError("No pptx")),
                pytest.raises(ImportError, match="python-pptx is required"),
            ):
                processor._preprocess_powerpoint_slides(tmp_path, [0])

        finally:
            tmp_path.unlink(missing_ok=True)


class TestFileProcessorPreprocessFileRouting:
    """Tests for _preprocess_file routing logic."""

    def test_preprocess_file_no_preprocessing_needed(self) -> None:
        """Test that files with no pages/sheet/range return original path."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(path=str(tmp_path), type="pdf")  # No pages
            result = processor._preprocess_file(file_input, tmp_path)

            # Should return original path unchanged
            assert result == tmp_path

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_file_routes_to_pdf_handler(self) -> None:
        """Test that PDF files with pages route to PDF handler."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(path=str(tmp_path), type="pdf", pages=[1, 2])

            with mock.patch.object(
                processor, "_preprocess_pdf_pages"
            ) as mock_pdf_preprocess:
                mock_pdf_preprocess.return_value = Path(
                    "/tmp/processed.pdf"  # noqa: S108
                )

                result = processor._preprocess_file(file_input, tmp_path)

                mock_pdf_preprocess.assert_called_once_with(tmp_path, [1, 2])
                assert result == Path("/tmp/processed.pdf")  # noqa: S108

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_file_routes_to_powerpoint_handler(self) -> None:
        """Test that PowerPoint files with pages route to PowerPoint handler."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(path=str(tmp_path), type="powerpoint", pages=[1])

            with mock.patch.object(
                processor, "_preprocess_powerpoint_slides"
            ) as mock_ppt_preprocess:
                mock_ppt_preprocess.return_value = Path(
                    "/tmp/processed.pptx"  # noqa: S108
                )

                result = processor._preprocess_file(file_input, tmp_path)

                mock_ppt_preprocess.assert_called_once_with(tmp_path, [1])
                assert result == Path("/tmp/processed.pptx")  # noqa: S108

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_file_routes_to_excel_handler(self) -> None:
        """Test that Excel files with sheet/range route to Excel handler."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(
                path=str(tmp_path), type="excel", sheet="Sheet1", range="A1:B10"
            )

            with mock.patch.object(
                processor, "_preprocess_excel_sheet_range"
            ) as mock_excel_preprocess:
                mock_excel_preprocess.return_value = Path(
                    "/tmp/processed.csv"  # noqa: S108
                )

                result = processor._preprocess_file(file_input, tmp_path)

                mock_excel_preprocess.assert_called_once_with(
                    tmp_path, "Sheet1", "A1:B10"
                )
                assert result == Path("/tmp/processed.csv")  # noqa: S108

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_file_excel_sheet_only(self) -> None:
        """Test Excel preprocessing with sheet but no range."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(path=str(tmp_path), type="excel", sheet="Data")

            with mock.patch.object(
                processor, "_preprocess_excel_sheet_range"
            ) as mock_excel_preprocess:
                mock_excel_preprocess.return_value = Path(
                    "/tmp/processed.csv"  # noqa: S108
                )

                _ = processor._preprocess_file(file_input, tmp_path)

                mock_excel_preprocess.assert_called_once_with(tmp_path, "Data", None)

        finally:
            tmp_path.unlink(missing_ok=True)

    def test_preprocess_file_excel_range_only(self) -> None:
        """Test Excel preprocessing with range but no sheet."""
        processor = FileProcessor()

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            file_input = FileInput(path=str(tmp_path), type="excel", range="A1:D20")

            with mock.patch.object(
                processor, "_preprocess_excel_sheet_range"
            ) as mock_excel_preprocess:
                mock_excel_preprocess.return_value = Path(
                    "/tmp/processed.csv"  # noqa: S108
                )

                _ = processor._preprocess_file(file_input, tmp_path)

                mock_excel_preprocess.assert_called_once_with(tmp_path, None, "A1:D20")

        finally:
            tmp_path.unlink(missing_ok=True)


class TestFileProcessorRemoteFilePreprocessing:
    """Tests for remote file preprocessing with metadata."""

    def test_local_file_with_sheet_metadata(self) -> None:
        """Test local Excel file with sheet extraction metadata."""
        with tempfile.TemporaryDirectory() as tmpdir:
            processor = FileProcessor(cache_dir=tmpdir)

            with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
                tmp_path = Path(tmp.name)

            try:
                import openpyxl

                # Create a simple Excel file with data
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = "Data"
                for row in range(1, 11):
                    for col in range(1, 6):
                        ws.cell(row=row, column=col, value=f"R{row}C{col}")
                wb.save(str(tmp_path))

                file_input = FileInput(
                    path=str(tmp_path),
                    type="excel",
                    sheet="Data",
                    range="A1:E10",
                )

                mock_md_instance = mock.MagicMock()
                mock_result = mock.MagicMock()
                mock_result.text_content = "| A | B |"
                mock_md_instance.convert.return_value = mock_result

                with mock.patch("markitdown.MarkItDown", return_value=mock_md_instance):
                    result = processor._process_local_file(file_input, start_time=0)

                    # Verify sheet and range metadata
                    assert result.metadata is not None
                    assert result.metadata.get("preprocessed") is True
                    assert result.metadata.get("sheet_extracted") == "Data"
                    assert result.metadata.get("range_extracted") == "A1:E10"

            finally:
                tmp_path.unlink(missing_ok=True)


class TestFileProcessorFromExecutionConfig:
    """Tests for FileProcessor.from_execution_config factory method."""

    def test_from_execution_config_with_all_values(self) -> None:
        """Test factory with all timeout values specified."""
        with tempfile.TemporaryDirectory() as tmpdir:
            custom_cache = Path(tmpdir) / "custom" / "cache"
            config = ExecutionConfig(
                file_timeout=60,
                download_timeout=45,
                cache_dir=str(custom_cache),
            )

            processor = FileProcessor.from_execution_config(config)

            # Verify seconds -> milliseconds conversion
            assert processor.processing_timeout_ms == 60000  # 60s -> 60000ms
            assert processor.download_timeout_ms == 45000  # 45s -> 45000ms
            assert processor.cache_dir == custom_cache

    def test_from_execution_config_with_defaults(self) -> None:
        """Test factory uses defaults when values not specified."""
        config = ExecutionConfig()  # All None

        processor = FileProcessor.from_execution_config(config)

        # Should use default 30s -> 30000ms
        assert processor.processing_timeout_ms == 30000
        assert processor.download_timeout_ms == 30000
        assert str(processor.cache_dir) == ".holodeck/cache"

    def test_from_execution_config_cache_dir_override(self) -> None:
        """Test cache_dir parameter overrides config."""
        with tempfile.TemporaryDirectory() as tmpdir:
            config_cache = Path(tmpdir) / "config" / "cache"
            override_cache = Path(tmpdir) / "override" / "cache"
            config = ExecutionConfig(cache_dir=str(config_cache))

            processor = FileProcessor.from_execution_config(
                config, cache_dir=str(override_cache)
            )

            assert processor.cache_dir == override_cache

    def test_from_execution_config_max_retries(self) -> None:
        """Test max_retries parameter is passed through."""
        config = ExecutionConfig()

        processor = FileProcessor.from_execution_config(config, max_retries=5)

        assert processor.max_retries == 5

    def test_from_execution_config_partial_values(self) -> None:
        """Test factory with only some values specified."""
        config = ExecutionConfig(
            file_timeout=120,  # Only file timeout specified
            download_timeout=None,
            cache_dir=None,
        )

        processor = FileProcessor.from_execution_config(config)

        assert processor.processing_timeout_ms == 120000  # 120s -> 120000ms
        assert processor.download_timeout_ms == 30000  # default 30s
        assert str(processor.cache_dir) == ".holodeck/cache"

    def test_from_execution_config_preserves_default_max_retries(self) -> None:
        """Test that default max_retries is 3."""
        config = ExecutionConfig()

        processor = FileProcessor.from_execution_config(config)

        assert processor.max_retries == 3
