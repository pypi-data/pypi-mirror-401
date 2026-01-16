"""PPTX extractor for PowerPoint presentations."""

import os
from datetime import datetime
from pathlib import Path

from semplex_cli.types import TextMetadata
from semplex_cli.utils.file_utils import get_file_owner

# Optional dependency - will be imported at runtime
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PptxExtractor:
    """Extract text content from PowerPoint presentations (.pptx)."""

    SUPPORTED_EXTENSIONS = {".pptx"}

    def __init__(self):
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPTX extraction. "
                "Install it with: pip install python-pptx"
            )

    def extract(self, file_path: Path) -> TextMetadata:
        """
        Extract text content from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            TextMetadata with file content and metadata
        """
        file_path = Path(file_path)
        stat = file_path.stat()

        # Read PPTX
        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        # Extract text from all slides
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []

            for shape in slide.shapes:
                # Extract text from shapes with text frames
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text_parts.append(" | ".join(row_text))

            if slide_text_parts:
                # Add slide separator with slide number
                text_parts.append(f"--- Slide {slide_num} ---")
                text_parts.extend(slide_text_parts)

        content = "\n\n".join(text_parts)

        # Calculate metrics
        lines = content.split("\n")
        words = content.split()

        return TextMetadata(
            filename=file_path.name,
            filepath=str(file_path.absolute()),
            file_type=".pptx",
            file_size=stat.st_size,
            file_owner=get_file_owner(file_path),
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
            extracted_at=datetime.now().isoformat(),
            content=content,
            char_count=len(content),
            word_count=len(words),
            line_count=len(lines),
            page_count=slide_count,
            section_count=slide_count,  # Each slide is a section
        )
