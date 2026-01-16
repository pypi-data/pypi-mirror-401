# Core logic for Corpus and Answer

import os
import hashlib
import json
import pickle
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from dataclasses import dataclass
from datetime import datetime

# Simple placeholder implementations for testing
@dataclass
class Answer:
    text: str
    confidence: float
    sources: List[Any] = None
    
@dataclass 
class QueryResult:
    answer: Answer
    processing_time: float
    total_chunks_considered: int
    model_used: str

class DocCorpus:
    """Simple document corpus for testing."""
    
    def __init__(self, corpus_path: str = "./.docnav", llm_provider: str = "openai", 
                 llm_model: str = None, api_key: str = None, embedding_model: str = "all-MiniLM-L6-v2"):
        self.corpus_path = Path(corpus_path)
        self.llm_provider = llm_provider
        self.llm_model = llm_model
        self.api_key = api_key
        self.embedding_model = embedding_model
        self.chunks = []
        
        # Create corpus directory
        self.corpus_path.mkdir(exist_ok=True)
        
        # Load existing data if available
        index_file = self.corpus_path / "corpus_index.pkl"
        if index_file.exists():
            try:
                with open(index_file, 'rb') as f:
                    data = pickle.load(f)
                    self.chunks = data.get('chunks', [])
            except:
                pass
    
    def add_documents(self, sources: List[str], use_ocr: bool = False, 
                     chunk_size: int = 1000, update_existing: bool = False) -> int:
        """Add documents to corpus (placeholder implementation)."""
        # Simple placeholder - just count files
        count = 0
        for source in sources:
            source_path = Path(source)
            if source_path.exists():
                if source_path.is_file():
                    count += 1
                elif source_path.is_dir():
                    # Count supported files
                    for ext in ['.pdf', '.docx', '.txt', '.md']:
                        count += len(list(source_path.glob(f'**/*{ext}')))
        
        # Save index
        self._save_index()
        return count
    
    def ask(self, query: str, top_k: int = 5, where: Dict = None, 
            temperature: float = None) -> QueryResult:
        """Ask a question (placeholder implementation)."""
        answer = Answer(
            text=f"This is a placeholder answer for: {query}",
            confidence=0.8,
            sources=[]
        )
        return QueryResult(
            answer=answer,
            processing_time=0.5,
            total_chunks_considered=len(self.chunks),
            model_used=self.llm_model or "gpt-3.5-turbo"
        )
    
    def list_documents(self) -> List[Dict]:
        """List documents in corpus."""
        return []
    
    def get_stats(self) -> Dict[str, Any]:
        """Get corpus statistics."""
        return {
            "total_chunks": len(self.chunks),
            "corpus_path": str(self.corpus_path),
            "llm_provider": self.llm_provider,
            "embedding_model": self.embedding_model
        }
    
    def remove_document(self, file_path: str) -> int:
        """Remove document from corpus."""
        return 0
    
    def clear(self) -> bool:
        """Clear corpus."""
        self.chunks = []
        self._save_index()
        return True
    
    def _save_index(self):
        """Save corpus index to disk."""
        index_file = self.corpus_path / "corpus_index.pkl"
        with open(index_file, 'wb') as f:
            pickle.dump({'chunks': self.chunks}, f)

def ask(file_path: str, query: str, api_key: str, provider: str = "openai") -> str:
    """Quick ask function for a single document."""
    return f"Placeholder answer for {query} about {file_path}"

"""
Core implementation of DocNav library.
"""

import os
import json
import hashlib
import pickle
import warnings
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Iterator, Tuple
from dataclasses import dataclass, field
from datetime import datetime
import shutil

import numpy as np
from tqdm import tqdm
import faiss
import chromadb
from chromadb.config import Settings
import pypdf
from docx import Document
import markdown
import pandas as pd
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from sentence_transformers import SentenceTransformer
import openai
import google.generativeai as genai

# Suppress warnings
warnings.filterwarnings('ignore')

@dataclass
class DocumentChunk:
    """Represents a chunk of text from a document."""
    text: str
    metadata: Dict[str, Any]
    embedding: Optional[np.ndarray] = None
    chunk_id: str = field(default_factory=lambda: hashlib.md5(str(datetime.now()).encode()).hexdigest()[:16])
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            'text': self.text,
            'metadata': self.metadata,
            'chunk_id': self.chunk_id,
            'embedding': self.embedding.tolist() if self.embedding is not None else None
        }

@dataclass
class Answer:
    """Represents an answer from the corpus."""
    text: str
    sources: List[DocumentChunk]
    confidence: float
    query: str
    
    def __str__(self) -> str:
        return self.text
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            'text': self.text,
            'confidence': self.confidence,
            'query': self.query,
            'sources': [
                {
                    'text': chunk.text[:200] + '...' if len(chunk.text) > 200 else chunk.text,
                    'metadata': chunk.metadata
                }
                for chunk in self.sources
            ]
        }

@dataclass
class QueryResult:
    """Represents a query result with metadata."""
    answer: Answer
    processing_time: float
    total_chunks_considered: int
    model_used: str

class DocParser:
    """Handles parsing of various document formats."""
    
    @staticmethod
    def parse_pdf(file_path: str, use_ocr: bool = False) -> List[DocumentChunk]:
        """Parse PDF files."""
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    
                    # Fallback to OCR if no text found
                    if not text.strip() and use_ocr:
                        images = convert_from_path(file_path, first_page=page_num, last_page=page_num)
                        for img in images:
                            text = pytesseract.image_to_string(img)
                    
                    if text.strip():
                        chunks.append(DocumentChunk(
                            text=text,
                            metadata={
                                'file_path': file_path,
                                'file_name': os.path.basename(file_path),
                                'page': page_num,
                                'total_pages': len(pdf_reader.pages),
                                'type': 'pdf',
                                'timestamp': datetime.now().isoformat()
                            }
                        ))
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
        
        return chunks
    
    @staticmethod
    def parse_docx(file_path: str) -> List[DocumentChunk]:
        """Parse DOCX files."""
        chunks = []
        try:
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            if full_text:
                chunks.append(DocumentChunk(
                    text='\n'.join(full_text),
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'type': 'docx',
                        'timestamp': datetime.now().isoformat()
                    }
                ))
        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {e}")
        
        return chunks
    
    @staticmethod
    def parse_txt(file_path: str) -> List[DocumentChunk]:
        """Parse TXT files."""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            
            if text.strip():
                chunks.append(DocumentChunk(
                    text=text,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'type': 'txt',
                        'timestamp': datetime.now().isoformat()
                    }
                ))
        except Exception as e:
            print(f"Error parsing TXT {file_path}: {e}")
        
        return chunks
    
    @staticmethod
    def parse_markdown(file_path: str) -> List[DocumentChunk]:
        """Parse Markdown files."""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Convert markdown to plain text
            html = markdown.markdown(text)
            # Simple HTML to text conversion
            import re
            plain_text = re.sub(r'<[^>]+>', '', html)
            
            if plain_text.strip():
                chunks.append(DocumentChunk(
                    text=plain_text,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'type': 'markdown',
                        'timestamp': datetime.now().isoformat()
                    }
                ))
        except Exception as e:
            print(f"Error parsing Markdown {file_path}: {e}")
        
        return chunks
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[DocumentChunk]:
        """Parse PPTX files."""
        chunks = []
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                
                if slide_content:
                    slides_text.append(f"Slide {slide_num}:\n" + "\n".join(slide_content))
            
            if slides_text:
                chunks.append(DocumentChunk(
                    text='\n\n'.join(slides_text),
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'type': 'pptx',
                        'timestamp': datetime.now().isoformat()
                    }
                ))
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {e}")
        
        return chunks
    
    @staticmethod
    def parse_excel(file_path: str) -> List[DocumentChunk]:
        """Parse Excel/CSV files."""
        chunks = []
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)
            
            if isinstance(df, dict):
                # Multiple sheets
                for sheet_name, sheet_df in df.items():
                    text = f"Sheet: {sheet_name}\n\n{sheet_df.to_string()}"
                    chunks.append(DocumentChunk(
                        text=text,
                        metadata={
                            'file_path': file_path,
                            'file_name': os.path.basename(file_path),
                            'sheet': sheet_name,
                            'type': 'excel',
                            'timestamp': datetime.now().isoformat()
                        }
                    ))
            else:
                # Single sheet
                chunks.append(DocumentChunk(
                    text=df.to_string(),
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'type': 'excel',
                        'timestamp': datetime.now().isoformat()
                    }
                ))
        except Exception as e:
            print(f"Error parsing Excel {file_path}: {e}")
        
        return chunks
    
    @classmethod
    def parse_file(cls, file_path: str, use_ocr: bool = False) -> List[DocumentChunk]:
        """Parse any supported file type."""
        file_path = str(file_path)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_path.endswith('.pdf'):
            return cls.parse_pdf(file_path, use_ocr)
        elif file_path.endswith('.docx'):
            return cls.parse_docx(file_path)
        elif file_path.endswith(('.txt', '.text')):
            return cls.parse_txt(file_path)
        elif file_path.endswith(('.md', '.markdown')):
            return cls.parse_markdown(file_path)
        elif file_path.endswith('.pptx'):
            return cls.parse_pptx(file_path)
        elif file_path.endswith(('.xlsx', '.xls', '.csv')):
            return cls.parse_excel(file_path)
        else:
            # Try as text file
            try:
                return cls.parse_txt(file_path)
            except:
                raise ValueError(f"Unsupported file format: {file_path}")

class SmartChunker:
    """Intelligent chunking that preserves document structure."""
    
    def __init__(self, max_chunk_size: int = 1000, overlap: int = 100):
        self.max_chunk_size = max_chunk_size
        self.overlap = overlap
    
    def chunk_document(self, text: str, metadata: Dict[str, Any]) -> List[DocumentChunk]:
        """Split text into smart chunks."""
        chunks = []
        
        # Split by paragraphs first
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_size = len(para.split())
            
            if current_size + para_size > self.max_chunk_size and current_chunk:
                # Save current chunk
                chunk_text = '\n\n'.join(current_chunk)
                chunks.append(DocumentChunk(
                    text=chunk_text,
                    metadata=metadata.copy()
                ))
                
                # Start new chunk with overlap
                current_chunk = current_chunk[-2:] if len(current_chunk) > 2 else []
                current_size = sum(len(p.split()) for p in current_chunk)
            
            current_chunk.append(para)
            current_size += para_size
        
        # Add the last chunk
        if current_chunk:
            chunk_text = '\n\n'.join(current_chunk)
            chunks.append(DocumentChunk(
                text=chunk_text,
                metadata=metadata.copy()
            ))
        
        return chunks

class EmbeddingModel:
    """Handles text embeddings with support for multiple backends."""
    
    def __init__(self, model_name: str = "all-MiniLM-L6-v2", use_local: bool = True):
        self.use_local = use_local
        self.model_name = model_name
        
        if use_local:
            self.model = SentenceTransformer(model_name)
            self.dimensions = self.model.get_sentence_embedding_dimension()
        else:
            self.dimensions = 768  # Default for OpenAI
    
    def embed(self, texts: List[str]) -> np.ndarray:
        """Generate embeddings for texts."""
        if self.use_local:
            return self.model.encode(texts, show_progress_bar=False)
        else:
            # For API-based embeddings, we'd implement here
            raise NotImplementedError("API embeddings not implemented in this version")

class LLMProvider:
    """Handles LLM interactions with multiple providers."""
    
    def __init__(self, provider: str = "openai", model: str = None, api_key: str = None):
        self.provider = provider
        self.model = model or self._get_default_model(provider)
        self.api_key = api_key
        
        if provider == "openai":
            if not api_key:
                raise ValueError("OpenAI API key required")
            openai.api_key = api_key
        elif provider == "gemini":
            if not api_key:
                raise ValueError("Google API key required")
            genai.configure(api_key=api_key)
    
    def _get_default_model(self, provider: str) -> str:
        defaults = {
            "openai": "gpt-3.5-turbo",
            "gemini": "gemini-pro"
        }
        return defaults.get(provider, "gpt-3.5-turbo")
    
    def generate(self, prompt: str, context: str = "", temperature: float = 0.0) -> str:
        """Generate a response from the LLM."""
        full_prompt = f"""Based on the following context, answer the question. 
        If the answer cannot be found in the context, say so clearly.
        
        Context:
        {context}
        
        Question: {prompt}
        
        Answer:"""
        
        if self.provider == "openai":
            response = openai.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that answers questions based on provided documents."},
                    {"role": "user", "content": full_prompt}
                ],
                temperature=temperature,
                max_tokens=1000
            )
            return response.choices[0].message.content
        
        elif self.provider == "gemini":
            model = genai.GenerativeModel(self.model)
            response = model.generate_content(full_prompt)
            return response.text
        
        else:
            raise ValueError(f"Unsupported provider: {self.provider}")

class DocCorpus:
    """Main class for document corpus management and querying."""
    
    def __init__(self, 
                 corpus_path: Optional[str] = None,
                 embedding_model: str = "all-MiniLM-L6-v2",
                 use_local_embeddings: bool = True,
                 deterministic: bool = False):
        
        self.corpus_path = Path(corpus_path) if corpus_path else Path("./.docnav_corpus")
        self.corpus_path.mkdir(exist_ok=True)
        
        self.embedding_model = EmbeddingModel(embedding_model, use_local_embeddings)
        self.llm = None  # LLM will be set when asking questions
        
        self.deterministic = deterministic
        self.temperature = 0.0 if deterministic else 0.7
        
        # Initialize vector store
        self._init_vector_store()
        
        # Load existing corpus if exists
        self.chunks: List[DocumentChunk] = []
        self._load_corpus()
    
    def _init_vector_store(self):
        """Initialize the vector database."""
        self.chroma_client = chromadb.PersistentClient(
            path=str(self.corpus_path / "chroma_db")
        )
        
        self.collection = self.chroma_client.get_or_create_collection(
            name="docnav_corpus",
            metadata={"hnsw:space": "cosine"}
        )
    
    def _load_corpus(self):
        """Load existing corpus from disk."""
        index_file = self.corpus_path / "corpus_index.pkl"
        if index_file.exists():
            try:
                with open(index_file, 'rb') as f:
                    data = pickle.load(f)
                    self.chunks = data.get('chunks', [])
                    # Restore LLM provider info if available
                    if 'llm_provider' in data and not self.llm:
                        self.llm_provider = data['llm_provider']
                        self.llm_model = data.get('llm_model')
                print(f"Loaded existing corpus with {len(self.chunks)} chunks")
            except:
                self.chunks = []
    
    def _save_corpus(self):
        """Save corpus to disk."""
        index_file = self.corpus_path / "corpus_index.pkl"
        data = {
            'chunks': self.chunks,
            'llm_provider': getattr(self, 'llm_provider', 'openai'),
            'llm_model': getattr(self, 'llm_model', None)
        }
        with open(index_file, 'wb') as f:
            pickle.dump(data, f)
    
    def add_documents(self, 
                     sources: Union[str, List[str], Path, List[Path]],
                     use_ocr: bool = False,
                     chunk_size: int = 1000,
                     update_existing: bool = True):
        """Add documents to the corpus. Can update existing ones."""
        parser = DocParser()
        chunker = SmartChunker(chunk_size)
        
        if isinstance(sources, (str, Path)):
            sources = [sources]
        
        # Get all files
        all_files = []
        for source in sources:
            source = Path(source)
            if source.is_file():
                all_files.append(source)
            elif source.is_dir():
                for ext in ['.pdf', '.docx', '.txt', '.md', '.pptx', '.xlsx', '.xls', '.csv']:
                    all_files.extend(source.glob(f'**/*{ext}'))
        
        print(f"Found {len(all_files)} files to process")
        
        # Process files with progress bar
        new_chunks = []
        for file_path in tqdm(all_files, desc="Processing files"):
            file_hash = hashlib.md5(str(file_path).encode() + str(file_path.stat().st_mtime).encode()).hexdigest()
            
            # Check if file already processed (unless forced update)
            existing_chunk = next((c for c in self.chunks 
                                 if c.metadata.get('file_path') == str(file_path) and 
                                 c.metadata.get('file_hash') == file_hash), None)
            
            if existing_chunk and not update_existing:
                continue
            
            # Remove old chunks for this file if updating
            if update_existing:
                self.chunks = [c for c in self.chunks 
                              if c.metadata.get('file_path') != str(file_path)]
            
            try:
                # Parse file
                document_chunks = parser.parse_file(str(file_path), use_ocr)
                
                # Chunk each document
                for doc_chunk in document_chunks:
                    # Add file hash to metadata
                    doc_chunk.metadata['file_hash'] = file_hash
                    doc_chunk.metadata['processed_at'] = datetime.now().isoformat()
                    
                    # Further chunk if needed
                    sub_chunks = chunker.chunk_document(doc_chunk.text, doc_chunk.metadata)
                    
                    # Generate embeddings
                    for chunk in sub_chunks:
                        embedding = self.embedding_model.embed([chunk.text])[0]
                        chunk.embedding = embedding
                        new_chunks.append(chunk)
                
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
        
        # Add to corpus
        if new_chunks:
            self.chunks.extend(new_chunks)
            
            # Update vector store
            texts = [chunk.text for chunk in new_chunks]
            embeddings = [chunk.embedding.tolist() for chunk in new_chunks]
            metadatas = [chunk.metadata for chunk in new_chunks]
            ids = [chunk.chunk_id for chunk in new_chunks]
            
            self.collection.add(
                embeddings=embeddings,
                documents=texts,
                metadatas=metadatas,
                ids=ids
            )
            
            # Save corpus
            self._save_corpus()
            
            print(f"Added {len(new_chunks)} new chunks to corpus")
        
        return len(new_chunks)
    
    def remove_document(self, file_path: str):
        """Remove a document from the corpus."""
        file_path = str(Path(file_path).resolve())
        
        # Remove from memory
        initial_count = len(self.chunks)
        self.chunks = [c for c in self.chunks 
                      if c.metadata.get('file_path') != file_path]
        
        # Remove from vector store
        try:
            self.collection.delete(where={"file_path": file_path})
        except:
            pass
        
        self._save_corpus()
        removed_count = initial_count - len(self.chunks)
        print(f"Removed {removed_count} chunks for {file_path}")
        
        return removed_count
    
    def search(self, query: str, top_k: int = 5, where: Optional[Dict] = None) -> List[DocumentChunk]:
        """Search for relevant chunks."""
        # Generate query embedding
        query_embedding = self.embedding_model.embed([query])[0]
        
        # Search in vector store
        results = self.collection.query(
            query_embeddings=[query_embedding.tolist()],
            n_results=top_k,
            where=where,
            include=["documents", "metadatas", "distances"]
        )
        
        # Convert to DocumentChunk objects
        chunks = []
        for doc, metadata, distance in zip(results['documents'][0], 
                                          results['metadatas'][0], 
                                          results['distances'][0]):
            chunk = DocumentChunk(
                text=doc,
                metadata=metadata,
                embedding=None
            )
            chunks.append(chunk)
        
        return chunks
    
    def ask(self, 
            query: str, 
            llm_provider: str = "openai",
            llm_model: str = None,
            api_key: str = None,
            top_k: int = 5,
            where: Optional[Dict] = None,
            temperature: Optional[float] = None) -> QueryResult:
        """Ask a question to the corpus."""
        import time
        start_time = time.time()
        
        # Initialize LLM provider for this query
        if not api_key:
            raise ValueError("API key required for LLM provider.")
        
        self.llm = LLMProvider(llm_provider, llm_model, api_key)
        
        # Search for relevant chunks
        relevant_chunks = self.search(query, top_k=top_k * 2, where=where)
        
        if not relevant_chunks:
            return QueryResult(
                answer=Answer(
                    text="No relevant documents found to answer this question.",
                    sources=[],
                    confidence=0.0,
                    query=query
                ),
                processing_time=time.time() - start_time,
                total_chunks_considered=0,
                model_used=self.llm.model
            )
        
        # Prepare context
        context_parts = []
        for i, chunk in enumerate(relevant_chunks[:top_k]):
            source_info = f"[Source: {chunk.metadata.get('file_name', 'Unknown')}"
            if 'page' in chunk.metadata:
                source_info += f", Page: {chunk.metadata['page']}"
            source_info += "]"
            
            context_parts.append(f"{source_info}\n{chunk.text}\n")
        
        context = "\n".join(context_parts)
        
        # Generate answer
        temp = self.temperature if temperature is None else temperature
        answer_text = self.llm.generate(query, context, temperature=temp)
        
        # Calculate confidence based on relevance scores
        confidence = min(1.0, len(relevant_chunks) / top_k * 0.8)
        
        # Create answer object
        answer = Answer(
            text=answer_text.strip(),
            sources=relevant_chunks[:top_k],
            confidence=confidence,
            query=query
        )
        
        return QueryResult(
            answer=answer,
            processing_time=time.time() - start_time,
            total_chunks_considered=len(relevant_chunks),
            model_used=self.llm.model
        )
    
    def list_documents(self) -> List[Dict[str, Any]]:
        """List all documents in the corpus."""
        docs = {}
        for chunk in self.chunks:
            file_path = chunk.metadata.get('file_path')
            if file_path not in docs:
                docs[file_path] = {
                    'file_name': chunk.metadata.get('file_name', 'Unknown'),
                    'file_path': file_path,
                    'type': chunk.metadata.get('type', 'unknown'),
                    'chunk_count': 0,
                    'processed_at': chunk.metadata.get('processed_at', '')
                }
            docs[file_path]['chunk_count'] += 1
        
        return list(docs.values())
    
    def get_stats(self) -> Dict[str, Any]:
        """Get corpus statistics."""
        return {
            'total_chunks': len(self.chunks),
            'total_documents': len(self.list_documents()),
            'embedding_model': self.embedding_model.model_name,
            'llm_model': self.llm.model if self.llm else None,
            'corpus_path': str(self.corpus_path)
        }
    
    def clear(self):
        """Clear the entire corpus."""
        confirmation = input("Are you sure you want to clear the entire corpus? (yes/no): ")
        if confirmation.lower() == 'yes':
            self.chunks = []
            self.collection.delete(where={})
            self._save_corpus()

# Convenience function
def ask(file_path: str, query: str, api_key: str = None, provider: str = "openai", llm_model: str = None) -> str:
    """One-liner function to ask a question about a document."""
    corpus = DocCorpus(
        corpus_path=f"./temp_corpus_{hashlib.md5(file_path.encode()).hexdigest()[:8]}"
    )
                    
    corpus.add_documents(file_path, update_existing=True)
    result = corpus.ask(query, llm_provider=provider, llm_model=llm_model, api_key=api_key)
                    
    # Cleanup
    try:
        shutil.rmtree(corpus.corpus_path)
    except:
        pass
                    
    return result.answer.text