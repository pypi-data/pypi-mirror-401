"""
DocNav Core Module

Core functionality for document management, parsing, chunking, embedding, and querying.
"""

import os
import sys
import json
import pickle
import hashlib
import tempfile
import warnings
import mimetypes
import shutil
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
import concurrent.futures

# Suppress warnings
warnings.filterwarnings('ignore')

# ============================================================================
# TQDM Progress Bar (with fallback)
# ============================================================================
try:
    from tqdm import tqdm
    HAS_TQDM = True
except ImportError:
    HAS_TQDM = False
    
    class tqdm:
        def __init__(self, iterable=None, total=None, desc="", unit="it", **kwargs):
            self.iterable = iterable
            self.total = total or (len(iterable) if iterable else None)
            self.desc = desc
            self.n = 0
            self.last_percent = 0
            
            if desc:
                print(f"\n{desc}...")
            if self.total:
                print(f"0/{self.total} (0%)")
        
        def __enter__(self):
            return self
        
        def __exit__(self, *args):
            if self.total:
                print(f"{self.total}/{self.total} (100%)")
        
        def __iter__(self):
            if self.iterable is None:
                return self
            
            for item in self.iterable:
                yield item
                self.update(1)
        
        def update(self, n=1):
            self.n += n
            if self.total:
                percent = (self.n * 100) // self.total
                if percent > self.last_percent:
                    print(f"{self.n}/{self.total} ({percent}%)")
                    self.last_percent = percent

# ============================================================================
# Colors for Terminal Output
# ============================================================================
class Colors:
    """ANSI color codes for terminal output."""
    RED = '\033[91m'
    GREEN = '\033[92m'
    YELLOW = '\033[93m'
    BLUE = '\033[94m'
    MAGENTA = '\033[95m'
    CYAN = '\033[96m'
    WHITE = '\033[97m'
    BOLD = '\033[1m'
    UNDERLINE = '\033[4m'
    END = '\033[0m'
    
    # Aliases for semantic naming
    ERROR = RED
    WARNING = YELLOW
    SUCCESS = GREEN
    INFO = CYAN
    PROGRESS = BLUE
    
    @staticmethod
    def success(text: str) -> str:
        return f"{Colors.SUCCESS}{text}{Colors.END}"
    
    @staticmethod
    def error(text: str) -> str:
        return f"{Colors.ERROR}{text}{Colors.END}"
    
    @staticmethod
    def warning(text: str) -> str:
        return f"{Colors.WARNING}{text}{Colors.END}"
    
    @staticmethod
    def info(text: str) -> str:
        return f"{Colors.INFO}{text}{Colors.END}"
    
    @staticmethod
    def progress(text: str) -> str:
        return f"{Colors.PROGRESS}{text}{Colors.END}"

# ============================================================================
# Data Structures
# ============================================================================
@dataclass
class DocumentChunk:
    """A chunk of document text with metadata."""
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    embedding: Optional[List[float]] = None

@dataclass
class Answer:
    """Answer from LLM with sources."""
    text: str
    sources: List[DocumentChunk]
    confidence: float
    query: str
    model_used: str
    processing_time: float
    
    def print(self, show_sources: bool = True, show_details: bool = False):
        """Print answer with optional sources and details."""
        print(f"\n{Colors.BOLD}{Colors.BLUE}ANSWER:{Colors.END}")
        print(self.text)
        
        if show_details:
            print(f"\n{Colors.BOLD}Details:{Colors.END}")
            print(f"  Confidence: {self.confidence:.1f}%")
            print(f"  Model: {self.model_used}")
            print(f"  Time: {self.processing_time:.2f}s")
            print(f"  Sources: {len(self.sources)}")
        
        if show_sources and self.sources:
            print(f"\n{Colors.BOLD}Sources:{Colors.END}")
            for i, source in enumerate(self.sources, 1):
                print(f"  {i}. {Colors.CYAN}{source.metadata.get('file_name', 'Unknown')}{Colors.END} [{source.metadata.get('type', 'unknown').upper()}]")
                print(f"     {source.text[:200]}...")

# ============================================================================
# Document Parser
# ============================================================================
class DocumentParser:
    """Parse various document formats."""
    
    @staticmethod
    def get_file_type(file_path: str) -> str:
        """Get file type from extension."""
        ext = Path(file_path).suffix.lower()
        type_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'docx',
            '.txt': 'text',
            '.md': 'text',
            '.csv': 'spreadsheet',
            '.xlsx': 'spreadsheet',
            '.xls': 'spreadsheet',
            '.pptx': 'pptx',
            '.ppt': 'pptx'
        }
        return type_map.get(ext, 'text')
    
    @staticmethod
    def parse_text_file(file_path: str) -> List[DocumentChunk]:
        """Parse plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            return [DocumentChunk(
                text=content,
                metadata={
                    'file_path': str(Path(file_path).resolve()),
                    'file_name': Path(file_path).name,
                    'type': 'text',
                    'size': len(content),
                    'processed_at': datetime.now().isoformat()
                }
            )]
        except Exception as e:
            print(Colors.error(f"Error reading text file {file_path}: {e}"))
            return []
    
    @staticmethod
    def parse_pdf(file_path: str, use_ocr: bool = False) -> List[DocumentChunk]:
        """Parse PDF files."""
        try:
            import PyPDF2
            from PyPDF2 import PdfReader
            
            chunks = []
            reader = PdfReader(file_path)
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        chunks.append(DocumentChunk(
                            text=text,
                            metadata={
                                'file_path': str(Path(file_path).resolve()),
                                'file_name': Path(file_path).name,
                                'type': 'pdf',
                                'page': page_num,
                                'total_pages': len(reader.pages),
                                'processed_at': datetime.now().isoformat()
                            }
                        ))
                except Exception as e:
                    print(Colors.warning(f"Error extracting page {page_num}: {e}"))
            
            return chunks
            
        except ImportError:
            print(Colors.warning("PyPDF2 not installed. Install with: pip install PyPDF2"))
            return []
        except Exception as e:
            print(Colors.error(f"Error parsing PDF {file_path}: {e}"))
            return []
    
    @staticmethod
    def parse_docx(file_path: str) -> List[DocumentChunk]:
        """Parse DOCX files."""
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            if full_text:
                return [DocumentChunk(
                    text='\n'.join(full_text),
                    metadata={
                        'file_path': str(Path(file_path).resolve()),
                        'file_name': Path(file_path).name,
                        'type': 'docx',
                        'processed_at': datetime.now().isoformat()
                    }
                )]
        
        except ImportError:
            print(Colors.warning("python-docx not installed. Install with: pip install python-docx"))
        except Exception as e:
            print(Colors.error(f"Error parsing DOCX {file_path}: {e}"))
        
        return DocumentParser.parse_text_file(file_path)  # Fallback
    
    @staticmethod
    def parse_spreadsheet(file_path: str) -> List[DocumentChunk]:
        """Parse Excel/CSV files."""
        try:
            import pandas as pd
            
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
                sheet_name = 'CSV'
            else:
                # Try to read all sheets
                xl = pd.ExcelFile(file_path)
                sheet_names = xl.sheet_names
                
                chunks = []
                for sheet in sheet_names:
                    df = xl.parse(sheet)
                    
                    # Convert to text representation
                    text = f"Sheet: {sheet}\n\n{df.to_string()}"
                    
                    chunks.append(DocumentChunk(
                        text=text,
                        metadata={
                            'file_path': str(Path(file_path).resolve()),
                            'file_name': Path(file_path).name,
                            'type': 'spreadsheet',
                            'sheet': sheet,
                            'processed_at': datetime.now().isoformat()
                        }
                    ))
                
                return chunks
        
        except ImportError:
            print(Colors.warning("pandas not installed. Install with: pip install pandas"))
        except Exception as e:
            print(Colors.error(f"Error parsing spreadsheet {file_path}: {e}"))
        
        return DocumentParser.parse_text_file(file_path)  # Fallback
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[DocumentChunk]:
        """Parse PowerPoint files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if slide_content:
                    slides_text.append(f"Slide {i+1}:\n" + "\n".join(slide_content))
            
            if slides_text:
                return [DocumentChunk(
                    text='\n\n'.join(slides_text),
                    metadata={
                        'file_path': str(Path(file_path).resolve()),
                        'file_name': Path(file_path).name,
                        'type': 'pptx',
                        'slides': len(slides_text),
                        'processed_at': datetime.now().isoformat()
                    }
                )]
        
        except ImportError:
            print(Colors.warning("python-pptx not installed. Install with: pip install python-pptx"))
        except Exception as e:
            print(Colors.error(f"Error parsing PPTX {file_path}: {e}"))
        
        return DocumentParser.parse_text_file(file_path)  # Fallback
    
    @classmethod
    def parse_file(cls, file_path: str, use_ocr: bool = False) -> List[DocumentChunk]:
        """Parse any supported file type."""
        file_path = str(Path(file_path).resolve())
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_type = cls.get_file_type(file_path)
        
        if file_type == 'pdf':
            return cls.parse_pdf(file_path, use_ocr)
        elif file_type == 'docx':
            return cls.parse_docx(file_path)
        elif file_type == 'text':
            return cls.parse_text_file(file_path)
        elif file_type == 'spreadsheet':
            return cls.parse_spreadsheet(file_path)
        elif file_type == 'pptx':
            return cls.parse_pptx(file_path)
        else:
            # Try as text file
            try:
                return cls.parse_text_file(file_path)
            except:
                print(Colors.warning(f"Unsupported file type: {file_path}"))
                return []

# ============================================================================
# Smart Chunking
# ============================================================================
class SmartChunker:
    """Intelligent document chunking that preserves structure."""
    
    def __init__(self, chunk_size: int = 1000, overlap: int = 100):
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def chunk_text(self, text: str, metadata: Dict[str, Any]) -> List[DocumentChunk]:
        """Split text into coherent chunks."""
        # Simple paragraph-based chunking
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        if not paragraphs:
            return []
        
        chunks = []
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_words = len(para.split())
            
            if current_size + para_words > self.chunk_size and current_chunk:
                # Save current chunk
                chunk_text = '\n\n'.join(current_chunk)
                chunk_meta = metadata.copy()
                chunk_meta['chunk_index'] = len(chunks)
                
                chunks.append(DocumentChunk(
                    text=chunk_text,
                    metadata=chunk_meta
                ))
                
                # Start new chunk with overlap
                overlap_paras = current_chunk[-2:] if len(current_chunk) > 2 else []
                current_chunk = overlap_paras + [para]
                current_size = sum(len(p.split()) for p in current_chunk)
            else:
                current_chunk.append(para)
                current_size += para_words
        
        # Add the last chunk
        if current_chunk:
            chunk_text = '\n\n'.join(current_chunk)
            chunk_meta = metadata.copy()
            chunk_meta['chunk_index'] = len(chunks)
            
            chunks.append(DocumentChunk(
                text=chunk_text,
                metadata=chunk_meta
            ))
        
        return chunks

# ============================================================================
# Embedding Models
# ============================================================================
class EmbeddingModel:
    """Handle text embeddings with multiple backends."""
    
    def __init__(self, model_name: str = "all-MiniLM-L6-v2"):
        self.model_name = model_name
        self._model = None
        
        # Force fallback mode for now to avoid hanging import
        print(Colors.warning("Using simple embeddings (fallback mode)."))
        print(Colors.info("To use sentence-transformers: pip install sentence-transformers"))
    
    def embed(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings for texts."""
        if self._model:
            return self._model.encode(texts, show_progress_bar=False).tolist()
        else:
            # Fallback: simple deterministic embedding
            import re
            from collections import Counter
            
            embeddings = []
            for text in texts:
                # Simple bag-of-words style embedding
                words = re.findall(r'\w+', text.lower())
                word_counts = Counter(words)
                
                # Create a deterministic embedding based on word hashes
                embedding = []
                for word in sorted(set(words)):
                    hash_val = int(hashlib.md5(word.encode()).hexdigest()[:8], 16)
                    normalized = (hash_val % 10000) / 10000.0
                    embedding.append(normalized)
                
                # Pad to 384 dimensions
                while len(embedding) < 384:
                    embedding.append(0.0)
                embedding = embedding[:384]
                
                embeddings.append(embedding)
            
            return embeddings
    
    def embed_single(self, text: str) -> List[float]:
        """Embed a single text."""
        return self.embed([text])[0]

# ============================================================================
# LLM Providers
# ============================================================================
class LLMProvider:
    """Handle LLM interactions with multiple providers."""
    
    def __init__(self, provider: str = "openai", model: str = None, api_key: Optional[str] = None):
        self.provider = provider.lower()
        self.api_key = api_key or self._get_api_key()
        
        # Set default models
        if not model:
            if self.provider == "openai":
                model = "gpt-3.5-turbo"
            elif self.provider == "gemini":
                model = "gemini-2.5-flash"
            elif self.provider == "claude":
                model = "claude-3-haiku-20240307"
            else:
                model = "gpt-3.5-turbo"
        
        self.model = model
        
        if not self.api_key:
            print(Colors.warning(f"No API key provided for {self.provider}. Please set it for LLM queries."))
    
    def _get_api_key(self) -> Optional[str]:
        """Get API key from environment variables."""
        env_vars = {
            "openai": "OPENAI_API_KEY",
            "gemini": "GOOGLE_API_KEY",
            "anthropic": "ANTHROPIC_API_KEY",
            "claude": "ANTHROPIC_API_KEY"
        }
        
        env_var = env_vars.get(self.provider)
        if env_var:
            return os.getenv(env_var)
        return None
    
    def generate(self, prompt: str, context: str = "", system_prompt: Optional[str] = None) -> str:
        """Generate a response from the LLM."""
        if not self.api_key:
            return "API key not provided. Please set your API key."
        
        full_prompt = f"""Based on the following context, answer the question accurately.
If the answer cannot be found in the context, say "I cannot find this information in the provided documents."

Context from documents:
{context}

Question: {prompt}

Provide a clear, accurate answer based only on the context above. Cite sources when possible."""

        if system_prompt:
            full_prompt = f"{system_prompt}\n\n{full_prompt}"
        
        try:
            if self.provider == "openai":
                try:
                    import openai
                    openai.api_key = self.api_key
                    
                    response = openai.chat.completions.create(
                        model=self.model,
                        messages=[
                            {"role": "system", "content": "You are a helpful assistant that answers questions based on provided documents."},
                            {"role": "user", "content": full_prompt}
                        ],
                        temperature=0.0,
                        max_tokens=4096
                    )
                    return response.choices[0].message.content
                except Exception as e:
                    return f"OpenAI Error: {str(e)}"
            
            elif self.provider == "gemini":
                try:
                    import google.generativeai as genai
                    genai.configure(api_key=self.api_key)
                    
                    model = genai.GenerativeModel(self.model)
                    response = model.generate_content(full_prompt)
                    return response.text
                except Exception as e:
                    return f"Gemini Error: {str(e)}"
            
            elif self.provider in ["anthropic", "claude"]:
                try:
                    import anthropic
                    client = anthropic.Anthropic(api_key=self.api_key)
                    
                    response = client.messages.create(
                        model=self.model,
                        max_tokens=65536,
                        temperature=0.0,
                        messages=[
                            {"role": "user", "content": full_prompt}
                        ]
                    )
                    return response.content[0].text
                except Exception as e:
                    return f"Claude Error: {str(e)}"
            
            else:
                return f"Unsupported provider: {self.provider}"
        
        except Exception as e:
            return f"Error calling {self.provider}: {str(e)}"

# ============================================================================
# Main Corpus Class
# ============================================================================
class Corpus:
    """Main document corpus for querying and management."""
    
    def __init__(self, 
                 name: str,
                 base_path: Optional[str] = None,
                 embedding_model: str = "all-MiniLM-L6-v2",
                 chunk_size: int = 1000):
        """Initialize a document corpus.
        
        Args:
            name: Name of the corpus
            base_path: Base directory to store corpora (default: ~/.docnav)
            embedding_model: Name of embedding model
            chunk_size: Size of document chunks
        """
        # Set up corpus directory
        if base_path:
            self.base_path = Path(base_path)
        else:
            self.base_path = Path.home() / ".docnav"
        
        self.corpus_path = self.base_path / "corpora" / name
        self.corpus_path.mkdir(parents=True, exist_ok=True)
        
        self.name = name
        self.embedding_model_name = embedding_model
        self.chunk_size = chunk_size
        
        # Initialize components
        self.embedder = EmbeddingModel(model_name=embedding_model)
        self.chunker = SmartChunker(chunk_size=chunk_size)
        self.parser = DocumentParser()
        
        # Load existing corpus
        self.chunks: List[DocumentChunk] = []
        self._load_corpus()
        
        print(Colors.success(f"Corpus '{name}' loaded"))
        print(f"   Path: {self.corpus_path}")
        print(f"   Documents: {len(set(chunk.metadata.get('file_path') for chunk in self.chunks))}")
        print(f"   Chunks: {len(self.chunks)}")
    
    def _load_corpus(self):
        """Load corpus from disk."""
        index_file = self.corpus_path / "corpus_index.pkl"
        if index_file.exists():
            try:
                with open(index_file, 'rb') as f:
                    self.chunks = pickle.load(f)
                print(Colors.info(f"Loaded {len(self.chunks)} chunks from existing corpus"))
            except Exception as e:
                print(Colors.warning(f"Failed to load corpus index: {e}"))
                self.chunks = []
        else:
            self.chunks = []
            print(Colors.info("Starting new corpus"))
    
    def _save_corpus(self):
        """Save corpus to disk."""
        index_file = self.corpus_path / "corpus_index.pkl"
        try:
            with open(index_file, 'wb') as f:
                pickle.dump(self.chunks, f)
        except Exception as e:
            print(Colors.error(f"Failed to save corpus: {e}"))
    
    def add(self, 
           paths: List[str], 
           use_ocr: bool = False, 
           update_existing: bool = False,
           chunk_size: Optional[int] = None) -> int:
        """Add documents to corpus."""
        if chunk_size:
            self.chunker = SmartChunker(chunk_size=chunk_size)
        
        # Collect all files
        all_files = []
        for path in paths:
            path_obj = Path(path)
            if path_obj.is_file():
                all_files.append(str(path_obj))
            elif path_obj.is_dir():
                for ext in ['*.pdf', '*.docx', '*.doc', '*.txt', '*.md', '*.csv', '*.xlsx', '*.xls', '*.pptx', '*.ppt']:
                    all_files.extend(str(f) for f in path_obj.glob(ext))
        
        if not all_files:
            print(Colors.warning("No files found to process"))
            return 0
        
        print(Colors.info(f"Found {len(all_files)} files to process"))
        
        # Process files
        new_chunks = []
        added_count = 0
        
        with tqdm(total=len(all_files), desc="Processing files") as pbar:
            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                future_to_file = {
                    executor.submit(self.parser.parse_file, file_path, use_ocr): file_path 
                    for file_path in all_files
                }
                
                for future in concurrent.futures.as_completed(future_to_file):
                    file_path = future_to_file[future]
                    try:
                        chunks = future.result()
                        if chunks:
                            for chunk in chunks:
                                # Check if already exists
                                file_hash = hashlib.md5(
                                    f"{chunk.metadata.get('file_path')}{chunk.text}".encode()
                                ).hexdigest()
                                
                                existing_hashes = set(
                                    hashlib.md5(f"{c.metadata.get('file_path')}{c.text}".encode()).hexdigest()
                                    for c in self.chunks
                                )
                                
                                if file_hash not in existing_hashes or update_existing:
                                    new_chunks.append(chunk)
                                    added_count += 1
                    except Exception as e:
                        print(Colors.warning(f"Error processing {file_path}: {e}"))
                    
                    pbar.update(1)
        
        # Generate embeddings for new chunks
        if new_chunks:
            print(Colors.info("Generating embeddings..."))
            texts = [chunk.text for chunk in new_chunks]
            embeddings = self.embedder.embed(texts)
            
            for chunk, embedding in zip(new_chunks, embeddings):
                chunk.embedding = embedding
            
            # Add to corpus
            self.chunks.extend(new_chunks)
            self._save_corpus()
        
        print(Colors.success(f"Added {added_count} new chunks"))
        print(f"   Processed: {len(all_files)} files")
        
        return len(new_chunks)
    
    def ask(self, 
            query: str, 
            top_k: int = 5,
            where: Optional[Dict[str, Any]] = None,
            include_sources: bool = True,
            llm_provider: str = "openai",
            llm_model: Optional[str] = None,
            api_key: Optional[str] = None,
            temperature: float = 0.0) -> Answer:
        """Ask a question to corpus."""
        import time
        start_time = time.time()
        
        if not self.chunks:
            return Answer(
                text="No documents in corpus. Please add documents first.",
                sources=[],
                confidence=0.0,
                query=query,
                model_used=llm_model or "none",
                processing_time=time.time() - start_time
            )
        
        # Generate query embedding
        query_embedding = self.embedder.embed_single(query)
        
        # Find similar chunks (simple cosine similarity)
        scored_chunks = []
        for chunk in self.chunks:
            if chunk.embedding:
                # Simple dot product similarity
                import math
                chunk_vec = chunk.embedding
                dot_product = sum(q * c for q, c in zip(query_embedding, chunk_vec))
                norm_q = math.sqrt(sum(q * q for q in query_embedding))
                norm_c = math.sqrt(sum(c * c for c in chunk_vec))
                similarity = dot_product / (norm_q * norm_c + 1e-8)
                
                # Apply filters
                if where:
                    matches = True
                    for key, value in where.items():
                        if chunk.metadata.get(key) != value:
                            matches = False
                            break
                    if not matches:
                        continue
                
                scored_chunks.append((similarity, chunk))
        
        # Sort by similarity
        scored_chunks.sort(key=lambda x: x[0], reverse=True)
        relevant_chunks = [chunk for _, chunk in scored_chunks[:top_k]]
        
        if not relevant_chunks:
            return Answer(
                text="I couldn't find any relevant information in documents to answer this question.",
                sources=[],
                confidence=0.0,
                query=query,
                model_used=llm_model or "none",
                processing_time=time.time() - start_time
            )
        
        # Prepare context
        context_parts = []
        for i, chunk in enumerate(relevant_chunks):
            source_info = f"[Source {i+1}: {chunk.metadata.get('file_name', 'Unknown')}"
            if 'page' in chunk.metadata:
                source_info += f", page {chunk.metadata['page']}"
            source_info += "]"
            
            context_parts.append(f"{source_info}\n{chunk.text}\n")
        
        context = "\n".join(context_parts)
        
        # Create LLM provider
        llm = LLMProvider(
            provider=llm_provider,
            model=llm_model,
            api_key=api_key
        )
        
        # Generate answer
        answer_text = llm.generate(query, context)
        
        # Calculate confidence
        confidence = min(1.0, len(relevant_chunks) / max(top_k, 1) * 0.9)
        
        # Create answer
        processing_time = time.time() - start_time
        
        answer = Answer(
            text=answer_text.strip(),
            sources=relevant_chunks if include_sources else [],
            confidence=confidence,
            query=query,
            model_used=llm.model,
            processing_time=processing_time
        )
        
        return answer
    
    def list(self, details: bool = False) -> List[Dict[str, Any]]:
        """List all documents in corpus."""
        docs = {}
        for chunk in self.chunks:
            file_path = chunk.metadata.get('file_path')
            if file_path not in docs:
                docs[file_path] = {
                    'file_name': chunk.metadata.get('file_name', 'Unknown'),
                    'file_path': file_path,
                    'type': chunk.metadata.get('type', 'unknown'),
                    'chunks': 0,
                    'size': 0,
                    'processed_at': chunk.metadata.get('processed_at')
                }
            docs[file_path]['chunks'] += 1
            docs[file_path]['size'] += len(chunk.text)
        
        return list(docs.values())
    
    def stats(self) -> Dict[str, Any]:
        """Get corpus statistics."""
        return {
            'name': self.name,
            'total_chunks': len(self.chunks),
            'total_documents': len(set(chunk.metadata.get('file_path') for chunk in self.chunks)),
            'total_size_mb': sum(len(chunk.text) for chunk in self.chunks) / (1024 * 1024),
            'embedding_model': self.embedding_model_name,
            'chunk_size': self.chunk_size,
            'corpus_path': str(self.corpus_path)
        }
    
    def remove(self, file_path: str) -> int:
        """Remove document from corpus."""
        file_path = str(Path(file_path).resolve())
        original_count = len(self.chunks)
        
        self.chunks = [
            chunk for chunk in self.chunks 
            if chunk.metadata.get('file_path') != file_path
        ]
        
        removed_count = original_count - len(self.chunks)
        if removed_count > 0:
            self._save_corpus()
            print(Colors.success(f"Removed {removed_count} chunks"))
        else:
            print(Colors.warning("No matching documents found"))
        
        return removed_count
    
    def clear(self, confirm: bool = True) -> bool:
        """Clear entire corpus."""
        if confirm:
            response = input(f"Are you sure you want to clear corpus '{self.name}'? [y/N]: ")
            if response.lower() != 'y':
                return False
        
        self.chunks = []
        self._save_corpus()
        return True

def load_corpus(name: str, base_path: Optional[str] = None) -> Corpus:
    """Load a corpus by name."""
    try:
        return Corpus(name=name, base_path=base_path)
    except Exception as e:
        raise FileNotFoundError(f"Corpus '{name}' not found. Create it first with 'docnav new {name}'")
