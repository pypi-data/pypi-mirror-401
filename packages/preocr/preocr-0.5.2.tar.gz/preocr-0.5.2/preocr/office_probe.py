"""Office document text extraction (DOCX, PPTX, XLSX)."""

from pathlib import Path
from typing import Any, Dict, Optional

from .exceptions import OfficeDocumentError, TextExtractionError
from .logger import get_logger

logger = get_logger(__name__)

# Declare these as Optional[Any] so mypy knows they can be None
Document: Optional[Any]
Presentation: Optional[Any]
load_workbook: Optional[Any]

try:
    from docx import Document as _Document

    Document = _Document
except ImportError:
    Document = None

try:
    from pptx import Presentation as _Presentation

    Presentation = _Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook as _load_workbook

    load_workbook = _load_workbook
except ImportError:
    load_workbook = None


def extract_office_text(file_path: str, mime_type: str) -> Dict[str, Any]:
    """
    Extract text from Office documents (DOCX, PPTX, XLSX).

    Args:
        file_path: Path to the office document
        mime_type: MIME type of the file

    Returns:
        Dictionary with keys:
            - text_length: Number of characters in extracted text
            - text: Extracted text (may be truncated for large files)
            - document_type: Type of document ("docx", "pptx", "xlsx")
    """
    path = Path(file_path)

    if "wordprocessingml" in mime_type or path.suffix.lower() == ".docx":
        return _extract_docx(path)
    elif "presentationml" in mime_type or path.suffix.lower() == ".pptx":
        return _extract_pptx(path)
    elif "spreadsheetml" in mime_type or path.suffix.lower() == ".xlsx":
        return _extract_xlsx(path)
    else:
        return {
            "text_length": 0,
            "text": "",
            "document_type": None,
        }


def _extract_docx(path: Path) -> Dict[str, Any]:
    """Extract text from DOCX file."""
    if Document is None:
        return {
            "text_length": 0,
            "text": "",
            "document_type": "docx",
        }

    try:
        doc = Document(str(path))
        text_parts = []

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)

        full_text = "\n".join(text_parts)

        return {
            "text_length": len(full_text),
            "text": full_text[:1000] if len(full_text) > 1000 else full_text,
            "document_type": "docx",
        }
    except (IOError, OSError, PermissionError) as e:
        logger.warning(f"Failed to read DOCX file: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "docx",
        }
    except Exception as e:
        logger.warning(f"DOCX text extraction failed: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "docx",
        }


def _extract_pptx(path: Path) -> Dict[str, Any]:
    """Extract text from PPTX file."""
    if Presentation is None:
        return {
            "text_length": 0,
            "text": "",
            "document_type": "pptx",
        }

    try:
        prs = Presentation(str(path))
        text_parts = []

        # Extract text from slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

        full_text = "\n".join(text_parts)

        return {
            "text_length": len(full_text),
            "text": full_text[:1000] if len(full_text) > 1000 else full_text,
            "document_type": "pptx",
        }
    except (IOError, OSError, PermissionError) as e:
        logger.warning(f"Failed to read PPTX file: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "pptx",
        }
    except Exception as e:
        logger.warning(f"PPTX text extraction failed: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "pptx",
        }


def _extract_xlsx(path: Path) -> Dict[str, Any]:
    """Extract text from XLSX file."""
    if not load_workbook:
        return {
            "text_length": 0,
            "text": "",
            "document_type": "xlsx",
        }

    try:
        wb = load_workbook(path, data_only=True)
        text_parts = []

        # Extract text from all sheets
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)

        wb.close()
        full_text = "\n".join(text_parts)

        return {
            "text_length": len(full_text),
            "text": full_text[:1000] if len(full_text) > 1000 else full_text,
            "document_type": "xlsx",
        }
    except (IOError, OSError, PermissionError) as e:
        logger.warning(f"Failed to read XLSX file: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "xlsx",
        }
    except Exception as e:
        logger.warning(f"XLSX text extraction failed: {e}")
        return {
            "text_length": 0,
            "text": "",
            "document_type": "xlsx",
        }
