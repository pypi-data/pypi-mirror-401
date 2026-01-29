# -*- coding: utf-8 -*-
"""
Custom Report Builder Panel
===========================
Manage report queue items and generate custom reports with selected plots and tables.
"""

import tkinter as tk
from tkinter import ttk, messagebox, filedialog, simpledialog
import threading
from typing import Dict, List, Optional, Any
from pathlib import Path
import os
import io
from datetime import datetime

from biblium.gui.config import FONTS, LAYOUT, get_theme
from biblium.gui.core.events import event_bus, EventBus
from biblium.gui.core.state import report_queue, ReportItem
from biblium.gui.panels.base import BasePanel
from biblium.gui.widgets.cards import Card, CollapsibleCard
from biblium.gui.widgets.buttons import ThemedButton, ActionButton


class CustomReportPanel(BasePanel):
    """Panel for building custom reports from collected plots and tables."""
    
    title = "Custom Report Builder"
    icon = "🎨"
    description = "Build custom reports from collected plots and tables"
    requires_data = False  # Can work without data loaded
    
    # Info tab content
    info_content = """
## Custom Report Builder

Build personalized reports by collecting plots and tables from various panels.

### How It Works

1. **Collect Items**: Right-click on any plot and select "Add to Report", or use the 
   "📄 Add to Report" button that appears below plots in many panels
2. **Manage Queue**: View, reorder, edit, and remove items from your collection
3. **Generate Report**: Export your custom report to DOCX, PPTX, or PDF

### Queue Management

| Action | Description |
|--------|-------------|
| **⬆️ Move Up** | Move selected item up in order |
| **⬇️ Move Down** | Move selected item down in order |
| **✏️ Edit Title** | Change the title of selected item |
| **🗑️ Remove** | Remove selected item from queue |
| **🔄 Clear All** | Remove all items from queue |

### Output Formats

- **DOCX (Word)**: Each item as a section with title and content
- **PPTX (PowerPoint)**: Each item as a slide
- **PDF**: Single document with all items

### Preview

- Double-click an item in the queue to preview it
- Plots show as images, tables show data preview

### Tips

1. **Order Matters**: Items appear in report in queue order
2. **Edit Titles**: Give items meaningful titles for better reports
3. **Add Descriptions**: Include captions that explain the content
4. **Mix Content**: Combine plots and tables in one report
5. **Collect Progressively**: Build your report as you explore data

### Methods Used

The custom report builder uses the global `report_queue` from the GUI core:

```python
from biblium.gui.core.state import report_queue

# Add items programmatically
report_queue.add_plot(figure, "My Plot Title", "Source Panel")
report_queue.add_table(dataframe, "My Table Title", "Source Panel")

# Manage queue
report_queue.move_up(index)
report_queue.move_down(index)
report_queue.remove_by_index(index)
report_queue.clear()
```
"""
    
    def __init__(self, parent, bib=None, theme: str = "light", **kwargs):
        super().__init__(parent, bib=bib, theme=theme, **kwargs)
        
        # Register for queue updates
        report_queue.on_change(self._on_queue_change)
        
        # Bind cleanup on destroy
        self.bind("<Destroy>", self._on_panel_destroy)
    
    def _on_panel_destroy(self, event):
        """Clean up when panel is destroyed."""
        if event.widget == self:
            try:
                report_queue.off_change(self._on_queue_change)
            except:
                pass
    
    def _on_queue_change(self, action: str, item: Optional[ReportItem], items: List[ReportItem]):
        """Handle queue changes."""
        # Update the list view
        self.after(0, self._refresh_queue_list)
    
    def _create_options(self):
        """Create the options panel."""
        self._add_title()
        
        # Queue Status Card
        status_card = Card(self.options_content, title="📊 Queue Status", theme=self.theme_name)
        status_card.pack(fill=tk.X, padx=8, pady=8)
        
        self.status_label = tk.Label(
            status_card.content,
            text="No items in queue",
            font=FONTS.get_font("body"),
            bg=self.theme["bg_card"], fg=self.theme["text_secondary"],
            anchor=tk.W,
        )
        self.status_label.pack(fill=tk.X, pady=4)
        
        # Stats frame
        self.stats_frame = tk.Frame(status_card.content, bg=self.theme["bg_card"])
        self.stats_frame.pack(fill=tk.X, pady=4)
        
        self.plots_label = tk.Label(
            self.stats_frame,
            text="📈 Plots: 0",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_card"], fg=self.theme["text_muted"],
        )
        self.plots_label.pack(side=tk.LEFT, padx=(0, 16))
        
        self.tables_label = tk.Label(
            self.stats_frame,
            text="📋 Tables: 0",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_card"], fg=self.theme["text_muted"],
        )
        self.tables_label.pack(side=tk.LEFT)
        
        # Queue Management Card
        mgmt_card = Card(self.options_content, title="🔧 Queue Management", theme=self.theme_name)
        mgmt_card.pack(fill=tk.X, padx=8, pady=8)
        
        # Buttons grid
        btn_frame1 = tk.Frame(mgmt_card.content, bg=self.theme["bg_card"])
        btn_frame1.pack(fill=tk.X, pady=4)
        
        ThemedButton(
            btn_frame1, text="⬆️ Move Up", style="ghost", size="small",
            command=self._move_up, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        ThemedButton(
            btn_frame1, text="⬇️ Move Down", style="ghost", size="small",
            command=self._move_down, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        btn_frame2 = tk.Frame(mgmt_card.content, bg=self.theme["bg_card"])
        btn_frame2.pack(fill=tk.X, pady=4)
        
        ThemedButton(
            btn_frame2, text="✏️ Edit Title", style="ghost", size="small",
            command=self._edit_title, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        ThemedButton(
            btn_frame2, text="📝 Add Description", style="ghost", size="small",
            command=self._edit_description, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        btn_frame3 = tk.Frame(mgmt_card.content, bg=self.theme["bg_card"])
        btn_frame3.pack(fill=tk.X, pady=4)
        
        ThemedButton(
            btn_frame3, text="🗑️ Remove Selected", style="danger", size="small",
            command=self._remove_selected, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        ThemedButton(
            btn_frame3, text="🔄 Clear All", style="danger", size="small",
            command=self._clear_all, theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=2, expand=True, fill=tk.X)
        
        # Output Format Card
        format_card = Card(self.options_content, title="📄 Output Formats", theme=self.theme_name)
        format_card.pack(fill=tk.X, padx=8, pady=8)
        
        tk.Label(
            format_card.content,
            text="Select one or more output formats:",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_card"], fg=self.theme["text_secondary"],
            anchor=tk.W,
        ).pack(fill=tk.X, pady=(0, 4))
        
        self.format_vars = {}
        formats_frame = tk.Frame(format_card.content, bg=self.theme["bg_card"])
        formats_frame.pack(fill=tk.X, pady=4)
        
        # Create checkboxes for each format
        self.format_vars["docx"] = tk.BooleanVar(value=True)
        self.format_vars["pptx"] = tk.BooleanVar(value=False)
        self.format_vars["pdf"] = tk.BooleanVar(value=False)
        
        for fmt, label in [("docx", "Word (DOCX)"), ("pptx", "PowerPoint (PPTX)"), ("pdf", "PDF")]:
            cb = tk.Checkbutton(
                formats_frame, text=label, variable=self.format_vars[fmt],
                font=FONTS.get_font("body"),
                bg=self.theme["bg_card"], fg=self.theme["text_primary"],
                selectcolor=self.theme["bg_secondary"],
                activebackground=self.theme["bg_card"],
            )
            cb.pack(anchor=tk.W, padx=8)
        
        # Report title entry
        title_frame = tk.Frame(format_card.content, bg=self.theme["bg_card"])
        title_frame.pack(fill=tk.X, pady=(8, 4))
        
        tk.Label(
            title_frame, text="Report Title:",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_card"], fg=self.theme["text_secondary"],
        ).pack(side=tk.LEFT, padx=(0, 8))
        
        self.report_title_var = tk.StringVar(value="Custom Report")
        self.report_title_entry = tk.Entry(
            title_frame,
            textvariable=self.report_title_var,
            font=FONTS.get_font("body"),
            bg=self.theme["bg_secondary"],
            fg=self.theme["text_primary"],
            relief=tk.FLAT,
        )
        self.report_title_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, ipady=4)
        
        # Action Buttons
        btn_frame = tk.Frame(self.options_content, bg=self.theme["bg_secondary"])
        btn_frame.pack(fill=tk.X, padx=8, pady=16)
        
        ActionButton(
            btn_frame, text="Generate Custom Report", icon="📝",
            command=self._generate_report, theme=self.theme_name,
        ).pack(fill=tk.X)
        
        ThemedButton(
            btn_frame, text="Preview Selected Item", style="secondary",
            command=self._preview_selected, theme=self.theme_name,
        ).pack(fill=tk.X, pady=(8, 0))
    
    def _create_results(self):
        """Create the results panel."""
        self.results_card = tk.Frame(
            self.results_frame, bg=self.theme["bg_card"],
            highlightbackground=self.theme["border"], highlightthickness=1,
        )
        self.results_card.pack(fill=tk.BOTH, expand=True, padx=4, pady=4)
        
        # Header
        header = tk.Frame(self.results_card, bg=self.theme["bg_card"])
        header.pack(fill=tk.X)
        
        tk.Label(
            header, text="Report Queue",
            font=FONTS.get_font("heading", bold=True),
            bg=self.theme["bg_card"], fg=self.theme["text_primary"],
            padx=12, pady=8,
        ).pack(side=tk.LEFT)
        
        # Refresh button
        ThemedButton(
            header, text="🔄", style="ghost", size="small",
            command=self._refresh_queue_list, theme=self.theme_name,
        ).pack(side=tk.RIGHT, padx=8)
        
        ttk.Separator(self.results_card, orient=tk.HORIZONTAL).pack(fill=tk.X, padx=12)
        
        # Create queue list with scrollbar
        list_container = tk.Frame(self.results_card, bg=self.theme["bg_card"])
        list_container.pack(fill=tk.BOTH, expand=True, padx=12, pady=12)
        
        # Scrollbar
        scrollbar = ttk.Scrollbar(list_container, orient=tk.VERTICAL)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Listbox for queue items
        self.queue_listbox = tk.Listbox(
            list_container,
            font=FONTS.get_font("body"),
            bg=self.theme["bg_secondary"],
            fg=self.theme["text_primary"],
            selectmode=tk.SINGLE,
            selectbackground=self.theme["accent_primary"],
            selectforeground="white",
            activestyle='none',
            yscrollcommand=scrollbar.set,
            height=15,
        )
        self.queue_listbox.pack(fill=tk.BOTH, expand=True)
        scrollbar.config(command=self.queue_listbox.yview)
        
        # Bind double-click for preview
        self.queue_listbox.bind("<Double-1>", lambda e: self._preview_selected())
        
        # Preview frame below list
        ttk.Separator(self.results_card, orient=tk.HORIZONTAL).pack(fill=tk.X, padx=12, pady=(0, 8))
        
        preview_header = tk.Frame(self.results_card, bg=self.theme["bg_card"])
        preview_header.pack(fill=tk.X, padx=12)
        
        tk.Label(
            preview_header, text="Preview",
            font=FONTS.get_font("body", bold=True),
            bg=self.theme["bg_card"], fg=self.theme["text_primary"],
        ).pack(side=tk.LEFT)
        
        self.preview_container = tk.Frame(self.results_card, bg=self.theme["bg_secondary"])
        self.preview_container.pack(fill=tk.BOTH, expand=True, padx=12, pady=(4, 12))
        
        self.preview_label = tk.Label(
            self.preview_container,
            text="Select an item and click 'Preview' or double-click to preview",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
            wraplength=350,
        )
        self.preview_label.pack(pady=30)
        
        # Initial load
        self._refresh_queue_list()
    
    def _refresh_queue_list(self):
        """Refresh the queue listbox."""
        try:
            self.queue_listbox.delete(0, tk.END)
        except tk.TclError:
            return  # Widget destroyed
        
        items = report_queue.get_items()
        
        for i, item in enumerate(items):
            icon = "📈" if item.item_type == "plot" else "📋"
            source = item.source_panel[:20] + "..." if len(item.source_panel) > 20 else item.source_panel
            self.queue_listbox.insert(tk.END, f"{i+1}. {icon} {item.title} ({source})")
        
        # Update status
        total = len(items)
        plots = len([i for i in items if i.item_type == "plot"])
        tables = len([i for i in items if i.item_type == "table"])
        
        if total == 0:
            self.status_label.config(text="No items in queue")
        else:
            self.status_label.config(text=f"Total items: {total}")
        
        self.plots_label.config(text=f"📈 Plots: {plots}")
        self.tables_label.config(text=f"📋 Tables: {tables}")
    
    def _get_selected_index(self) -> Optional[int]:
        """Get currently selected item index."""
        selection = self.queue_listbox.curselection()
        if not selection:
            return None
        return selection[0]
    
    def _move_up(self):
        """Move selected item up."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to move.")
            return
        
        if report_queue.move_up(idx):
            self._refresh_queue_list()
            # Reselect the moved item
            self.queue_listbox.selection_set(idx - 1)
    
    def _move_down(self):
        """Move selected item down."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to move.")
            return
        
        if report_queue.move_down(idx):
            self._refresh_queue_list()
            # Reselect the moved item
            self.queue_listbox.selection_set(idx + 1)
    
    def _edit_title(self):
        """Edit the title of selected item."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to edit.")
            return
        
        items = report_queue.get_items()
        if idx >= len(items):
            return
        
        item = items[idx]
        new_title = simpledialog.askstring(
            "Edit Title",
            "Enter new title:",
            initialvalue=item.title,
            parent=self,
        )
        
        if new_title and new_title.strip():
            item.title = new_title.strip()
            self._refresh_queue_list()
    
    def _edit_description(self):
        """Edit the description of selected item."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to edit.")
            return
        
        items = report_queue.get_items()
        if idx >= len(items):
            return
        
        item = items[idx]
        new_desc = simpledialog.askstring(
            "Edit Description",
            "Enter description/caption:",
            initialvalue=item.description or "",
            parent=self,
        )
        
        if new_desc is not None:  # Allow empty string to clear description
            item.description = new_desc.strip()
            self._refresh_queue_list()
    
    def _remove_selected(self):
        """Remove selected item from queue."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to remove.")
            return
        
        items = report_queue.get_items()
        if idx >= len(items):
            return
        
        item = items[idx]
        if messagebox.askyesno("Confirm Remove", f"Remove '{item.title}' from queue?"):
            report_queue.remove_by_index(idx)
            self._refresh_queue_list()
    
    def _clear_all(self):
        """Clear all items from queue."""
        if len(report_queue) == 0:
            messagebox.showinfo("Empty Queue", "Queue is already empty.")
            return
        
        if messagebox.askyesno("Confirm Clear", f"Remove all {len(report_queue)} items from the queue?"):
            report_queue.clear()
            self._refresh_queue_list()
            self._clear_preview()
    
    def _preview_selected(self):
        """Preview the selected item."""
        idx = self._get_selected_index()
        if idx is None:
            messagebox.showinfo("No Selection", "Please select an item to preview.")
            return
        
        items = report_queue.get_items()
        if idx >= len(items):
            return
        
        item = items[idx]
        self._show_preview(item)
    
    def _clear_preview(self):
        """Clear the preview area."""
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        self.preview_label = tk.Label(
            self.preview_container,
            text="Select an item and click 'Preview' or double-click to preview",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
            wraplength=350,
        )
        self.preview_label.pack(pady=30)
    
    def _show_preview(self, item: ReportItem):
        """Show preview of an item."""
        # Clear previous preview
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        # Title
        tk.Label(
            self.preview_container,
            text=item.title,
            font=FONTS.get_font("body", bold=True),
            bg=self.theme["bg_secondary"], fg=self.theme["text_primary"],
            anchor=tk.W,
        ).pack(fill=tk.X, padx=8, pady=(8, 2))
        
        # Type and source
        tk.Label(
            self.preview_container,
            text=f"Type: {item.item_type.title()} | Source: {item.source_panel}",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
            anchor=tk.W,
        ).pack(fill=tk.X, padx=8, pady=2)
        
        # Description if present
        if item.description:
            tk.Label(
                self.preview_container,
                text=f"Description: {item.description}",
                font=FONTS.get_font("small"),
                bg=self.theme["bg_secondary"], fg=self.theme["text_secondary"],
                anchor=tk.W,
                wraplength=350,
            ).pack(fill=tk.X, padx=8, pady=2)
        
        ttk.Separator(self.preview_container, orient=tk.HORIZONTAL).pack(fill=tk.X, padx=8, pady=8)
        
        if item.item_type == "plot":
            self._show_plot_preview(item)
        else:
            self._show_table_preview(item)
    
    def _show_plot_preview(self, item: ReportItem):
        """Show plot preview."""
        try:
            from PIL import Image, ImageTk
            
            # Create image from PNG bytes
            img_data = io.BytesIO(item.data)
            img = Image.open(img_data)
            
            # Scale to fit preview area
            max_width = 350
            max_height = 200
            
            ratio = min(max_width / img.width, max_height / img.height)
            new_size = (int(img.width * ratio), int(img.height * ratio))
            img = img.resize(new_size, Image.Resampling.LANCZOS)
            
            photo = ImageTk.PhotoImage(img)
            
            img_label = tk.Label(
                self.preview_container,
                image=photo,
                bg=self.theme["bg_secondary"],
            )
            img_label.image = photo  # Keep reference
            img_label.pack(padx=8, pady=8)
            
        except Exception as e:
            tk.Label(
                self.preview_container,
                text=f"Could not display plot preview: {e}",
                font=FONTS.get_font("small"),
                bg=self.theme["bg_secondary"], fg=self.theme["danger"],
            ).pack(pady=20)
    
    def _show_table_preview(self, item: ReportItem):
        """Show table preview."""
        try:
            import pandas as pd
            
            df = item.data
            if not isinstance(df, pd.DataFrame):
                raise ValueError("Data is not a DataFrame")
            
            # Show shape
            tk.Label(
                self.preview_container,
                text=f"Shape: {df.shape[0]} rows × {df.shape[1]} columns",
                font=FONTS.get_font("small"),
                bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
            ).pack(anchor=tk.W, padx=8, pady=4)
            
            # Show first few rows
            preview_text = df.head(5).to_string()
            
            text_widget = tk.Text(
                self.preview_container,
                font=("Courier", 9),
                bg=self.theme["bg_card"],
                fg=self.theme["text_primary"],
                wrap=tk.NONE,
                height=8,
                padx=8, pady=8,
            )
            text_widget.pack(fill=tk.BOTH, expand=True, padx=8, pady=8)
            text_widget.insert("1.0", preview_text)
            text_widget.config(state=tk.DISABLED)
            
            # Add horizontal scrollbar
            h_scroll = ttk.Scrollbar(self.preview_container, orient=tk.HORIZONTAL, command=text_widget.xview)
            h_scroll.pack(fill=tk.X, padx=8)
            text_widget.config(xscrollcommand=h_scroll.set)
            
        except Exception as e:
            tk.Label(
                self.preview_container,
                text=f"Could not display table preview: {e}",
                font=FONTS.get_font("small"),
                bg=self.theme["bg_secondary"], fg=self.theme["danger"],
            ).pack(pady=20)
    
    def _generate_report(self):
        """Generate custom report from queue items."""
        items = report_queue.get_items()
        
        if not items:
            messagebox.showwarning("Empty Queue", "No items in the report queue.\n\nAdd plots and tables using the right-click menu.")
            return
        
        # Get selected formats
        selected_formats = [fmt for fmt, var in self.format_vars.items() if var.get()]
        
        if not selected_formats:
            messagebox.showwarning("No Format", "Please select at least one output format.")
            return
        
        report_title = self.report_title_var.get() or "Custom Report"
        
        # Get output folder and base filename
        folder = filedialog.askdirectory(
            title="Select Output Folder for Custom Report",
        )
        
        if not folder:
            return
        
        # Ask for base filename
        base_name = simpledialog.askstring(
            "Report Filename",
            "Enter base filename (without extension):",
            initialvalue="custom_report",
            parent=self,
        )
        
        if not base_name:
            return
        
        base_name = base_name.strip()
        
        self._show_generating_status()
        
        def do_generate():
            results = {}
            try:
                for fmt in selected_formats:
                    ext = {"docx": ".docx", "pptx": ".pptx", "pdf": ".pdf"}[fmt]
                    filename = os.path.join(folder, f"{base_name}{ext}")
                    
                    try:
                        if fmt == "docx":
                            self._generate_docx(filename, items, report_title)
                        elif fmt == "pptx":
                            self._generate_pptx(filename, items, report_title)
                        elif fmt == "pdf":
                            self._generate_pdf(filename, items, report_title)
                        
                        results[fmt] = filename
                    except Exception as e:
                        results[fmt] = f"Error: {e}"
                
                self.after(0, lambda: self._show_success_multiple(results, folder))
                
            except Exception as e:
                import traceback
                error_msg = f"{e}\n\n{traceback.format_exc()}"
                self.after(0, lambda msg=error_msg: self._show_error(msg))
        
        threading.Thread(target=do_generate, daemon=True).start()
    
    def _show_generating_status(self):
        """Show generating status."""
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        tk.Label(
            self.preview_container,
            text="⏳ Generating report...",
            font=FONTS.get_font("body"),
            bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
        ).pack(pady=50)
    
    def _show_success_multiple(self, results: Dict[str, str], folder: str):
        """Show success message for multiple generated files."""
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        tk.Label(
            self.preview_container,
            text="✓ Reports Generated Successfully",
            font=FONTS.get_font("body", bold=True),
            bg=self.theme["bg_secondary"], fg=self.theme["success"],
        ).pack(pady=(20, 10))
        
        # Show each generated file
        files_generated = []
        for fmt, path in results.items():
            if not str(path).startswith("Error"):
                files_generated.append((fmt, path))
                tk.Label(
                    self.preview_container,
                    text=f"  ✓ {fmt.upper()}: {Path(path).name}",
                    font=FONTS.get_font("small"),
                    bg=self.theme["bg_secondary"], fg=self.theme["text_secondary"],
                    anchor=tk.W,
                ).pack(fill=tk.X, padx=20)
            else:
                tk.Label(
                    self.preview_container,
                    text=f"  ✗ {fmt.upper()}: {path}",
                    font=FONTS.get_font("small"),
                    bg=self.theme["bg_secondary"], fg=self.theme["danger"],
                    anchor=tk.W,
                ).pack(fill=tk.X, padx=20)
        
        # Buttons
        btn_frame = tk.Frame(self.preview_container, bg=self.theme["bg_secondary"])
        btn_frame.pack(pady=15)
        
        ThemedButton(
            btn_frame, text="Open Folder", style="primary",
            command=lambda: self._open_folder(folder),
            theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=4)
        
        # Add individual file open buttons
        for fmt, path in files_generated:
            ThemedButton(
                btn_frame, text=f"Open {fmt.upper()}", style="secondary",
                command=lambda p=path: self._open_file(p),
                theme=self.theme_name,
            ).pack(side=tk.LEFT, padx=4)
    
    def _show_success(self, filename: str):
        """Show success message."""
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        tk.Label(
            self.preview_container,
            text="✓ Report Generated Successfully",
            font=FONTS.get_font("body", bold=True),
            bg=self.theme["bg_secondary"], fg=self.theme["success"],
        ).pack(pady=(30, 10))
        
        tk.Label(
            self.preview_container,
            text=filename,
            font=FONTS.get_font("small"),
            bg=self.theme["bg_secondary"], fg=self.theme["text_muted"],
            wraplength=350,
        ).pack(pady=5)
        
        btn_frame = tk.Frame(self.preview_container, bg=self.theme["bg_secondary"])
        btn_frame.pack(pady=15)
        
        ThemedButton(
            btn_frame, text="Open File", style="primary",
            command=lambda: self._open_file(filename),
            theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=4)
        
        ThemedButton(
            btn_frame, text="Open Folder", style="secondary",
            command=lambda: self._open_folder(str(Path(filename).parent)),
            theme=self.theme_name,
        ).pack(side=tk.LEFT, padx=4)
    
    def _show_error(self, message: str):
        """Show error message."""
        for widget in self.preview_container.winfo_children():
            widget.destroy()
        
        tk.Label(
            self.preview_container,
            text=f"❌ Error\n\n{message[:500]}",
            font=FONTS.get_font("small"),
            bg=self.theme["bg_secondary"], fg=self.theme["danger"],
            wraplength=350, justify=tk.LEFT,
        ).pack(pady=20, padx=10)
    
    def _format_value(self, val, max_length: int = 100) -> str:
        """Format a value for display, rounding floats to 3 decimal places."""
        if val is None:
            return ""
        if isinstance(val, float):
            # Round to 3 decimal places
            if val == int(val):
                return str(int(val))
            return f"{val:.3f}"
        # Truncate strings if needed
        result = str(val)
        if len(result) > max_length:
            return result[:max_length-3] + "..."
        return result
    
    def _generate_docx(self, filename: str, items: List[ReportItem], title: str):
        """Generate Word document."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            raise ImportError("python-docx is required for DOCX export. Install with: pip install python-docx")
        
        doc = Document()
        
        # Title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Subtitle with date
        subtitle = doc.add_paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()  # Space
        
        for i, item in enumerate(items):
            # Section heading
            doc.add_heading(f"{i+1}. {item.title}", level=1)
            
            # Source info
            source_para = doc.add_paragraph()
            source_run = source_para.add_run(f"Source: {item.source_panel}")
            source_run.font.size = Pt(9)
            source_run.font.italic = True
            
            # Description if present
            if item.description:
                desc_para = doc.add_paragraph(item.description)
                desc_para.style = 'Quote'
            
            if item.item_type == "plot":
                # Add image
                img_stream = io.BytesIO(item.data)
                doc.add_picture(img_stream, width=Inches(6))
            else:
                # Add table
                import pandas as pd
                df = item.data
                if isinstance(df, pd.DataFrame):
                    # Limit rows for Word
                    df_display = df.head(50)
                    
                    table = doc.add_table(rows=1, cols=len(df_display.columns))
                    table.style = 'Table Grid'
                    
                    # Header row
                    header_cells = table.rows[0].cells
                    for j, col in enumerate(df_display.columns):
                        header_cells[j].text = str(col)
                    
                    # Data rows
                    for idx, row in df_display.iterrows():
                        row_cells = table.add_row().cells
                        for j, val in enumerate(row):
                            row_cells[j].text = self._format_value(val, 100)
                    
                    if len(df) > 50:
                        doc.add_paragraph(f"(Showing first 50 of {len(df)} rows)")
            
            doc.add_paragraph()  # Space between items
        
        doc.save(filename)
    
    def _generate_pptx(self, filename: str, items: List[ReportItem], title: str):
        """Generate PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx is required for PPTX export. Install with: pip install python-pptx")
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        for item in items:
            if item.item_type == "plot":
                # Content slide with image
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tf = title_box.text_frame
                tf.text = item.title
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.bold = True
                
                # Add image
                img_stream = io.BytesIO(item.data)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))
                
                # Add caption/description if present
                if item.description:
                    caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
                    tf = caption_box.text_frame
                    tf.text = item.description
                    tf.paragraphs[0].font.size = Pt(12)
                    tf.paragraphs[0].font.italic = True
            else:
                # Table slide
                import pandas as pd
                df = item.data
                if isinstance(df, pd.DataFrame):
                    slide_layout = prs.slide_layouts[5]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                    tf = title_box.text_frame
                    tf.text = item.title
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.bold = True
                    
                    # Limit size for slide
                    df_display = df.head(10).iloc[:, :8]
                    
                    # Add table
                    rows, cols = df_display.shape
                    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * (rows + 1)))
                    
                    # Header
                    for j, col in enumerate(df_display.columns):
                        cell = table.table.cell(0, j)
                        cell.text = str(col)[:20]
                    
                    # Data
                    for i, (idx, row) in enumerate(df_display.iterrows()):
                        for j, val in enumerate(row):
                            cell = table.table.cell(i + 1, j)
                            cell.text = self._format_value(val, 20)
        
        prs.save(filename)
    
    def _generate_pdf(self, filename: str, items: List[ReportItem], title: str):
        """Generate PDF document."""
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
            from reportlab.lib.enums import TA_CENTER
        except ImportError:
            raise ImportError("reportlab is required for PDF export. Install with: pip install reportlab")
        
        doc = SimpleDocTemplate(filename, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=12,
        )
        
        # Title
        story.append(Paragraph(title, title_style))
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        for i, item in enumerate(items):
            # Section heading
            story.append(Paragraph(f"{i+1}. {item.title}", styles['Heading2']))
            story.append(Paragraph(f"<i>Source: {item.source_panel}</i>", styles['Normal']))
            
            if item.description:
                story.append(Paragraph(item.description, styles['Italic']))
            
            story.append(Spacer(1, 0.2*inch))
            
            if item.item_type == "plot":
                # Add image
                from PIL import Image as PILImage
                img_stream = io.BytesIO(item.data)
                pil_img = PILImage.open(img_stream)
                
                # Scale to fit page
                max_width = 6.5 * inch
                max_height = 5 * inch
                
                ratio = min(max_width / pil_img.width, max_height / pil_img.height)
                new_width = pil_img.width * ratio
                new_height = pil_img.height * ratio
                
                img_stream.seek(0)
                img = Image(img_stream, width=new_width, height=new_height)
                story.append(img)
            else:
                # Add table
                import pandas as pd
                df = item.data
                if isinstance(df, pd.DataFrame):
                    df_display = df.head(20).iloc[:, :6]  # Limit for PDF
                    
                    # Convert to table data with formatted values
                    data = [list(df_display.columns)]
                    for idx, row in df_display.iterrows():
                        data.append([self._format_value(v, 30) for v in row])
                    
                    table = Table(data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ]))
                    story.append(table)
                    
                    if len(df) > 20:
                        story.append(Paragraph(f"(Showing first 20 of {len(df)} rows)", styles['Italic']))
            
            story.append(Spacer(1, 0.5*inch))
        
        doc.build(story)
    
    def _open_file(self, path: str):
        """Open file with default application."""
        import platform
        import subprocess
        
        try:
            if platform.system() == "Windows":
                os.startfile(path)
            elif platform.system() == "Darwin":
                subprocess.run(["open", path])
            else:
                subprocess.run(["xdg-open", path])
        except Exception as e:
            messagebox.showerror("Error", f"Could not open file: {e}")
    
    def _open_folder(self, path: str):
        """Open folder in file explorer."""
        import platform
        import subprocess
        
        try:
            if platform.system() == "Windows":
                os.startfile(path)
            elif platform.system() == "Darwin":
                subprocess.run(["open", path])
            else:
                subprocess.run(["xdg-open", path])
        except Exception as e:
            messagebox.showerror("Error", f"Could not open folder: {e}")
