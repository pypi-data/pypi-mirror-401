# -*- coding: utf-8 -*-
"""
Report Generator Module.

Main orchestrator for generating reports in multiple formats.
"""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Union

import numpy as np
import pandas as pd

from biblium.reports.config import ReportConfig, ReportFormat, ReportLevel
from biblium.reports.template import ReportTemplate, ReportItem

if TYPE_CHECKING:
    from biblium.bibstats import BiblioStats


class ReportGenerator:
    """
    Generate bibliometric analysis reports in multiple formats.
    
    Supports: XLSX, DOCX, PPTX, LaTeX, HTML
    
    Parameters
    ----------
    biblio : BiblioStats
        The BiblioStats/BiblioAnalysis instance with analysis results.
    config : ReportConfig, optional
        Report configuration. Uses default if not provided.
    template : ReportTemplate, optional
        Custom template. Uses default if not provided.
    
    Examples
    --------
    >>> ba = BiblioAnalysis("data.csv", db="scopus")
    >>> report = ReportGenerator(ba)
    >>> 
    >>> # Generate single format
    >>> report.generate("xlsx", level="standard")
    >>> 
    >>> # Generate with custom config
    >>> config = ReportConfig(title="My Analysis", level=ReportLevel.FULL)
    >>> report.generate("docx", config=config)
    >>> 
    >>> # Generate all formats
    >>> report.generate_all(level="standard")
    """
    
    def __init__(
        self,
        biblio: "BiblioStats",
        config: Optional[ReportConfig] = None,
        template: Optional[ReportTemplate] = None,
    ):
        self.biblio = biblio
        self.config = config or ReportConfig()
        self.template = template or ReportTemplate()
        
        # Set output directory
        if self.config.output_dir:
            self.output_dir = Path(self.config.output_dir)
        elif hasattr(biblio, "res_folder") and biblio.res_folder:
            self.output_dir = Path(biblio.res_folder) / "reports"
        else:
            self.output_dir = Path.cwd() / "reports"
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate(
        self,
        format: Union[str, ReportFormat],
        level: Optional[Union[str, ReportLevel]] = None,
        config: Optional[ReportConfig] = None,
        filename: Optional[str] = None,
    ) -> Path:
        """
        Generate report in specified format.
        
        Parameters
        ----------
        format : str or ReportFormat
            Output format: "xlsx", "docx", "pptx", "tex", "html"
        level : str or ReportLevel, optional
            Report level. Overrides config if provided.
        config : ReportConfig, optional
            Custom config for this generation.
        filename : str, optional
            Custom output filename (without extension).
        
        Returns
        -------
        Path
            Path to generated report file.
        """
        if isinstance(format, str):
            format = ReportFormat(format.lower())
        
        cfg = config or self.config
        if level:
            if isinstance(level, str):
                level = ReportLevel(level)
            cfg = ReportConfig(**{**cfg.__dict__, "level": level})
        
        # Get items for this level
        items = self.template.filter_available(self.biblio, cfg.level.value)
        
        if not items:
            raise ValueError(f"No data available for report level '{cfg.level.value}'")
        
        # Generate filename
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M")
            filename = f"{cfg.filename_prefix}_{cfg.level.value}_{timestamp}"
        
        # Dispatch to format-specific generator
        generators = {
            ReportFormat.XLSX: self._generate_xlsx,
            ReportFormat.DOCX: self._generate_docx,
            ReportFormat.PPTX: self._generate_pptx,
            ReportFormat.LATEX: self._generate_latex,
            ReportFormat.HTML: self._generate_html,
        }
        
        if format not in generators:
            raise ValueError(f"Unsupported format: {format}")
        
        output_path = generators[format](items, cfg, filename)
        print(f"✓ Generated: {output_path}")
        return output_path
    
    def generate_all(
        self,
        level: Optional[Union[str, ReportLevel]] = None,
        formats: Optional[List[str]] = None,
    ) -> Dict[str, Path]:
        """
        Generate reports in multiple formats.
        
        Parameters
        ----------
        level : str or ReportLevel, optional
            Report level for all formats.
        formats : list of str, optional
            Formats to generate. Default: all supported formats.
        
        Returns
        -------
        dict
            Mapping of format name to output path.
        """
        if formats is None:
            formats = ["xlsx", "docx", "pptx", "html", "tex"]
        
        results = {}
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        level_str = level if isinstance(level, str) else (level.value if level else self.config.level.value)
        base_filename = f"{self.config.filename_prefix}_{level_str}_{timestamp}"
        
        for fmt in formats:
            try:
                path = self.generate(fmt, level=level, filename=base_filename)
                results[fmt] = path
            except Exception as e:
                print(f"✗ Failed to generate {fmt}: {e}")
        
        return results
    
    # =========================================================================
    # XLSX GENERATION
    # =========================================================================
    
    def _generate_xlsx(
        self,
        items: List[ReportItem],
        config: ReportConfig,
        filename: str,
    ) -> Path:
        """Generate Excel report."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils.dataframe import dataframe_to_rows
        from openpyxl.drawing.image import Image as XLImage
        
        wb = Workbook()
        del wb["Sheet"]  # Remove default sheet
        
        # Style definitions
        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill("solid", fgColor=config.style.primary_color)
        thin_border = Border(
            left=Side(style='thin', color='CCCCCC'),
            right=Side(style='thin', color='CCCCCC'),
            top=Side(style='thin', color='CCCCCC'),
            bottom=Side(style='thin', color='CCCCCC')
        )
        
        def format_value(val, dtype=None):
            """Format value based on column dtype."""
            if pd.isna(val):
                return val
            # Check column dtype first
            if dtype is not None:
                if pd.api.types.is_integer_dtype(dtype):
                    return int(val) if not pd.isna(val) else val
                elif pd.api.types.is_float_dtype(dtype):
                    return round(float(val), 2)
            # Fallback to value type
            if isinstance(val, (int, np.integer)):
                return int(val)
            if isinstance(val, (float, np.floating)):
                return round(val, 2)
            return val
        
        # Create TOC sheet
        if config.include_toc:
            ws_toc = wb.create_sheet("Contents")
            ws_toc["A1"] = config.title
            ws_toc["A1"].font = Font(bold=True, size=16)
            ws_toc["A2"] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            ws_toc["A3"] = f"Documents: {self.biblio.n}"
            
            row = 5
            for item in items:
                ws_toc.cell(row=row, column=1, value=item.level1)
                ws_toc.cell(row=row, column=2, value=item.level2)
                ws_toc.cell(row=row, column=3, value=item.item_type)
                row += 1
        
        # Group items by section
        sections = self.template.get_items_by_section(config.level.value)
        
        for section_name, section_items in sections.items():
            # Filter to available items
            available = [i for i in section_items if i in items]
            if not available:
                continue
            
            # Create sheet for section
            sheet_name = section_name[:31]  # Excel limit
            ws = wb.create_sheet(sheet_name)
            
            current_row = 1
            
            for item in available:
                if item.is_table and item.data_attr:
                    data = item.get_data(self.biblio)
                    if data is not None and len(data) > 0:
                        # Limit rows
                        data = data.head(config.top_n_tables)
                        
                        # Add caption
                        ws.cell(row=current_row, column=1, value=item.caption)
                        ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                        current_row += 1
                        
                        # Add headers
                        for col_idx, col_name in enumerate(data.columns, 1):
                            cell = ws.cell(row=current_row, column=col_idx, value=col_name)
                            cell.font = header_font
                            cell.fill = header_fill
                            cell.border = thin_border
                        current_row += 1
                        
                        # Add data with proper formatting based on column dtype
                        col_dtypes = data.dtypes
                        for _, row_data in data.iterrows():
                            for col_idx, (col_name, value) in enumerate(row_data.items(), 1):
                                formatted = format_value(value, col_dtypes[col_name])
                                cell = ws.cell(row=current_row, column=col_idx, value=formatted)
                                cell.border = thin_border
                            current_row += 1
                        
                        current_row += 2  # Space between tables
                
                elif item.is_plot and item.plot_filename:
                    # Check if plot file exists
                    plot_path = None
                    if hasattr(self.biblio, "res_folder") and self.biblio.res_folder:
                        potential_path = Path(self.biblio.res_folder) / "figures" / item.plot_filename
                        if potential_path.exists():
                            plot_path = potential_path
                    
                    if plot_path:
                        try:
                            img = XLImage(str(plot_path))
                            img.width = 600
                            img.height = 400
                            ws.add_image(img, f"A{current_row}")
                            current_row += 25  # Space for image
                        except Exception:
                            pass
            
            # Auto-fit columns (approximation)
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except:
                        pass
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
        
        output_path = self.output_dir / f"{filename}.xlsx"
        wb.save(output_path)
        return output_path
    
    # =========================================================================
    # DOCX GENERATION
    # =========================================================================
    
    def _generate_docx(
        self,
        items: List[ReportItem],
        config: ReportConfig,
        filename: str,
    ) -> Path:
        """Generate Word document report."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
        except ImportError:
            raise ImportError("python-docx required. Install with: pip install python-docx")
        
        doc = Document()
        
        # Set margins
        for section in doc.sections:
            section.top_margin = Inches(config.style.margin_top)
            section.bottom_margin = Inches(config.style.margin_bottom)
            section.left_margin = Inches(config.style.margin_left)
            section.right_margin = Inches(config.style.margin_right)
        
        # Title
        title = doc.add_heading(config.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        if config.subtitle:
            subtitle = doc.add_paragraph(config.subtitle)
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata
        meta = doc.add_paragraph()
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if config.author:
            meta.add_run(f"Author: {config.author}\n")
        meta.add_run(f"Generated: {datetime.now().strftime('%Y-%m-%d')}\n")
        meta.add_run(f"Documents analyzed: {self.biblio.n}")
        
        doc.add_page_break()
        
        # Table of Contents placeholder
        if config.include_toc:
            doc.add_heading("Table of Contents", level=1)
            doc.add_paragraph("(Update field to generate TOC)")
            doc.add_page_break()
        
        # Content by section
        current_section = None
        
        for item in items:
            # Section heading
            if item.level1 != current_section:
                current_section = item.level1
                doc.add_heading(item.level1, level=1)
            
            # Item heading
            doc.add_heading(item.caption, level=2)
            
            # Description
            if item.description:
                doc.add_paragraph(item.description)
            
            if item.is_table and item.data_attr:
                data = item.get_data(self.biblio)
                if data is not None and len(data) > 0:
                    data = data.head(config.top_n_tables)
                    
                    # Create table
                    table = doc.add_table(rows=1, cols=len(data.columns))
                    table.style = 'Table Grid'
                    
                    # Headers
                    header_cells = table.rows[0].cells
                    for i, col in enumerate(data.columns):
                        header_cells[i].text = str(col)
                        header_cells[i].paragraphs[0].runs[0].bold = True
                    
                    # Data rows with dtype-aware formatting
                    col_dtypes = data.dtypes
                    for _, row in data.iterrows():
                        row_cells = table.add_row().cells
                        for i, (col_name, value) in enumerate(row.items()):
                            if pd.notna(value):
                                dtype = col_dtypes[col_name]
                                if pd.api.types.is_integer_dtype(dtype):
                                    row_cells[i].text = str(int(value))
                                elif pd.api.types.is_float_dtype(dtype):
                                    row_cells[i].text = f"{value:.2f}"
                                else:
                                    row_cells[i].text = str(value)
                            else:
                                row_cells[i].text = ""
                    
                    doc.add_paragraph()  # Space
            
            elif item.is_plot and item.plot_filename:
                plot_path = None
                if hasattr(self.biblio, "res_folder") and self.biblio.res_folder:
                    potential_path = Path(self.biblio.res_folder) / "figures" / item.plot_filename
                    if potential_path.exists():
                        plot_path = potential_path
                
                if plot_path:
                    try:
                        doc.add_picture(str(plot_path), width=Inches(6))
                        last_para = doc.paragraphs[-1]
                        last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    except Exception:
                        doc.add_paragraph(f"[Image: {item.plot_filename}]")
                else:
                    doc.add_paragraph(f"[Image not found: {item.plot_filename}]")
            
            if item.page_break:
                doc.add_page_break()
        
        output_path = self.output_dir / f"{filename}.docx"
        doc.save(output_path)
        return output_path
    
    # =========================================================================
    # PPTX GENERATION
    # =========================================================================
    
    def _generate_pptx(
        self,
        items: List[ReportItem],
        config: ReportConfig,
        filename: str,
    ) -> Path:
        """Generate PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError("python-pptx required. Install with: pip install python-pptx")
        
        prs = Presentation()
        
        # Set slide size (16:9)
        if config.pptx_settings.get("aspect_ratio") == "16:9":
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        
        # Title slide
        title_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(title_layout)
        
        # Add title text box
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(12.33)
        height = Inches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = config.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if config.subtitle:
            top = Inches(4)
            txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = config.subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
        
        # Metadata
        top = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.biblio.n} documents | Generated {datetime.now().strftime('%Y-%m-%d')}"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        
        # Content slides
        slide_count = 1
        max_slides = config.pptx_settings.get("max_slides", 50)
        
        current_section = None
        
        for item in items:
            if slide_count >= max_slides:
                break
            
            # Section divider
            if item.level1 != current_section:
                current_section = item.level1
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = item.level1
                p.font.size = Pt(40)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                slide_count += 1
            
            if slide_count >= max_slides:
                break
            
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.caption
            p.font.size = Pt(28)
            p.font.bold = True
            
            if item.is_table and item.data_attr:
                data = item.get_data(self.biblio)
                if data is not None and len(data) > 0:
                    # Limit for slides
                    data = data.head(min(10, config.top_n_tables))
                    
                    # Simple table representation
                    rows = min(len(data) + 1, 11)
                    cols = min(len(data.columns), 5)
                    
                    left = Inches(0.5)
                    top = Inches(1.3)
                    width = Inches(12.33)
                    height = Inches(5.5)
                    
                    # Get column dtypes for proper formatting
                    col_dtypes = data.dtypes
                    col_list = list(data.columns[:cols])
                    
                    def format_cell_value(val, col_name):
                        """Format cell value based on column dtype."""
                        if pd.isna(val):
                            return ""
                        dtype = col_dtypes.get(col_name)
                        if dtype is not None:
                            if pd.api.types.is_integer_dtype(dtype):
                                return str(int(val))[:30]
                            elif pd.api.types.is_float_dtype(dtype):
                                return f"{val:.2f}"[:30]
                        return str(val)[:30]
                    
                    try:
                        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                        table = shape.table
                        
                        # Headers
                        for i, col in enumerate(col_list):
                            cell = table.cell(0, i)
                            cell.text = str(col)[:30]
                        
                        # Data
                        for row_idx, (_, row) in enumerate(data.head(rows-1).iterrows()):
                            for col_idx, col_name in enumerate(col_list):
                                cell = table.cell(row_idx + 1, col_idx)
                                cell.text = format_cell_value(row[col_name], col_name)
                    except Exception as e:
                        # Fallback to text if table fails
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = f"[Table: {item.caption}]"
            
            elif item.is_plot and item.plot_filename:
                plot_path = None
                if hasattr(self.biblio, "res_folder") and self.biblio.res_folder:
                    potential_path = Path(self.biblio.res_folder) / "figures" / item.plot_filename
                    if potential_path.exists():
                        plot_path = potential_path
                
                if plot_path:
                    try:
                        slide.shapes.add_picture(
                            str(plot_path),
                            Inches(0.5), Inches(1.3),
                            width=Inches(12.33), height=Inches(5.8)
                        )
                    except Exception:
                        pass
            
            slide_count += 1
        
        # Add slide numbers if configured
        if config.pptx_settings.get("include_slide_numbers"):
            for i, slide in enumerate(list(prs.slides)[1:], 2):  # Skip title slide
                try:
                    txBox = slide.shapes.add_textbox(
                        Inches(12.5), Inches(7), Inches(0.8), Inches(0.3)
                    )
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = str(i)
                    p.font.size = Pt(10)
                    p.alignment = PP_ALIGN.RIGHT
                except Exception:
                    pass  # Skip if can't add slide number
        
        output_path = self.output_dir / f"{filename}.pptx"
        prs.save(output_path)
        return output_path
    
    # =========================================================================
    # LATEX GENERATION
    # =========================================================================
    
    def _generate_latex(
        self,
        items: List[ReportItem],
        config: ReportConfig,
        filename: str,
    ) -> Path:
        """Generate LaTeX document."""
        
        settings = config.latex_settings
        
        lines = [
            f"\\documentclass[{settings.get('font_size', '11pt')},{settings.get('paper_size', 'a4paper')}]{{{settings.get('document_class', 'article')}}}",
            "",
            "% Packages",
            "\\usepackage[utf8]{inputenc}",
            "\\usepackage[T1]{fontenc}",
            "\\usepackage{graphicx}",
            "\\usepackage{booktabs}",
            "\\usepackage{longtable}",
            "\\usepackage{hyperref}",
            "\\usepackage{geometry}",
            "\\usepackage{float}",
            "\\usepackage{caption}",
            "",
            f"\\geometry{{margin={config.style.margin_left}in}}",
            "",
            f"\\title{{{self._escape_latex(config.title)}}}",
        ]
        
        if config.author:
            lines.append(f"\\author{{{self._escape_latex(config.author)}}}")
        
        lines.extend([
            f"\\date{{{datetime.now().strftime('%B %d, %Y')}}}",
            "",
            "\\begin{document}",
            "",
            "\\maketitle",
            "",
        ])
        
        # Abstract/Summary
        if config.include_summary:
            lines.extend([
                "\\begin{abstract}",
                f"This report presents a bibliometric analysis of {self.biblio.n} documents. ",
                "The analysis includes publication trends, author productivity, keyword analysis, ",
                "and various bibliometric indicators.",
                "\\end{abstract}",
                "",
            ])
        
        # TOC
        if config.include_toc:
            lines.extend([
                "\\tableofcontents",
                "\\newpage",
                "",
            ])
        
        # Content
        current_section = None
        
        for item in items:
            if item.level1 != current_section:
                current_section = item.level1
                lines.append(f"\\section{{{self._escape_latex(item.level1)}}}")
                lines.append("")
            
            lines.append(f"\\subsection{{{self._escape_latex(item.caption)}}}")
            
            if item.description:
                lines.append(self._escape_latex(item.description))
                lines.append("")
            
            if item.is_table and item.data_attr:
                data = item.get_data(self.biblio)
                if data is not None and len(data) > 0:
                    data = data.head(config.top_n_tables)
                    lines.append(self._df_to_latex(data, item.caption))
                    lines.append("")
            
            elif item.is_plot and item.plot_filename:
                plot_path = None
                if hasattr(self.biblio, "res_folder") and self.biblio.res_folder:
                    potential_path = Path(self.biblio.res_folder) / "figures" / item.plot_filename
                    if potential_path.exists():
                        plot_path = potential_path
                
                if plot_path:
                    lines.extend([
                        "\\begin{figure}[H]",
                        "\\centering",
                        f"\\includegraphics[width=0.9\\textwidth]{{{plot_path}}}",
                        f"\\caption{{{self._escape_latex(item.caption)}}}",
                        "\\end{figure}",
                        "",
                    ])
            
            if item.page_break:
                lines.append("\\newpage")
                lines.append("")
        
        lines.extend([
            "",
            "\\end{document}",
        ])
        
        output_path = self.output_dir / f"{filename}.tex"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        
        return output_path
    
    def _escape_latex(self, text: str) -> str:
        """Escape special LaTeX characters."""
        if not text:
            return ""
        chars = {
            '&': r'\&',
            '%': r'\%',
            '$': r'\$',
            '#': r'\#',
            '_': r'\_',
            '{': r'\{',
            '}': r'\}',
            '~': r'\textasciitilde{}',
            '^': r'\^{}',
        }
        for char, replacement in chars.items():
            text = text.replace(char, replacement)
        return text
    
    def _df_to_latex(self, df: pd.DataFrame, caption: str) -> str:
        """Convert DataFrame to LaTeX table."""
        # Limit columns for readability
        df = df.iloc[:, :6] if len(df.columns) > 6 else df
        col_dtypes = df.dtypes
        
        def format_value(val, col_name):
            """Format value based on column dtype."""
            if pd.isna(val):
                return ""
            dtype = col_dtypes.get(col_name)
            if dtype is not None:
                if pd.api.types.is_integer_dtype(dtype):
                    return str(int(val))
                elif pd.api.types.is_float_dtype(dtype):
                    return f"{val:.2f}"
            return str(val)[:30]
        
        lines = [
            "\\begin{table}[H]",
            "\\centering",
            f"\\caption{{{self._escape_latex(caption)}}}",
            "\\begin{tabular}{" + "l" * len(df.columns) + "}",
            "\\toprule",
        ]
        
        # Header
        headers = " & ".join(self._escape_latex(str(c)) for c in df.columns)
        lines.append(f"{headers} \\\\")
        lines.append("\\midrule")
        
        # Data with dtype-aware formatting
        for _, row in df.iterrows():
            row_str = " & ".join(self._escape_latex(format_value(v, col)) for col, v in row.items())
            lines.append(f"{row_str} \\\\")
        
        lines.extend([
            "\\bottomrule",
            "\\end{tabular}",
            "\\end{table}",
        ])
        
        return "\n".join(lines)
    
    # =========================================================================
    # HTML GENERATION
    # =========================================================================
    
    def _generate_html(
        self,
        items: List[ReportItem],
        config: ReportConfig,
        filename: str,
    ) -> Path:
        """Generate HTML report."""
        
        settings = config.html_settings
        is_dark = settings.get("theme") == "dark"
        
        # CSS
        css = self._get_html_css(config, is_dark)
        
        # Build HTML
        html_parts = [
            "<!DOCTYPE html>",
            "<html lang='en'>",
            "<head>",
            "<meta charset='UTF-8'>",
            "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
            f"<title>{config.title}</title>",
            "<style>",
            css,
            "</style>",
            "</head>",
            "<body>",
            "<div class='container'>",
        ]
        
        # Header
        html_parts.extend([
            "<header>",
            f"<h1>{config.title}</h1>",
        ])
        if config.subtitle:
            html_parts.append(f"<p class='subtitle'>{config.subtitle}</p>")
        html_parts.extend([
            f"<p class='meta'>{self.biblio.n} documents | Generated {datetime.now().strftime('%Y-%m-%d')}</p>",
            "</header>",
        ])
        
        # Navigation/TOC
        if config.include_toc:
            html_parts.append("<nav class='toc'>")
            html_parts.append("<h2>Contents</h2>")
            html_parts.append("<ul>")
            current_section = None
            for item in items:
                if item.level1 != current_section:
                    current_section = item.level1
                    section_id = current_section.lower().replace(" ", "-")
                    html_parts.append(f"<li><a href='#{section_id}'>{current_section}</a></li>")
            html_parts.append("</ul>")
            html_parts.append("</nav>")
        
        # Content
        html_parts.append("<main>")
        current_section = None
        
        for item in items:
            if item.level1 != current_section:
                if current_section is not None:
                    html_parts.append("</section>")
                current_section = item.level1
                section_id = current_section.lower().replace(" ", "-")
                html_parts.append(f"<section id='{section_id}'>")
                html_parts.append(f"<h2>{current_section}</h2>")
            
            html_parts.append("<article>")
            html_parts.append(f"<h3>{item.caption}</h3>")
            
            if item.description:
                html_parts.append(f"<p class='description'>{item.description}</p>")
            
            if item.is_table and item.data_attr:
                data = item.get_data(self.biblio)
                if data is not None and len(data) > 0:
                    data = data.head(config.top_n_tables)
                    html_parts.append(self._df_to_html(data))
            
            elif item.is_plot and item.plot_filename:
                plot_path = None
                if hasattr(self.biblio, "res_folder") and self.biblio.res_folder:
                    potential_path = Path(self.biblio.res_folder) / "figures" / item.plot_filename
                    if potential_path.exists():
                        plot_path = potential_path
                
                if plot_path and settings.get("single_file"):
                    # Embed image as base64
                    import base64
                    with open(plot_path, "rb") as f:
                        img_data = base64.b64encode(f.read()).decode()
                    ext = plot_path.suffix.lower().replace(".", "")
                    mime = {"png": "image/png", "jpg": "image/jpeg", "svg": "image/svg+xml"}.get(ext, "image/png")
                    html_parts.append(f"<img src='data:{mime};base64,{img_data}' alt='{item.caption}' class='plot'>")
                elif plot_path:
                    html_parts.append(f"<img src='{plot_path}' alt='{item.caption}' class='plot'>")
            
            html_parts.append("</article>")
        
        if current_section is not None:
            html_parts.append("</section>")
        
        html_parts.extend([
            "</main>",
            "<footer>",
            f"<p>Generated by Biblium | {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>",
            "</footer>",
            "</div>",
            "</body>",
            "</html>",
        ])
        
        output_path = self.output_dir / f"{filename}.html"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_parts))
        
        return output_path
    
    def _get_html_css(self, config: ReportConfig, is_dark: bool) -> str:
        """Generate CSS for HTML report."""
        style = config.style
        
        if is_dark:
            bg = "#1a1a2e"
            text = "#eaeaea"
            card_bg = "#16213e"
            border = "#0f3460"
        else:
            bg = "#f5f5f5"
            text = "#333333"
            card_bg = "#ffffff"
            border = "#dddddd"
        
        return f"""
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: {style.body_font}, -apple-system, BlinkMacSystemFont, sans-serif;
            font-size: {style.body_size}pt;
            line-height: 1.6;
            color: {text};
            background: {bg};
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            padding: 2rem;
        }}
        header {{
            text-align: center;
            margin-bottom: 3rem;
            padding: 2rem;
            background: linear-gradient(135deg, #{style.primary_color}, #{style.secondary_color});
            color: white;
            border-radius: 10px;
        }}
        h1 {{ font-size: {style.title_size}pt; margin-bottom: 0.5rem; }}
        h2 {{ font-size: {style.heading1_size}pt; color: #{style.primary_color}; margin: 2rem 0 1rem; }}
        h3 {{ font-size: {style.heading2_size}pt; margin: 1rem 0 0.5rem; }}
        .subtitle {{ font-size: 1.2rem; opacity: 0.9; }}
        .meta {{ font-size: 0.9rem; opacity: 0.8; }}
        nav.toc {{
            background: {card_bg};
            padding: 1.5rem;
            border-radius: 8px;
            margin-bottom: 2rem;
            border: 1px solid {border};
        }}
        nav.toc ul {{ list-style: none; display: flex; flex-wrap: wrap; gap: 1rem; }}
        nav.toc a {{ color: #{style.primary_color}; text-decoration: none; }}
        nav.toc a:hover {{ text-decoration: underline; }}
        section {{
            background: {card_bg};
            padding: 1.5rem;
            border-radius: 8px;
            margin-bottom: 1.5rem;
            border: 1px solid {border};
        }}
        article {{ margin-bottom: 2rem; }}
        .description {{ color: #666; margin-bottom: 1rem; }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 1rem 0;
            font-size: 0.9rem;
        }}
        th, td {{
            padding: 0.75rem;
            text-align: left;
            border-bottom: 1px solid {border};
        }}
        th {{
            background: #{style.primary_color};
            color: white;
            font-weight: 600;
        }}
        tr:hover {{ background: rgba(0,0,0,0.02); }}
        img.plot {{
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            margin: 1rem 0;
        }}
        footer {{
            text-align: center;
            padding: 2rem;
            color: #666;
            font-size: 0.85rem;
        }}
        @media (max-width: 768px) {{
            .container {{ padding: 1rem; }}
            h1 {{ font-size: 1.5rem; }}
        }}
        """
    
    def _df_to_html(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to HTML table."""
        # Limit columns
        df = df.iloc[:, :8] if len(df.columns) > 8 else df
        col_dtypes = df.dtypes
        
        def format_value(val, col_name):
            """Format value based on column dtype."""
            if pd.isna(val):
                return ""
            dtype = col_dtypes.get(col_name)
            if dtype is not None:
                if pd.api.types.is_integer_dtype(dtype):
                    return str(int(val))
                elif pd.api.types.is_float_dtype(dtype):
                    return f"{val:.2f}"
            return str(val)[:50]
        
        html = ["<table>", "<thead>", "<tr>"]
        for col in df.columns:
            html.append(f"<th>{col}</th>")
        html.extend(["</tr>", "</thead>", "<tbody>"])
        
        for _, row in df.iterrows():
            html.append("<tr>")
            for col_name, val in row.items():
                display_val = format_value(val, col_name)
                html.append(f"<td>{display_val}</td>")
            html.append("</tr>")
        
        html.extend(["</tbody>", "</table>"])
        return "\n".join(html)
