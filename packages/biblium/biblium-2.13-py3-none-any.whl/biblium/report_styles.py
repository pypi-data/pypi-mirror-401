"""
Modern Report Styling Module for Biblium

Provides professional, modern styling for Word, PowerPoint, Excel, and LaTeX reports.
Includes color palettes, typography settings, and formatting utilities.
"""

from __future__ import annotations
from typing import Dict, List, Optional, Tuple, Any
from pathlib import Path
from datetime import datetime

# =============================================================================
# COLOR PALETTES
# =============================================================================

class ModernPalette:
    """Modern color palette for reports."""
    
    # Primary colors
    PRIMARY = "#2C3E50"       # Dark blue-gray (headings)
    SECONDARY = "#3498DB"     # Bright blue (accents)
    ACCENT = "#E74C3C"        # Red (highlights)
    SUCCESS = "#27AE60"       # Green (positive values)
    WARNING = "#F39C12"       # Orange (warnings)
    
    # Neutral colors
    DARK = "#2C3E50"
    MEDIUM = "#7F8C8D"
    LIGHT = "#BDC3C7"
    LIGHTER = "#ECF0F1"
    WHITE = "#FFFFFF"
    
    # Table colors
    TABLE_HEADER_BG = "#2C3E50"
    TABLE_HEADER_FG = "#FFFFFF"
    TABLE_ROW_ALT = "#F8F9FA"
    TABLE_BORDER = "#DEE2E6"
    
    # Chart colors (for consistent visualizations)
    CHART_COLORS = [
        "#3498DB",  # Blue
        "#E74C3C",  # Red
        "#2ECC71",  # Green
        "#F39C12",  # Orange
        "#9B59B6",  # Purple
        "#1ABC9C",  # Teal
        "#34495E",  # Dark gray
        "#E67E22",  # Dark orange
    ]
    
    @classmethod
    def hex_to_rgb(cls, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    @classmethod
    def get_gradient(cls, start_hex: str, end_hex: str, steps: int) -> List[str]:
        """Generate gradient between two colors."""
        start = cls.hex_to_rgb(start_hex)
        end = cls.hex_to_rgb(end_hex)
        
        gradient = []
        for i in range(steps):
            ratio = i / (steps - 1) if steps > 1 else 0
            r = int(start[0] + (end[0] - start[0]) * ratio)
            g = int(start[1] + (end[1] - start[1]) * ratio)
            b = int(start[2] + (end[2] - start[2]) * ratio)
            gradient.append(f"#{r:02x}{g:02x}{b:02x}")
        return gradient


# =============================================================================
# TYPOGRAPHY
# =============================================================================

class Typography:
    """Typography settings for modern reports."""
    
    # Font families
    HEADING_FONT = "Calibri Light"
    BODY_FONT = "Calibri"
    MONO_FONT = "Consolas"
    
    # Font sizes (in points)
    TITLE_SIZE = 28
    HEADING1_SIZE = 18
    HEADING2_SIZE = 14
    HEADING3_SIZE = 12
    BODY_SIZE = 11
    CAPTION_SIZE = 9
    SMALL_SIZE = 8
    
    # Line spacing
    LINE_SPACING = 1.15
    PARAGRAPH_SPACING = 6  # points after paragraph


# =============================================================================
# WORD DOCUMENT STYLING
# =============================================================================

def style_word_document(doc, title: str = "Bibliometric Analysis Report"):
    """
    Apply modern styling to a Word document.
    
    Parameters
    ----------
    doc : docx.Document
        The document to style.
    title : str
        Report title for cover page.
    """
    from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    from docx.enum.table import WD_TABLE_ALIGNMENT
    
    # Update default styles
    styles = doc.styles
    
    # Normal style
    normal = styles["Normal"]
    normal.font.name = Typography.BODY_FONT
    normal.font.size = DocxPt(Typography.BODY_SIZE)
    normal.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.DARK))
    
    # Heading 1
    if "Heading 1" in styles:
        h1 = styles["Heading 1"]
        h1.font.name = Typography.HEADING_FONT
        h1.font.size = DocxPt(Typography.HEADING1_SIZE)
        h1.font.bold = True
        h1.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.PRIMARY))
    
    # Heading 2
    if "Heading 2" in styles:
        h2 = styles["Heading 2"]
        h2.font.name = Typography.HEADING_FONT
        h2.font.size = DocxPt(Typography.HEADING2_SIZE)
        h2.font.bold = True
        h2.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.SECONDARY))
    
    # Caption style
    if "Caption" in styles:
        caption = styles["Caption"]
        caption.font.name = Typography.BODY_FONT
        caption.font.size = DocxPt(Typography.CAPTION_SIZE)
        caption.font.italic = True
        caption.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
    
    return doc


def create_word_cover_page(
    doc,
    title: str = "Bibliometric Analysis Report",
    subtitle: str = "",
    author: str = "",
    date: str = None,
    dataset_info: Dict[str, Any] = None,
):
    """
    Create a modern cover page for Word document.
    
    Parameters
    ----------
    doc : docx.Document
        The document.
    title : str
        Main title.
    subtitle : str
        Subtitle or description.
    author : str
        Author name.
    date : str
        Report date (defaults to today).
    dataset_info : dict
        Dataset summary info.
    """
    from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    
    if date is None:
        date = datetime.now().strftime("%B %d, %Y")
    
    # Add some vertical spacing
    for _ in range(4):
        doc.add_paragraph()
    
    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(title)
    title_run.font.name = Typography.HEADING_FONT
    title_run.font.size = DocxPt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.PRIMARY))
    
    # Decorative line
    line_para = doc.add_paragraph()
    line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    line_run = line_para.add_run("─" * 40)
    line_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.SECONDARY))
    
    # Subtitle
    if subtitle:
        sub_para = doc.add_paragraph()
        sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_para.add_run(subtitle)
        sub_run.font.name = Typography.BODY_FONT
        sub_run.font.size = DocxPt(14)
        sub_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
    
    # Spacing
    for _ in range(3):
        doc.add_paragraph()
    
    # Dataset info box
    if dataset_info:
        info_para = doc.add_paragraph()
        info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        info_lines = []
        if "n_documents" in dataset_info:
            info_lines.append(f"📄 {dataset_info['n_documents']:,} Documents")
        if "year_range" in dataset_info:
            info_lines.append(f"📅 {dataset_info['year_range']}")
        if "database" in dataset_info:
            info_lines.append(f"🗄️ {dataset_info['database']}")
        
        for line in info_lines:
            run = info_para.add_run(line + "\n")
            run.font.name = Typography.BODY_FONT
            run.font.size = DocxPt(12)
            run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.DARK))
    
    # More spacing
    for _ in range(4):
        doc.add_paragraph()
    
    # Author and date at bottom
    if author:
        author_para = doc.add_paragraph()
        author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        author_run = author_para.add_run(f"Prepared by: {author}")
        author_run.font.name = Typography.BODY_FONT
        author_run.font.size = DocxPt(11)
        author_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
    
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_run = date_para.add_run(date)
    date_run.font.name = Typography.BODY_FONT
    date_run.font.size = DocxPt(11)
    date_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
    
    # Page break after cover
    doc.add_page_break()
    
    return doc


def style_word_table(table, style: str = "modern", has_header: bool = True):
    """
    Apply modern styling to a Word table.
    
    Parameters
    ----------
    table : docx.table.Table
        The table to style.
    style : str
        Style type: "modern", "minimal", "striped", "bordered"
    has_header : bool
        Whether first row is a header.
    """
    from docx.shared import Pt as DocxPt, RGBColor as DocxColor
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import OxmlElement, parse_xml
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    
    # Set table alignment
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    # Process each row
    for i, row in enumerate(table.rows):
        is_header = (i == 0 and has_header)
        is_alt_row = (i % 2 == 1) and not is_header
        
        for cell in row.cells:
            # Set cell shading
            shading_elm = OxmlElement("w:shd")
            
            if is_header:
                # Header styling
                shading_elm.set(qn("w:fill"), ModernPalette.TABLE_HEADER_BG.lstrip("#"))
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.WHITE))
                        run.font.size = DocxPt(10)
            elif is_alt_row and style in ("modern", "striped"):
                # Alternating row
                shading_elm.set(qn("w:fill"), ModernPalette.TABLE_ROW_ALT.lstrip("#"))
            else:
                shading_elm.set(qn("w:fill"), ModernPalette.WHITE.lstrip("#"))
            
            cell._tc.get_or_add_tcPr().append(shading_elm)
            
            # Cell padding and font
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                for run in paragraph.runs:
                    run.font.name = Typography.BODY_FONT
                    if not is_header:
                        run.font.size = DocxPt(9)
                        run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.DARK))
    
    return table


def add_word_toc(doc):
    """
    Add a table of contents to Word document.
    
    Parameters
    ----------
    doc : docx.Document
        The document.
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt as DocxPt, RGBColor as DocxColor
    
    # TOC heading
    toc_heading = doc.add_paragraph()
    toc_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    toc_run = toc_heading.add_run("Table of Contents")
    toc_run.font.name = Typography.HEADING_FONT
    toc_run.font.size = DocxPt(Typography.HEADING1_SIZE)
    toc_run.font.bold = True
    toc_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.PRIMARY))
    
    # Add TOC field
    paragraph = doc.add_paragraph()
    run = paragraph.add_run()
    
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = 'TOC \\o "1-3" \\h \\z \\u'
    
    fld_char_sep = OxmlElement("w:fldChar")
    fld_char_sep.set(qn("w:fldCharType"), "separate")
    
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    
    run._r.append(fld_char_begin)
    run._r.append(instr_text)
    run._r.append(fld_char_sep)
    run._r.append(fld_char_end)
    
    # Note for user
    note_para = doc.add_paragraph()
    note_run = note_para.add_run("(Right-click and select 'Update Field' to refresh)")
    note_run.font.size = DocxPt(8)
    note_run.font.italic = True
    note_run.font.color.rgb = DocxColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
    
    doc.add_page_break()
    
    return doc


# =============================================================================
# POWERPOINT STYLING
# =============================================================================

def create_modern_pptx_template():
    """
    Create a modern PowerPoint presentation template.
    
    Returns
    -------
    pptx.Presentation
        Styled presentation object.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PPTColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    return prs


def style_pptx_title_slide(slide, title: str, subtitle: str = ""):
    """
    Style a PowerPoint title slide.
    
    Parameters
    ----------
    slide : pptx.slide.Slide
        The slide to style.
    title : str
        Main title.
    subtitle : str
        Subtitle text.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PPTColor
    from pptx.enum.text import PP_ALIGN
    
    # Style title
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = Typography.HEADING_FONT
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.PRIMARY))
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Find and style subtitle placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            if subtitle:
                shape.text = subtitle
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = Typography.BODY_FONT
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.MEDIUM))
                paragraph.alignment = PP_ALIGN.CENTER
    
    return slide


def style_pptx_content_slide(slide, title: str = None):
    """
    Style a PowerPoint content slide.
    
    Parameters
    ----------
    slide : pptx.slide.Slide
        The slide to style.
    title : str
        Slide title.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as PPTColor
    from pptx.enum.text import PP_ALIGN
    
    if slide.shapes.title and title:
        title_shape = slide.shapes.title
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = Typography.HEADING_FONT
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.PRIMARY))
    
    return slide


def style_pptx_table(table, has_header: bool = True):
    """
    Apply modern styling to a PowerPoint table.
    
    Parameters
    ----------
    table : pptx.table.Table
        The table to style.
    has_header : bool
        Whether first row is a header.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as PPTColor
    from pptx.enum.text import PP_ALIGN
    
    for i, row in enumerate(table.rows):
        is_header = (i == 0 and has_header)
        is_alt = (i % 2 == 1) and not is_header
        
        for cell in row.cells:
            # Cell fill
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.TABLE_HEADER_BG))
            elif is_alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.TABLE_ROW_ALT))
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.WHITE))
            
            # Text styling
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = Typography.BODY_FONT
                paragraph.font.size = Pt(10)
                
                if is_header:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.WHITE))
                else:
                    paragraph.font.color.rgb = PPTColor(*ModernPalette.hex_to_rgb(ModernPalette.DARK))
    
    return table


# =============================================================================
# EXCEL STYLING
# =============================================================================

def style_excel_workbook(wb):
    """
    Apply modern styling to Excel workbook.
    
    Parameters
    ----------
    wb : openpyxl.Workbook
        The workbook to style.
    """
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment, NamedStyle
    
    # Create named styles
    try:
        # Header style
        header_style = NamedStyle(name="modern_header")
        header_style.font = Font(
            name=Typography.BODY_FONT,
            size=11,
            bold=True,
            color="FFFFFF"
        )
        header_style.fill = PatternFill(
            start_color=ModernPalette.TABLE_HEADER_BG.lstrip("#"),
            end_color=ModernPalette.TABLE_HEADER_BG.lstrip("#"),
            fill_type="solid"
        )
        header_style.alignment = Alignment(horizontal="center", vertical="center")
        wb.add_named_style(header_style)
    except ValueError:
        pass  # Style already exists
    
    try:
        # Data style
        data_style = NamedStyle(name="modern_data")
        data_style.font = Font(name=Typography.BODY_FONT, size=10)
        data_style.alignment = Alignment(horizontal="left", vertical="center")
        wb.add_named_style(data_style)
    except ValueError:
        pass
    
    try:
        # Alternating row style
        alt_style = NamedStyle(name="modern_alt")
        alt_style.font = Font(name=Typography.BODY_FONT, size=10)
        alt_style.fill = PatternFill(
            start_color=ModernPalette.TABLE_ROW_ALT.lstrip("#"),
            end_color=ModernPalette.TABLE_ROW_ALT.lstrip("#"),
            fill_type="solid"
        )
        alt_style.alignment = Alignment(horizontal="left", vertical="center")
        wb.add_named_style(alt_style)
    except ValueError:
        pass
    
    return wb


def style_excel_worksheet(ws, df, start_row: int = 1, has_header: bool = True):
    """
    Apply modern styling to an Excel worksheet containing a DataFrame.
    
    Parameters
    ----------
    ws : openpyxl.worksheet.Worksheet
        The worksheet.
    df : pandas.DataFrame
        The DataFrame that was written.
    start_row : int
        Row where data starts.
    has_header : bool
        Whether header row exists.
    """
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
    from openpyxl.utils import get_column_letter
    from openpyxl.formatting.rule import FormulaRule, ColorScaleRule
    
    # Define styles
    header_font = Font(name=Typography.BODY_FONT, size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(
        start_color=ModernPalette.TABLE_HEADER_BG.lstrip("#"),
        end_color=ModernPalette.TABLE_HEADER_BG.lstrip("#"),
        fill_type="solid"
    )
    
    data_font = Font(name=Typography.BODY_FONT, size=10)
    alt_fill = PatternFill(
        start_color=ModernPalette.TABLE_ROW_ALT.lstrip("#"),
        end_color=ModernPalette.TABLE_ROW_ALT.lstrip("#"),
        fill_type="solid"
    )
    
    thin_border = Border(
        bottom=Side(style="thin", color=ModernPalette.TABLE_BORDER.lstrip("#"))
    )
    
    n_cols = len(df.columns)
    n_rows = len(df) + (1 if has_header else 0)
    
    # Style header row
    if has_header:
        for col in range(1, n_cols + 1):
            cell = ws.cell(row=start_row, column=col)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
    
    # Style data rows
    data_start = start_row + (1 if has_header else 0)
    for row_idx in range(data_start, start_row + n_rows):
        is_alt = (row_idx - data_start) % 2 == 1
        for col in range(1, n_cols + 1):
            cell = ws.cell(row=row_idx, column=col)
            cell.font = data_font
            cell.border = thin_border
            if is_alt:
                cell.fill = alt_fill
    
    # Auto-fit column widths (approximate)
    for col_idx, col_name in enumerate(df.columns, 1):
        max_length = len(str(col_name))
        for row in df[col_name].astype(str):
            max_length = max(max_length, len(row))
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[get_column_letter(col_idx)].width = adjusted_width
    
    # Freeze header row
    if has_header:
        ws.freeze_panes = ws.cell(row=start_row + 1, column=1)
    
    return ws


def add_excel_cover_sheet(wb, title: str, info: Dict[str, Any] = None):
    """
    Add a cover sheet to Excel workbook.
    
    Parameters
    ----------
    wb : openpyxl.Workbook
        The workbook.
    title : str
        Report title.
    info : dict
        Additional info to display.
    """
    from openpyxl.styles import Font, PatternFill, Alignment
    
    # Create cover sheet at the beginning
    ws = wb.create_sheet("Cover", 0)
    
    # Title
    ws.merge_cells("B3:F3")
    title_cell = ws["B3"]
    title_cell.value = title
    title_cell.font = Font(name=Typography.HEADING_FONT, size=28, bold=True, 
                           color=ModernPalette.PRIMARY.lstrip("#"))
    title_cell.alignment = Alignment(horizontal="center")
    
    # Decorative line
    ws.merge_cells("B5:F5")
    line_cell = ws["B5"]
    line_cell.value = "─" * 60
    line_cell.font = Font(color=ModernPalette.SECONDARY.lstrip("#"))
    line_cell.alignment = Alignment(horizontal="center")
    
    # Info section
    if info:
        row = 7
        for key, value in info.items():
            ws.cell(row=row, column=2).value = key
            ws.cell(row=row, column=2).font = Font(name=Typography.BODY_FONT, bold=True)
            ws.cell(row=row, column=3).value = str(value)
            ws.cell(row=row, column=3).font = Font(name=Typography.BODY_FONT)
            row += 1
    
    # Date
    ws.cell(row=row + 2, column=2).value = "Generated:"
    ws.cell(row=row + 2, column=2).font = Font(name=Typography.BODY_FONT, bold=True)
    ws.cell(row=row + 2, column=3).value = datetime.now().strftime("%B %d, %Y")
    ws.cell(row=row + 2, column=3).font = Font(name=Typography.BODY_FONT)
    
    # Set column widths
    ws.column_dimensions["A"].width = 3
    ws.column_dimensions["B"].width = 20
    ws.column_dimensions["C"].width = 30
    
    return wb


def add_excel_conditional_formatting(ws, df, start_row: int = 1):
    """
    Add conditional formatting to highlight values.
    
    Parameters
    ----------
    ws : openpyxl.worksheet.Worksheet
        The worksheet.
    df : pandas.DataFrame
        The DataFrame.
    start_row : int
        Starting row for data.
    """
    from openpyxl.formatting.rule import ColorScaleRule, DataBarRule
    from openpyxl.utils import get_column_letter
    
    # Find numeric columns that might benefit from data bars
    numeric_cols = df.select_dtypes(include=["int64", "float64"]).columns
    
    for col_name in numeric_cols:
        if any(kw in col_name.lower() for kw in ["number", "count", "total", "citations"]):
            col_idx = list(df.columns).index(col_name) + 1
            col_letter = get_column_letter(col_idx)
            
            # Apply data bar
            data_bar_rule = DataBarRule(
                start_type="min",
                end_type="max",
                color=ModernPalette.SECONDARY.lstrip("#"),
                showValue=True,
                minLength=None,
                maxLength=None
            )
            
            range_str = f"{col_letter}{start_row + 1}:{col_letter}{start_row + len(df)}"
            ws.conditional_formatting.add(range_str, data_bar_rule)
    
    return ws


# =============================================================================
# LATEX STYLING
# =============================================================================

def get_latex_preamble(
    title: str = "Bibliometric Analysis Report",
    author: str = "",
    date: str = None,
    modern: bool = True,
) -> str:
    """
    Generate modern LaTeX preamble.
    
    Parameters
    ----------
    title : str
        Document title.
    author : str
        Author name.
    date : str
        Date string.
    modern : bool
        Use modern packages and styling.
    
    Returns
    -------
    str
        LaTeX preamble code.
    """
    if date is None:
        date = datetime.now().strftime("%B %d, %Y")
    
    preamble = r"""\documentclass[11pt,a4paper]{article}

% Modern packages
\usepackage[utf8]{inputenc}
\usepackage[T1]{fontenc}
\usepackage{lmodern}
\usepackage[margin=1in]{geometry}
\usepackage{graphicx}
\usepackage{booktabs}
\usepackage{longtable}
\usepackage{array}
\usepackage{xcolor}
\usepackage{hyperref}
\usepackage{fancyhdr}
\usepackage{titlesec}
\usepackage{tocloft}

% Color definitions
\definecolor{primary}{HTML}{2C3E50}
\definecolor{secondary}{HTML}{3498DB}
\definecolor{accent}{HTML}{E74C3C}
\definecolor{tableheader}{HTML}{2C3E50}
\definecolor{tablealt}{HTML}{F8F9FA}

% Modern section styling
\titleformat{\section}
  {\normalfont\Large\bfseries\color{primary}}
  {\thesection}{1em}{}
\titleformat{\subsection}
  {\normalfont\large\bfseries\color{secondary}}
  {\thesubsection}{1em}{}

% Hyperlink styling
\hypersetup{
    colorlinks=true,
    linkcolor=secondary,
    urlcolor=secondary,
    citecolor=secondary
}

% Header/footer
\pagestyle{fancy}
\fancyhf{}
\fancyhead[L]{\small\textcolor{primary}{""" + title.replace("&", r"\&") + r"""}}
\fancyhead[R]{\small\textcolor{primary}{\thepage}}
\renewcommand{\headrulewidth}{0.4pt}
\renewcommand{\headrule}{\hbox to\headwidth{\color{secondary}\leaders\hrule height \headrulewidth\hfill}}

% Table styling
\newcommand{\tableheaderrow}{\rowcolor{tableheader}}
\newcommand{\tablealtrow}{\rowcolor{tablealt}}

% Document info
\title{\textcolor{primary}{\textbf{""" + title.replace("&", r"\&") + r"""}}}
\author{""" + author.replace("&", r"\&") + r"""}
\date{""" + date + r"""}

\begin{document}
"""
    
    return preamble


def get_latex_cover_page(
    title: str,
    subtitle: str = "",
    author: str = "",
    date: str = None,
    info: Dict[str, Any] = None,
) -> str:
    """
    Generate LaTeX cover page.
    
    Parameters
    ----------
    title : str
        Main title.
    subtitle : str
        Subtitle.
    author : str
        Author name.
    date : str
        Date.
    info : dict
        Additional info.
    
    Returns
    -------
    str
        LaTeX code for cover page.
    """
    if date is None:
        date = datetime.now().strftime("%B %d, %Y")
    
    # Escape special characters
    title = title.replace("&", r"\&").replace("_", r"\_")
    subtitle = subtitle.replace("&", r"\&").replace("_", r"\_") if subtitle else ""
    author = author.replace("&", r"\&").replace("_", r"\_") if author else ""
    
    cover = r"""
\begin{titlepage}
\centering
\vspace*{3cm}

{\Huge\bfseries\textcolor{primary}{""" + title + r"""}\par}
\vspace{0.5cm}
{\color{secondary}\rule{0.6\textwidth}{0.5pt}\par}
"""
    
    if subtitle:
        cover += r"""
\vspace{0.5cm}
{\Large\textcolor{gray}{""" + subtitle + r"""}\par}
"""
    
    cover += r"""
\vspace{2cm}
"""
    
    if info:
        cover += r"""
\begin{tabular}{rl}
"""
        for key, value in info.items():
            key_escaped = str(key).replace("&", r"\&").replace("_", r"\_")
            value_escaped = str(value).replace("&", r"\&").replace("_", r"\_")
            cover += rf"\textbf{{{key_escaped}:}} & {value_escaped} \\" + "\n"
        cover += r"""
\end{tabular}
"""
    
    cover += r"""
\vfill

"""
    if author:
        cover += r"""{\large Prepared by: """ + author + r"""}\par
\vspace{0.3cm}
"""
    
    cover += r"""{\large """ + date + r"""}\par

\end{titlepage}

\tableofcontents
\newpage
"""
    
    return cover


def latex_table_modern(df, caption: str = "", label: str = "") -> str:
    """
    Generate modern LaTeX table from DataFrame.
    
    Parameters
    ----------
    df : pandas.DataFrame
        The DataFrame.
    caption : str
        Table caption.
    label : str
        Table label for referencing.
    
    Returns
    -------
    str
        LaTeX table code.
    """
    import pandas as pd
    
    n_cols = len(df.columns)
    col_spec = "l" + "r" * (n_cols - 1)
    
    # Escape column names and values
    def escape_latex(s):
        if pd.isna(s):
            return ""
        s = str(s)
        for char in ["&", "%", "$", "#", "_", "{", "}"]:
            s = s.replace(char, f"\\{char}")
        return s
    
    table = r"""
\begin{table}[htbp]
\centering
\small
\begin{tabular}{""" + col_spec + r"""}
\toprule
\rowcolor{tableheader}
"""
    
    # Header row
    headers = [r"\textcolor{white}{\textbf{" + escape_latex(col) + "}}" for col in df.columns]
    table += " & ".join(headers) + r" \\" + "\n"
    table += r"\midrule" + "\n"
    
    # Data rows with alternating colors
    for i, (_, row) in enumerate(df.iterrows()):
        if i % 2 == 1:
            table += r"\rowcolor{tablealt}" + "\n"
        values = [escape_latex(val) for val in row.values]
        table += " & ".join(values) + r" \\" + "\n"
    
    table += r"""\bottomrule
\end{tabular}
"""
    
    if caption:
        caption_escaped = caption.replace("&", r"\&").replace("_", r"\_")
        table += r"\caption{" + caption_escaped + "}\n"
    
    if label:
        table += r"\label{" + label + "}\n"
    
    table += r"\end{table}" + "\n"
    
    return table


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def format_number(value, decimals: int = 0) -> str:
    """Format number with thousands separator."""
    if isinstance(value, (int, float)):
        if decimals == 0:
            return f"{int(value):,}"
        return f"{value:,.{decimals}f}"
    return str(value)


def add_trend_indicator(current, previous) -> str:
    """Add trend indicator arrow."""
    if current > previous:
        return "↑"
    elif current < previous:
        return "↓"
    return "→"


def truncate_text(text: str, max_length: int = 50) -> str:
    """Truncate text with ellipsis."""
    if len(text) <= max_length:
        return text
    return text[:max_length - 3] + "..."
