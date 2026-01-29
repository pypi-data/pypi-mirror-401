"""
report_export_utils_optimized

High-level exporters for Excel, PowerPoint, Word, and LaTeX reports driven by
a single Excel template.

Template columns (fixed)
------------------------
The template workbook is assumed to have at least these columns
(with *exactly* these names in Excel):

- Level 1
- Level 2
- Item Type           (Table / Plot / Text / Description)
- Data Attr
- Text
- Path
- Plot Filename
- Head                (used in Word / PPTX / TeX; ignored in Excel)
- Tab Color
- TOC Label Color
- Caption
- Description         (section-level text attribute for Word / TeX)
- Narrative           (text attribute, especially for PPTX)
- Page Break After
- Table Style

Item Type semantics
-------------------
- "Table"       → DataFrame from `Data Attr`.
- "Plot"        → Image from `Path` + `Plot Filename`.
- "Text"        → Literal text from the "Text" column (or Narrative attr in PPTX fallback).
- "Description" → Text from an attribute on `self` named in `Data Attr`.

Head semantics (Word / PPTX / TeX)
----------------------------------
For each table item:

- If Head is empty / NaN:
    * If the exporter has `top_n` set (e.g. top_n=5), show the first `top_n` rows.
    * If `top_n` is None, show the full table.
- If Head == 0 → show the full table.
- If Head > 0 → show the first Head rows.

Excel ignores Head and always prints the full DataFrame.

Notes
-----
- Strings use double quotes.
- All public functions include docstrings and type hints.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

import pandas as pd

# -----------------------------
# Common paths
# -----------------------------

FD: Path = Path(__file__).resolve().parent
ICONS_DIR: Path = FD / "additional files" / "icons"

# -----------------------------
# OpenPyXL
# -----------------------------

from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.worksheet.hyperlink import Hyperlink

# -----------------------------
# PowerPoint
# -----------------------------

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTColor

# -----------------------------
# Word
# -----------------------------

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.shared import RGBColor as DocxColor


# -----------------------------
# Color helpers
# -----------------------------

COLOR_MAP: Dict[str, str] = {
    "LightGreen": "90EE90",
    "LightBlue": "ADD8E6",
    "Khaki": "F0E68C",
    "Gold": "FFD700",
    "Orange": "FFA500",
    "Pink": "FFC0CB",
    "Red": "FF0000",
    "Blue": "0000FF",
    "Green": "008000",
    "Purple": "800080",
    "Yellow": "FFFF00",
    "Gray": "808080",
    "Cyan": "00FFFF",
    "Magenta": "FF00FF",
    "White": "FFFFFF",
    "Black": "000000",
}


def _normalize_hex(
    rgb_or_name: str,
    default: str = "FFFFFF",
) -> str:
    """
    Return a 6-digit upper-case hex string from a color name or hex string.

    Parameters
    ----------
    rgb_or_name : str
        Color name (key in COLOR_MAP) or hex string (with or without "#").
    default : str, default "FFFFFF"
        Fallback hex string if the input is not recognized.

    Returns
    -------
    str
        Normalized 6-character upper-case hex code (no "#").
    """
    if not rgb_or_name:
        return default

    s = str(rgb_or_name).strip()
    if s in COLOR_MAP:
        return COLOR_MAP[s]

    if s.startswith("#"):
        s = s[1:]

    s = s.upper()
    if len(s) == 6 and all(c in "0123456789ABCDEF" for c in s):
        return s

    return default


def _set_internal_link(
    cell,
    sheet_name: str,
    cell_ref: str = "A1",
) -> None:
    """
    Make `cell` a clickable internal link to `'sheet_name'!cell_ref`.

    Parameters
    ----------
    cell
        OpenPyXL cell object.
    sheet_name : str
        Target sheet name.
    cell_ref : str, default "A1"
        Target cell reference on the target sheet.
    """
    target = f"'{sheet_name}'!{cell_ref}"
    cell.hyperlink = Hyperlink(ref=cell.coordinate, location=target)
    cell.font = Font(color="0000FF", underline="single")


def _set_external_file_link(
    cell,
    file_path: Path,
) -> None:
    """
    Make `cell` a clickable hyperlink to a local file path.

    Parameters
    ----------
    cell
        OpenPyXL cell object.
    file_path : Path
        File path to link to.
    """
    try:
        url = file_path.as_uri()
    except Exception:
        url = "file://" + str(file_path.resolve()).replace("\\", "/")

    cell.hyperlink = Hyperlink(ref=cell.coordinate, target=url)
    cell.font = Font(color="0000FF", underline="single")


def _fit_image(
    img: XLImage,
    *,
    scale: float = 1.0,
    max_w: int = 640,
    max_h: int = 360,
) -> None:
    """
    Scale `img` by `scale` then constrain to max width/height (pixels).

    Parameters
    ----------
    img : XLImage
        OpenPyXL image to resize.
    scale : float, default 1.0
        Multiplicative scale factor before constraining by max_w / max_h.
    max_w : int, default 640
        Maximum width in pixels.
    max_h : int, default 360
        Maximum height in pixels.
    """
    w = max(1, int(img.width * scale))
    h = max(1, int(img.height * scale))

    if w > max_w or h > max_h:
        r = min(max_w / float(w), max_h / float(h))
        w = max(1, int(w * r))
        h = max(1, int(h * r))

    img.width, img.height = w, h


# -----------------------------
# Path helpers
# -----------------------------

def _resolve_image_path(
    self: Any,
    tpath: str,
    filename: str,
    results_base: Optional[Union[str, Path]] = None,
) -> Optional[Path]:
    """
    Resolve image path from (Path, Filename) using several base folders.

    The search order is:
    1. Absolute `Path` if `tpath` is absolute.
    2. Explicit `results_base`, if provided.
    3. `self.res_folder`, if present.
    4. `self.res_folder/plots`, if present.
    5. Module folder (FD).
    6. Current working directory.
    7. Current working directory/plots.

    Also tries filename variations:
    - Original filename
    - With spaces replaced by underscores
    - With underscores replaced by spaces
    - Without extension (trying .png, .jpg, .jpeg, .svg, .pdf)

    Parameters
    ----------
    self : Any
        Object that may have a `res_folder` attribute.
    tpath : str
        Template path cell value (may be empty / relative / absolute).
    filename : str
        File name of the image.
    results_base : str or Path, optional
        Explicit base folder for relative paths.

    Returns
    -------
    Path or None
        Resolved existing path, or None if nothing was found.
    """
    filename = str(filename or "").strip()
    tpath = str(tpath or "").strip()

    if not filename:
        return None

    p = Path(tpath) if tpath else None

    # Generate filename variations to try
    def _filename_variations(fname: str) -> List[str]:
        variations = [fname]
        # Use Path to split extension instead of os.path.splitext
        p_fname = Path(fname)
        base = p_fname.stem
        ext = p_fname.suffix
        
        # Try with spaces/underscores swapped
        if " " in base:
            variations.append(base.replace(" ", "_") + ext)
        if "_" in base:
            variations.append(base.replace("_", " ") + ext)
        
        # If no extension, try common image extensions
        if not ext:
            for e in [".png", ".jpg", ".jpeg", ".svg", ".pdf"]:
                variations.append(fname + e)
                if " " in fname:
                    variations.append(fname.replace(" ", "_") + e)
                if "_" in fname:
                    variations.append(fname.replace("_", " ") + e)
        
        return variations

    filename_variants = _filename_variations(filename)

    # 1) Absolute base path
    if p and p.is_absolute():
        for fn in filename_variants:
            cand = (p / fn).resolve()
            if cand.exists():
                return cand

    # Collect candidate bases
    bases: List[Path] = []

    # 2) Explicit results_base
    if results_base is not None:
        try:
            rb = Path(results_base)
            bases.append(rb)
            bases.append(rb / "plots")  # Also check plots subfolder
        except Exception:
            pass

    # 3) self.res_folder
    base_self = getattr(self, "res_folder", None)
    if base_self:
        bases.append(Path(base_self))
        bases.append(Path(base_self) / "plots")  # Also check plots subfolder

    # 4) Module folder
    bases.append(FD)

    # 5) Current working directory
    bases.append(Path.cwd())
    bases.append(Path.cwd() / "plots")  # Also check plots subfolder

    # Try combinations of bases and filename variants
    for base in bases:
        rel = (p or Path(""))
        for fn in filename_variants:
            cand = (base / rel / fn).resolve()
            if cand.exists():
                return cand

    return None


def _resolve_base(
    self: Any,
    results_base: Optional[Union[str, Path]] = None,
) -> Path:
    """
    Return preferred base folder for relative paths.

    Parameters
    ----------
    self : Any
        Object that may have `res_folder` attribute.
    results_base : str or Path, optional
        Explicit override for base folder.

    Returns
    -------
    Path
        Resolved base path.
    """
    if results_base is not None:
        try:
            return Path(results_base).resolve()
        except Exception:
            return FD

    base = getattr(self, "res_folder", None)
    return Path(base).resolve() if base else FD


def _resolve_path(
    self: Any,
    maybe_path: Union[str, Path],
    results_base: Optional[Union[str, Path]] = None,
) -> Path:
    """
    Resolve a path that might be absolute or relative.

    Parameters
    ----------
    self : Any
        Object that may have `res_folder` attribute.
    maybe_path : str or Path
        Path that may be absolute or relative.
    results_base : str or Path, optional
        Explicit base folder for relative paths.

    Returns
    -------
    Path
        Resolved path.
    """
    p = Path(str(maybe_path))
    if p.is_absolute():
        return p
    return (_resolve_base(self, results_base) / p).resolve()


# -----------------------------
# Template / data helpers
# -----------------------------

# Required columns for template validation
REQUIRED_TEMPLATE_COLUMNS = [
    "Level 1",
    "Level 2", 
    "Item Type",
    "Data Attr",
]

OPTIONAL_TEMPLATE_COLUMNS = [
    "Plot Method",
    "Plot Filename",
    "Caption",
    "Description",
    "Tab Color",
    "TOC Color",
    "Page Break",
    "Table Style",
    "Columns (docx/pptx/tex)",
    "Columns (xlsx)",
    "Head",
    "xlsx_only",
]


def validate_template(
    template_path: Union[str, Path],
    sheet_name: str = "all",
    raise_on_error: bool = False,
) -> Dict[str, Any]:
    """
    Validate that a template file has the required structure.
    
    Parameters
    ----------
    template_path : str or Path
        Path to the template workbook.
    sheet_name : str, default "all"
        Name of the sheet to validate.
    raise_on_error : bool, default False
        If True, raise ValueError on validation failure.
        
    Returns
    -------
    dict
        Validation result with keys:
        - 'valid': bool
        - 'errors': list of error messages
        - 'warnings': list of warning messages
        - 'columns_found': list of columns in template
        - 'sheets_found': list of sheets in workbook
        
    Examples
    --------
    >>> result = validate_template("my_template.xlsx")
    >>> if not result['valid']:
    ...     print("Template errors:", result['errors'])
    """
    result = {
        'valid': True,
        'errors': [],
        'warnings': [],
        'columns_found': [],
        'sheets_found': [],
    }
    
    template_path = Path(template_path)
    
    # Check file exists
    if not template_path.exists():
        result['valid'] = False
        result['errors'].append(f"Template file not found: {template_path}")
        if raise_on_error:
            raise ValueError(result['errors'][0])
        return result
    
    # Check file extension
    if template_path.suffix.lower() not in ['.xlsx', '.xls']:
        result['valid'] = False
        result['errors'].append(f"Template must be an Excel file (.xlsx or .xls)")
        if raise_on_error:
            raise ValueError(result['errors'][0])
        return result
    
    try:
        # Read available sheets
        xlsx = pd.ExcelFile(template_path)
        result['sheets_found'] = xlsx.sheet_names
        
        # Check if requested sheet exists
        if sheet_name not in xlsx.sheet_names:
            result['valid'] = False
            result['errors'].append(
                f"Sheet '{sheet_name}' not found. Available: {xlsx.sheet_names}"
            )
            if raise_on_error:
                raise ValueError(result['errors'][0])
            return result
        
        # Read the sheet
        df = pd.read_excel(xlsx, sheet_name=sheet_name)
        df.columns = [str(c).strip() for c in df.columns]
        result['columns_found'] = list(df.columns)
        
        # Check required columns
        missing_required = [
            col for col in REQUIRED_TEMPLATE_COLUMNS 
            if col not in df.columns
        ]
        
        if missing_required:
            result['valid'] = False
            result['errors'].append(
                f"Missing required columns: {missing_required}"
            )
        
        # Check for optional columns (warnings only)
        missing_optional = [
            col for col in OPTIONAL_TEMPLATE_COLUMNS
            if col not in df.columns
        ]
        
        if missing_optional:
            result['warnings'].append(
                f"Missing optional columns (will use defaults): {missing_optional[:5]}"
                + (f" and {len(missing_optional) - 5} more" if len(missing_optional) > 5 else "")
            )
        
        # Check for empty template
        if len(df) == 0:
            result['warnings'].append("Template sheet is empty")
        
        # Check Item Type values
        if 'Item Type' in df.columns:
            valid_types = {'Table', 'Plot', 'Text', 'Description'}
            invalid_types = set(df['Item Type'].dropna().unique()) - valid_types
            if invalid_types:
                result['warnings'].append(
                    f"Unknown Item Types (will be ignored): {invalid_types}"
                )
        
    except Exception as e:
        result['valid'] = False
        result['errors'].append(f"Error reading template: {type(e).__name__}: {str(e)}")
    
    if raise_on_error and not result['valid']:
        raise ValueError("; ".join(result['errors']))
    
    return result


def _read_template(
    template_path: Union[str, Path],
    sheet_name: str,
) -> pd.DataFrame:
    """
    Read the Excel template and strip whitespace from header cells.

    Parameters
    ----------
    template_path : str or Path
        Path to the template workbook.
    sheet_name : str
        Name of the sheet inside the template.

    Returns
    -------
    pandas.DataFrame
        Template definition dataframe.
    """
    df = pd.read_excel(template_path, sheet_name=sheet_name)
    df.columns = [str(c).strip() for c in df.columns]
    return df


def _get_df_attr(
    self: Any,
    name: str,
) -> Optional[pd.DataFrame]:
    """
    Fetch a DataFrame attribute on `self` by name if present and valid.

    Parameters
    ----------
    self : Any
        Object that may hold pandas DataFrame attributes.
    name : str
        Attribute name.

    Returns
    -------
    pandas.DataFrame or None
        DataFrame attribute, or None if not found / not a DataFrame.
    """
    if not name:
        return None
    obj = getattr(self, name, None)
    return obj if isinstance(obj, pd.DataFrame) else None


def _round_numeric(
    df: pd.DataFrame,
    ndigits: int = 3,
) -> pd.DataFrame:
    """
    Round floating numeric columns to `ndigits` (copy).

    Parameters
    ----------
    df : pandas.DataFrame
        Input DataFrame.
    ndigits : int, default 3
        Number of decimal places.

    Returns
    -------
    pandas.DataFrame
        DataFrame with numeric columns rounded.
    """
    if df is None or df.empty:
        return df

    out = df.copy()
    float_cols: List[str] = [
        c
        for c in out.columns
        if pd.api.types.is_float_dtype(out[c])
        or (
            pd.api.types.is_numeric_dtype(out[c])
            and not pd.api.types.is_integer_dtype(out[c])
        )
    ]
    if float_cols:
        out[float_cols] = out[float_cols].round(ndigits)
    return out


def _autosize_columns(
    ws,
) -> None:
    """
    Autosize columns on an openpyxl worksheet based on cell text length.

    Parameters
    ----------
    ws
        OpenPyXL worksheet.
    """
    for col_cells in ws.columns:
        try:
            idx = col_cells[0].column
        except Exception:
            continue
        max_len = 0
        for cell in col_cells:
            if cell.value is not None:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[get_column_letter(idx)].width = max_len + 2


def _parse_head_value(
    head_value: Any,
) -> Optional[int]:
    """
    Parse a 'Head' cell into an integer or None.

    Parameters
    ----------
    head_value : Any
        Cell value from the Head column.

    Returns
    -------
    int or None
        Parsed integer, or None if not meaningful.
    """
    if head_value is None:
        return None

    try:
        if isinstance(head_value, float) and pd.isna(head_value):
            return None
    except Exception:
        pass

    if isinstance(head_value, str):
        s = head_value.strip()
        if not s:
            return None
        try:
            return int(s)
        except Exception:
            return None

    if isinstance(head_value, (int, float)):
        try:
            return int(head_value)
        except Exception:
            return None

    return None


def _apply_head(
    df: Optional[pd.DataFrame],
    head_value: Any,
    default_top_n: Optional[int],
) -> Optional[pd.DataFrame]:
    """
    Return a possibly truncated dataframe according to Head and top_n.

    Parameters
    ----------
    df : pandas.DataFrame or None
        Input DataFrame.
    head_value : Any
        Cell from the Head column.
    default_top_n : int or None
        Default number of rows if Head is empty.

    Returns
    -------
    pandas.DataFrame or None
        Possibly truncated DataFrame.
    """
    if df is None or df.empty:
        return df

    k = _parse_head_value(head_value)
    if k is None:
        if default_top_n is None:
            return df
        return df.head(int(default_top_n))

    if k <= 0:
        return df

    return df.head(k)


def _filter_columns(
    df: Optional[pd.DataFrame],
    columns_spec: str,
) -> Optional[pd.DataFrame]:
    """
    Filter DataFrame to only include specified columns.
    
    Parameters
    ----------
    df : pandas.DataFrame or None
        Input DataFrame.
    columns_spec : str
        Semicolon-separated list of column names to keep.
        If empty, returns all columns.
        
    Returns
    -------
    pandas.DataFrame or None
        DataFrame with only specified columns (if they exist).
    """
    if df is None or df.empty:
        return df
    
    columns_spec = _clean_str(columns_spec)
    if not columns_spec:
        return df
    
    # Parse column names (semicolon or comma separated)
    if ';' in columns_spec:
        requested_cols = [c.strip() for c in columns_spec.split(';') if c.strip()]
    else:
        requested_cols = [c.strip() for c in columns_spec.split(',') if c.strip()]
    
    if not requested_cols:
        return df
    
    # Keep only columns that exist
    available_cols = [c for c in requested_cols if c in df.columns]
    
    if not available_cols:
        return df
    
    return df[available_cols].copy()


def _is_xlsx_only(row: pd.Series) -> bool:
    """
    Check if a template row should only be included in xlsx output.
    
    Parameters
    ----------
    row : pandas.Series
        Template row.
        
    Returns
    -------
    bool
        True if xlsx_only flag is set.
    """
    xlsx_only = row.get("xlsx_only", False)
    if isinstance(xlsx_only, bool):
        return xlsx_only
    if isinstance(xlsx_only, str):
        return xlsx_only.strip().lower() in ('true', 'yes', '1', 'x')
    try:
        return bool(xlsx_only)
    except:
        return False


def _infer_item_type(
    row: pd.Series,
) -> Optional[str]:
    """
    Infer item type from the "Item Type" column.

    Parameters
    ----------
    row : pandas.Series
        Template row.

    Returns
    -------
    str or None
        "table", "plot", "text", "description", or None if unrecognized.
    """
    raw = str(row.get("Item Type", "")).strip().lower()
    if raw in {"table", "plot", "text", "description"}:
        return raw
    return None


def _is_missing(
    value: Any,
) -> bool:
    """
    Return True if `value` is considered missing (None/NaN/pandas.NA/empty-like).

    Parameters
    ----------
    value : Any
        Value to test.

    Returns
    -------
    bool
        True if missing, False otherwise.
    """
    if value is None:
        return True

    if isinstance(value, str):
        s = value.strip()
        if not s:
            return True
        if s.lower() in {"nan", "na", "none", "<na>"}:
            return True
        return False

    try:
        return bool(pd.isna(value))
    except Exception:
        return False


def _clean_str(
    value: Any,
) -> str:
    """
    Convert a value to a clean string, treating NaN/NA as empty.

    Parameters
    ----------
    value : Any
        Value to convert.

    Returns
    -------
    str
        Cleaned string ("" if missing).
    """
    if _is_missing(value):
        return ""
    return str(value).strip()


def _df_display(
    df: Optional[pd.DataFrame],
    ndigits: int = 3,
) -> Optional[pd.DataFrame]:
    """
    Prepare a DataFrame for export to Excel/PPTX/Word/TeX.

    This function:
    - Rounds numeric columns to `ndigits`.
    - Converts all values to object dtype.
    - Replaces missing values (NaN/pandas.NA) with empty strings.

    Parameters
    ----------
    df : pandas.DataFrame or None
        Input DataFrame.
    ndigits : int, default 3
        Decimal places for numeric columns.

    Returns
    -------
    pandas.DataFrame or None
        Cleaned DataFrame, or None if input was None or empty.
    """
    if df is None or df.empty:
        return df

    out = _round_numeric(df, ndigits).copy()
    out = out.astype("object")
    mask = pd.isna(out)
    out[mask] = ""
    return out


def _flag_true(
    value: Any,
) -> bool:
    """
    Interpret a template flag cell as boolean True/False.

    Parameters
    ----------
    value : Any
        Cell value to interpret.

    Returns
    -------
    bool
        True for "TRUE", "Yes", "1", etc.; False otherwise.
    """
    s = _clean_str(value).lower()
    return s in {"true", "yes", "1", "y"}


# -----------------------------
# EXCEL
# -----------------------------

def save_excel_report_from_template(
    self: Any,
    output_path: str = "bibliometric_report.xlsx",
    template_path: Union[str, Path] = FD
    / "additional files"
    / "template for report output.xlsx",
    template_sheet: str = "all",
    striped_rows: bool = True,
    show_icons: bool = False,
    autofit: bool = True,
    show_gridlines: bool = False,
    freeze_header: bool = True,
    image_scale: float = 0.75,
    image_max_width_px: int = 640,
    image_max_height_px: int = 360,
    results_base: Union[str, Path] = "results",
    # New modern styling options
    modern_style: bool = True,
    include_cover_sheet: bool = True,
    report_title: str = "Bibliometric Analysis Report",
    data_bars: bool = True,
) -> Path:
    """
    Generate an Excel report from a template definition sheet.

    Excel always prints full DataFrames (Head is ignored here).

    Parameters
    ----------
    self : Any
        Object holding DataFrame attributes and optionally `res_folder`.
    output_path : str, default "bibliometric_report.xlsx"
        Output Excel file name or path (relative to `results_base` / `self.res_folder`).
    template_path : str or Path, default "<FD>/additional files/template for report output.xlsx"
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    striped_rows : bool, default False
        Apply striped fill to rows in the TOC.
    show_icons : bool, default False
        Show small icons for item types in the TOC.
    autofit : bool, default True
        Autosize columns.
    show_gridlines : bool, default True
        Show gridlines on sheets.
    freeze_header : bool, default True
        Freeze header row of each content sheet.
    image_scale : float, default 0.75
        Initial scaling factor for embedded images.
    image_max_width_px : int, default 640
        Maximum width for embedded images in pixels.
    image_max_height_px : int, default 360
        Maximum height for embedded images in pixels.
    results_base : str or Path, default "results"
        Base folder for relative paths and output.
    modern_style : bool, default True
        Apply modern table styling with colored headers and zebra stripes.
    include_cover_sheet : bool, default True
        Add a cover sheet with report info.
    report_title : str, default "Bibliometric Analysis Report"
        Title for the cover sheet.
    data_bars : bool, default True
        Add data bars to numeric columns.

    Returns
    -------
    Path
        Path to the saved Excel file.
    """
    from openpyxl.styles import Border, Side, Alignment
    from openpyxl.formatting.rule import DataBarRule
    
    # Modern color constants
    HEADER_BG = "2C3E50"  # Dark blue-gray
    HEADER_FG = "FFFFFF"  # White
    ALT_ROW_BG = "F8F9FA"  # Light gray
    ACCENT_COLOR = "3498DB"  # Blue
    
    template_path = _resolve_path(self, template_path, results_base)
    df_template = _read_template(template_path, template_sheet)

    wb = Workbook()
    
    # Add cover sheet if enabled
    if include_cover_sheet:
        cover_ws = wb.active
        cover_ws.title = "Cover"
        cover_ws.sheet_properties.tabColor = HEADER_BG
        
        # Title
        cover_ws.merge_cells("B4:F4")
        title_cell = cover_ws["B4"]
        title_cell.value = report_title
        title_cell.font = Font(name="Calibri Light", size=28, bold=True, color=HEADER_BG)
        title_cell.alignment = Alignment(horizontal="center")
        
        # Decorative line
        cover_ws.merge_cells("B6:F6")
        line_cell = cover_ws["B6"]
        line_cell.value = "─" * 60
        line_cell.font = Font(color=ACCENT_COLOR)
        line_cell.alignment = Alignment(horizontal="center")
        
        # Summary stats
        summary_stats = _get_summary_stats(self)
        row = 8
        for key, value in summary_stats.items():
            cover_ws.cell(row=row, column=2).value = key
            cover_ws.cell(row=row, column=2).font = Font(name="Calibri", bold=True, size=11)
            cover_ws.cell(row=row, column=3).value = str(value)
            cover_ws.cell(row=row, column=3).font = Font(name="Calibri", size=11)
            row += 1
        
        # Date
        from datetime import datetime
        cover_ws.cell(row=row + 1, column=2).value = "Generated:"
        cover_ws.cell(row=row + 1, column=2).font = Font(name="Calibri", bold=True, size=11)
        cover_ws.cell(row=row + 1, column=3).value = datetime.now().strftime("%B %d, %Y")
        cover_ws.cell(row=row + 1, column=3).font = Font(name="Calibri", size=11)
        
        # Column widths
        cover_ws.column_dimensions["A"].width = 3
        cover_ws.column_dimensions["B"].width = 25
        cover_ws.column_dimensions["C"].width = 35
        
        cover_ws.sheet_view.showGridLines = False
        
        # Create new TOC sheet
        toc_ws = wb.create_sheet("Table of Contents")
    else:
        toc_ws = wb.active
        toc_ws.title = "Table of Contents"
    
    toc_ws.freeze_panes = "A3"
    toc_tab_color = _normalize_hex("LightBlue")
    toc_ws.sheet_properties.tabColor = toc_tab_color

    header = ["Level 1", "Level 2", "Sheet Name"] + (["Icon"] if show_icons else [])
    toc_ws.append(header)
    
    # Modern header styling for TOC
    for col_idx in range(1, len(header) + 1):
        cell = toc_ws.cell(row=1, column=col_idx)
        if modern_style:
            cell.fill = PatternFill(start_color=HEADER_BG, end_color=HEADER_BG, fill_type="solid")
            cell.font = Font(bold=True, color=HEADER_FG, name="Calibri")
        else:
            cell.fill = PatternFill(start_color="D9EAD3", end_color="D9EAD3", fill_type="solid")
            cell.font = Font(bold=True)

    def _insert_icon(ws, row_idx: int, item_type: str) -> None:
        """
        Insert a small icon image into TOC row for visual cue.

        Parameters
        ----------
        ws
            TOC worksheet.
        row_idx : int
            Row index to place the icon.
        item_type : str
            Item type ("table", "plot", "text", "description").
        """
        icon_name: Optional[str] = None
        if item_type == "table":
            icon_name = "table.png"
        elif item_type == "plot":
            icon_name = "plot.png"
        elif item_type == "text":
            icon_name = "text.png"
        elif item_type == "description":
            icon_name = "info.png"

        if icon_name is None:
            return

        icon_path = ICONS_DIR / icon_name
        if icon_path.exists():
            img = XLImage(str(icon_path))
            img.width = int(img.width * 0.5)
            img.height = int(img.height * 0.5)
            ws.add_image(img, f"D{row_idx}")

    row_idx = 2

    for _, row in df_template.iterrows():
        level1 = _clean_str(row.get("Level 1", ""))
        level2 = _clean_str(row.get("Level 2", ""))
        item_type = _infer_item_type(row)
        if not item_type:
            continue
        it = item_type.lower()

        df_obj: Optional[pd.DataFrame] = None
        img_path: Optional[Path] = None
        text_value: str = ""
        has_content = False

        if it == "table":
            data_attr = _clean_str(row.get("Data Attr", ""))
            df_obj = _get_df_attr(self, data_attr)
            has_content = isinstance(df_obj, pd.DataFrame) and not df_obj.empty

        elif it == "plot":
            tpath = _clean_str(row.get("Path", ""))
            plot_filename = _clean_str(row.get("Plot Filename", ""))
            if plot_filename:
                img_path = _resolve_image_path(self, tpath, plot_filename, results_base)
                has_content = img_path is not None

        elif it == "text":
            text_value = _clean_str(row.get("Text", ""))
            has_content = bool(text_value)

        elif it == "description":
            data_attr = _clean_str(row.get("Data Attr", ""))
            if data_attr and hasattr(self, data_attr):
                desc = getattr(self, data_attr)
                if isinstance(desc, str):
                    text_value = desc
                    has_content = bool(text_value.strip())

        if not has_content:
            continue

        sheet_name = (level2 or level1 or "Sheet")[:31]
        tab_color = _normalize_hex(_clean_str(row.get("Tab Color", "")))
        toc_color = _normalize_hex(_clean_str(row.get("TOC Label Color", "")))

        ws = wb.create_sheet(sheet_name)
        ws.sheet_properties.tabColor = tab_color

        # Add home link with icon
        home_cell = ws["A1"]
        home_cell.value = "🏠 Home"
        _set_internal_link(home_cell, toc_ws.title, "A1")
        try:
            if modern_style:
                home_cell.fill = PatternFill(
                    start_color=ACCENT_COLOR,
                    end_color=ACCENT_COLOR,
                    fill_type="solid",
                )
                home_cell.font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
            else:
                home_cell.fill = PatternFill(
                    start_color="FFF2CC",
                    end_color="FFF2CC",
                    fill_type="solid",
                )
        except Exception:
            pass

        start_row = 3

        toc_ws.cell(row=row_idx, column=1, value=level1)
        toc_ws.cell(row=row_idx, column=2, value=level2)
        toc_cell = toc_ws.cell(row=row_idx, column=3, value=sheet_name)
        _set_internal_link(toc_cell, sheet_name, "A1")
        if toc_color:
            toc_cell.fill = PatternFill(
                start_color=toc_color,
                end_color=toc_color,
                fill_type="solid",
            )

        if show_icons:
            _insert_icon(toc_ws, row_idx, it)

        row_idx += 1

        # Content
        if it == "table" and df_obj is not None:
            if freeze_header:
                ws.freeze_panes = f"A{start_row+1}"

            df_print = _df_display(df_obj)
            if df_print is None or df_print.empty:
                continue

            n_cols = len(df_print.columns)
            n_data_rows = len(df_print)
            
            for r_idx, row_data in enumerate(
                dataframe_to_rows(df_print, index=False, header=True),
                start=start_row,
            ):
                is_header = (r_idx == start_row)
                is_alt_row = (r_idx - start_row) % 2 == 0 and not is_header
                
                for c_idx, value in enumerate(row_data, 1):
                    ws_cell = ws.cell(row=r_idx, column=c_idx, value=value)
                    
                    if modern_style:
                        if is_header:
                            ws_cell.fill = PatternFill(start_color=HEADER_BG, end_color=HEADER_BG, fill_type="solid")
                            ws_cell.font = Font(bold=True, color=HEADER_FG, name="Calibri", size=10)
                            ws_cell.alignment = Alignment(horizontal="center", vertical="center")
                        else:
                            if is_alt_row:
                                ws_cell.fill = PatternFill(start_color=ALT_ROW_BG, end_color=ALT_ROW_BG, fill_type="solid")
                            ws_cell.font = Font(name="Calibri", size=10)
                            ws_cell.border = Border(bottom=Side(style="thin", color="DEE2E6"))
                    else:
                        if is_header:
                            ws_cell.font = Font(bold=True)
            
            # Add data bars to numeric columns
            if modern_style and data_bars and n_data_rows > 1:
                # Find numeric columns that might benefit from data bars
                for col_idx, col_name in enumerate(df_print.columns, 1):
                    col_lower = str(col_name).lower()
                    if any(kw in col_lower for kw in ["number", "count", "total", "citations", "documents"]):
                        col_letter = get_column_letter(col_idx)
                        data_bar_rule = DataBarRule(
                            start_type="min",
                            end_type="max",
                            color=ACCENT_COLOR,
                            showValue=True,
                            minLength=None,
                            maxLength=None
                        )
                        range_str = f"{col_letter}{start_row + 1}:{col_letter}{start_row + n_data_rows}"
                        try:
                            ws.conditional_formatting.add(range_str, data_bar_rule)
                        except Exception:
                            pass  # Skip if conditional formatting fails

        elif it == "plot" and img_path is not None:
            try:
                img = XLImage(str(img_path))
                _fit_image(
                    img,
                    scale=image_scale,
                    max_w=int(image_max_width_px),
                    max_h=int(image_max_height_px),
                )
                ws.add_image(img, "B5")
                link_cell = ws["A2"]
                link_cell.value = str(img_path)
                _set_external_file_link(link_cell, img_path)
            except Exception:
                pass

        elif it in {"text", "description"} and text_value:
            ws.cell(row=start_row, column=1, value=text_value)

        if autofit:
            _autosize_columns(ws)

        if not show_gridlines:
            ws.sheet_view.showGridLines = False

    if striped_rows:
        for row in toc_ws.iter_rows(min_row=2, max_row=row_idx - 1):
            if row[0].row % 2 == 0:
                for cell in row:
                    cell.fill = PatternFill(
                        start_color="F2F2F2",
                        end_color="F2F2F2",
                        fill_type="solid",
                    )

    _autosize_columns(toc_ws)

    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]

    out_path = _resolve_path(self, output_path, results_base)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return out_path


# -----------------------------
# POWERPOINT
# -----------------------------

def _ppt_add_df_table(
    slide,
    placeholder,
    df: pd.DataFrame,
    *,
    round_ndigits: int = 3,
    header_fill: Optional[PPTColor] = None,
    header_font: Optional[PPTColor] = None,
    style_headers: bool = True,
    modern_style: bool = True,
    max_rows: int = 10,
) -> None:
    """
    Render a pandas DataFrame as a PPT table inside `placeholder` bounds.

    Parameters
    ----------
    slide
        pptx Slide object.
    placeholder
        Content placeholder shape for the table.
    df : pandas.DataFrame
        Input DataFrame to display.
    round_ndigits : int, default 3
        Decimal places for numeric columns.
    header_fill : PPTColor, optional
        Fill color for header row.
    header_font : PPTColor, optional
        Font color for header row.
    style_headers : bool, default True
        Whether to style header cells with fill + bold font.
    modern_style : bool, default True
        Apply modern styling with zebra stripes and better fonts.
    max_rows : int, default 10
        Maximum rows to show (for PowerPoint readability).
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    
    dfp = _df_display(df, ndigits=round_ndigits)
    if dfp is None or dfp.empty:
        return
    
    # Limit rows for PowerPoint readability
    if len(dfp) > max_rows:
        dfp = dfp.head(max_rows)

    rows, cols = dfp.shape
    table_shape = slide.shapes.add_table(
        rows + 1,
        cols,
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    ).table
    
    # Modern colors
    HEADER_BG = PPTColor(44, 62, 80)      # #2C3E50
    HEADER_FG = PPTColor(255, 255, 255)   # White
    ALT_ROW = PPTColor(248, 249, 250)     # #F8F9FA
    DATA_COLOR = PPTColor(52, 73, 94)     # #34495E
    
    if header_fill is None:
        header_fill = HEADER_BG
    if header_font is None:
        header_font = HEADER_FG

    # Header row
    for j, name in enumerate(dfp.columns):
        cell = table_shape.cell(0, j)
        cell.text = str(name)
        
        # Center align and style header
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.name = "Calibri"
        para.font.size = Pt(10)
        
        if style_headers or modern_style:
            para.font.bold = True
            para.font.color.rgb = header_font
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
        
        # Vertical alignment
        cell.text_frame.word_wrap = True
        try:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    # Data rows
    for i in range(rows):
        is_alt = (i % 2 == 1)
        for j in range(cols):
            cell = table_shape.cell(i + 1, j)
            cell.text = str(dfp.iloc[i, j])
            
            # Style data cell
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Calibri"
            para.font.size = Pt(9)
            para.font.color.rgb = DATA_COLOR
            
            # Apply zebra striping
            if modern_style and is_alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW
            
            # Vertical alignment
            try:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass


def save_powerpoint_report_from_template(
    self: Any,
    output_path: str = "bibliometric_report.pptx",
    template_path: Union[str, Path] = FD
    / "additional files"
    / "template for report output.xlsx",
    template_sheet: str = "all",
    theme_color: str = "FFFFFF",
    title_font: Dict[str, Any] | None = None,
    subtitle_font: Dict[str, Any] | None = None,
    bullet_font: Dict[str, Any] | None = None,
    image_width: float = 6.5,
    align_center: bool = True,
    include_toc_slide: bool = True,
    footer_text: str = "",
    top_n: Optional[int] = 5,
    presentation_title: str = "Bibliometric Report",
    presentation_subtitle: str = "",
    bold_slide_titles: bool = True,
    custom_template_path: str = "",
    color_by_section: bool = False,
    style_table_headers: bool = True,
    show_section_footer: bool = False,
    insert_logo: bool = False,
    logo_path: str = str(ICONS_DIR / "logo.png"),
    logo_position: str = "bottom-right",
    logo_width: float = 1.0,
    theme_palette: Optional[Dict[str, PPTColor]] = None,
    results_base: Union[str, Path] = "results",
    # New modern options
    modern_style: bool = True,
) -> Path:
    """
    Generate a PowerPoint deck from a template definition sheet.

    Behavior w.r.t. Level 1 / Level 2
    ----------------------------------
    - There is one global title slide (presentation_title / subtitle).
    - Optionally a Table-of-Contents slide listing unique Level 1 values.
    - For each *new* Level 1 value, a "section title" slide is inserted
      with Level 1 as the slide title.
    - For every template row with an item:
        * The content slide title is **Level 2** (only).
        * If Level 2 is empty, Level 1 is used as fallback.

    Parameters
    ----------
    self : Any
        Object holding DataFrame attributes and optionally `res_folder`.
    output_path : str, default "bibliometric_report.pptx"
        Output PPTX file name or path (relative to `results_base` / `self.res_folder`).
    template_path : str or Path, default "<FD>/additional files/template for report output.xlsx"
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    theme_color : str, default "FFFFFF"
        Reserved for future theme tuning (not currently used).
    title_font : dict, optional
        Font settings for title slide main title.
    subtitle_font : dict, optional
        Font settings for slide subtitles (currently only used for global title slide).
    bullet_font : dict, optional
        Font settings for text content.
    image_width : float, default 6.5
        Width in inches for inserted images (plots).
    align_center : bool, default True
        Reserved for future alignment customization.
    include_toc_slide : bool, default True
        Insert a Table-of-Contents slide summarizing Level 1 sections.
    footer_text : str, default ""
        Optional footer text for each content slide if `show_section_footer` is True.
    top_n : int or None, default 5
        Default max rows for tables when Head column is empty. None → full table.
    presentation_title : str, default "Bibliometric Report"
        Title slide main title text.
    presentation_subtitle : str, default ""
        Title slide subtitle text.
    bold_slide_titles : bool, default True
        Render slide titles in bold.
    custom_template_path : str, default ""
        Optional path to a PPTX template to use.
    color_by_section : bool, default False
        Reserved for future section-based theming.
    style_table_headers : bool, default False
        Apply header fill/font styling for tables.
    show_section_footer : bool, default False
        Add footer text at bottom of each content slide.
    insert_logo : bool, default False
        Insert logo on all slides.
    logo_path : str, default "<ICONS_DIR>/logo.png"
        Path to logo image file.
    logo_position : str, default "bottom-right"
        One of {"bottom-right", "bottom-left", "top-left", "top-right"}.
    logo_width : float, default 1.0
        Logo width in inches.
    theme_palette : dict, optional
        Palette dictionary with keys "header_fill", "header_font" as PPTColor.
    results_base : str or Path, default "results"
        Base folder for relative paths and output.

    Returns
    -------
    Path
        Path to the saved PPTX file.
    """
    template_path = _resolve_path(self, template_path, results_base)
    df_template = _read_template(template_path, template_sheet)

    if title_font is None:
        title_font = {"name": "Calibri Light", "size": 44, "bold": True, "color": "2C3E50"}
    if subtitle_font is None:
        subtitle_font = {"name": "Calibri", "size": 24, "italic": False, "color": "7F8C8D"}
    if bullet_font is None:
        bullet_font = {"name": "Calibri", "size": 16}

    # Modern color palette
    HEADER_BG_PPT = PPTColor(44, 62, 80)      # #2C3E50
    HEADER_FG_PPT = PPTColor(255, 255, 255)   # White
    ALT_ROW_PPT = PPTColor(248, 249, 250)     # #F8F9FA
    ACCENT_PPT = PPTColor(52, 152, 219)       # #3498DB

    prs = Presentation(custom_template_path) if custom_template_path else Presentation()
    
    # Set 16:9 aspect ratio for modern look
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_layout.shapes.title
    title_shape.text = presentation_title
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = title_font.get("name", "Calibri Light")
    p.font.size = Pt(title_font.get("size", 44))
    p.font.bold = bool(title_font.get("bold", True))
    if "color" in title_font:
        color_hex = title_font["color"].lstrip("#")
        p.font.color.rgb = PPTColor(
            int(color_hex[0:2], 16), 
            int(color_hex[2:4], 16), 
            int(color_hex[4:6], 16)
        )

    if presentation_subtitle and len(title_layout.placeholders) > 1:
        subtitle_shape = title_layout.placeholders[1]
        subtitle_shape.text = presentation_subtitle
        sp = subtitle_shape.text_frame.paragraphs[0]
        sp.font.name = subtitle_font.get("name", "Arial")
        sp.font.size = Pt(subtitle_font.get("size", 24))
        sp.font.italic = bool(subtitle_font.get("italic", False))

    title_and_content = prs.slide_layouts[1]
    section_title_layout = prs.slide_layouts[0]  # reuse title slide layout for sections

    # Optional ToC slide
    if include_toc_slide:
        toc_slide = prs.slides.add_slide(title_and_content)
        toc_slide.shapes.title.text = "Table of Contents"
        if bold_slide_titles:
            toc_slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        tf = toc_slide.placeholders[1].text_frame
        seen: List[str] = []
        if "Level 1" in df_template.columns:
            for lvl1 in df_template["Level 1"]:
                s = _clean_str(lvl1)
                if s and s not in seen:
                    seen.append(s)
            for s in seen:
                pr = tf.add_paragraph()
                pr.text = s
                pr.level = 0

    if theme_palette is None:
        theme_palette = {
            "header_fill": HEADER_BG_PPT,
            "header_font": HEADER_FG_PPT,
        }

    current_lvl1: Optional[str] = None

    for _, row in df_template.iterrows():
        # Skip rows marked as xlsx_only
        if _is_xlsx_only(row):
            continue
            
        level1 = _clean_str(row.get("Level 1", ""))
        level2 = _clean_str(row.get("Level 2", ""))
        item_type = _infer_item_type(row)

        # Insert a section title slide when Level 1 changes
        if level1 and level1 != current_lvl1:
            section_slide = prs.slides.add_slide(section_title_layout)
            section_title_shape = section_slide.shapes.title
            section_title_shape.text = level1
            if bold_slide_titles:
                section_title_shape.text_frame.paragraphs[0].font.bold = True
            current_lvl1 = level1

        if not item_type:
            # Still honor explicit page-break flag (no direct PPTX equivalent, so ignore).
            continue

        it = item_type.lower()

        df_obj: Optional[pd.DataFrame] = None
        img_path: Optional[Path] = None
        text_value: str = ""
        has_content = False

        if it == "table":
            data_attr = _clean_str(row.get("Data Attr", ""))
            df_obj = _get_df_attr(self, data_attr)
            if isinstance(df_obj, pd.DataFrame) and not df_obj.empty:
                # Filter columns for pptx output
                columns_spec = _clean_str(row.get("Columns (docx/pptx/tex)", ""))
                df_obj = _filter_columns(df_obj, columns_spec)
            has_content = isinstance(df_obj, pd.DataFrame) and not df_obj.empty

        elif it == "plot":
            tpath = _clean_str(row.get("Path", ""))
            filename = _clean_str(row.get("Plot Filename", ""))
            if filename:
                img_path = _resolve_image_path(self, tpath, filename, results_base)
                has_content = img_path is not None

        elif it == "text":
            text_value = _clean_str(row.get("Text", ""))
            if not text_value:
                narrative_attr = _clean_str(row.get("Narrative", ""))
                if narrative_attr and hasattr(self, narrative_attr):
                    narr = getattr(self, narrative_attr)
                    if isinstance(narr, str):
                        text_value = narr
            has_content = bool(text_value)

        elif it == "description":
            data_attr = _clean_str(row.get("Data Attr", ""))
            if data_attr and hasattr(self, data_attr):
                desc = getattr(self, data_attr)
                if isinstance(desc, str):
                    text_value = desc
                    has_content = bool(text_value.strip())

        if not has_content:
            continue

        # Content slide: title is always Level 2 (fallback = Level 1)
        slide = prs.slides.add_slide(title_and_content)
        title_text = level2 or level1 or ""
        slide.shapes.title.text = title_text
        if bold_slide_titles:
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        placeholder = slide.placeholders[1]

        if it == "plot" and img_path is not None:
            slide.shapes.add_picture(
                str(img_path),
                placeholder.left,
                placeholder.top,
                width=placeholder.width,
                height=placeholder.height,
            )

        elif it == "table" and df_obj is not None:
            df_print_raw = _apply_head(df_obj, row.get("Head", None), top_n)
            # Filter columns for PPTX
            columns_spec = _clean_str(row.get("Columns (docx/pptx/tex)", ""))
            df_print_filtered = _filter_columns(df_print_raw, columns_spec)
            if df_print_filtered is not None and not df_print_filtered.empty:
                _ppt_add_df_table(
                    slide,
                    placeholder,
                    df_print_filtered,
                    round_ndigits=3,
                    header_fill=theme_palette["header_fill"],
                    header_font=theme_palette["header_font"],
                    style_headers=bool(style_table_headers),
                    modern_style=modern_style,
                )

        elif it in {"text", "description"} and text_value:
            tf = placeholder.text_frame
            # Clear any existing content
            for para in tf.paragraphs:
                para.text = ""
            pbody = tf.paragraphs[0]
            pbody.text = text_value
            pbody.font.name = bullet_font.get("name", "Calibri")
            pbody.font.size = Pt(bullet_font.get("size", 18))

        if show_section_footer and level1:
            footer = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(6.8),
                Inches(8),
                Inches(0.3),
            )
            footer.text_frame.text = footer_text or level1

    if insert_logo:
        logo = _resolve_path(self, logo_path, results_base)
        if logo.exists():
            for slide in prs.slides:
                if logo_position == "bottom-right":
                    left, top = Inches(9), Inches(6.5)
                elif logo_position == "bottom-left":
                    left, top = Inches(0.5), Inches(6.5)
                elif logo_position == "top-left":
                    left, top = Inches(0.5), Inches(0.5)
                elif logo_position == "top-right":
                    left, top = Inches(9), Inches(0.5)
                else:
                    left, top = Inches(9), Inches(6.5)
                slide.shapes.add_picture(
                    str(logo),
                    left,
                    top,
                    width=Inches(float(logo_width)),
                )

    out_path = _resolve_path(self, output_path, results_base)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# -----------------------------
# WORD
# -----------------------------

# Modern color palette
MODERN_COLORS = {
    "primary": DocxColor(44, 62, 80),      # Dark blue-gray #2C3E50
    "secondary": DocxColor(52, 152, 219),  # Blue #3498DB
    "accent": DocxColor(46, 204, 113),     # Green #2ECC71
    "warning": DocxColor(241, 196, 15),    # Yellow #F1C40F
    "danger": DocxColor(231, 76, 60),      # Red #E74C3C
    "light": DocxColor(236, 240, 241),     # Light gray #ECF0F1
    "dark": DocxColor(52, 73, 94),         # Dark #34495E
    "white": DocxColor(255, 255, 255),
    "muted": DocxColor(149, 165, 166),     # Muted gray #95A5A6
}


def _apply_modern_table_style(table, header_color=None, zebra=True, autofit=True):
    """
    Apply modern styling to a Word table.
    
    Parameters
    ----------
    table : docx.table.Table
        The table to style.
    header_color : RGBColor, optional
        Background color for header row.
    zebra : bool
        Whether to apply zebra striping.
    autofit : bool
        Whether to autofit table to content.
    """
    from docx.oxml.ns import nsdecls, qn
    from docx.oxml import parse_xml, OxmlElement
    from docx.shared import Pt, Twips
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    
    if header_color is None:
        header_color = MODERN_COLORS["primary"]
    
    # Convert color to hex string (RGBColor's str() returns hex without #)
    header_hex = str(header_color)
    
    # Set table alignment to center
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    # Enable autofit
    if autofit:
        try:
            # Set table to autofit content
            tbl = table._tbl
            tblPr = tbl.tblPr if tbl.tblPr is not None else OxmlElement('w:tblPr')
            
            # Set table width to auto
            tblW = OxmlElement('w:tblW')
            tblW.set(qn('w:type'), 'auto')
            tblW.set(qn('w:w'), '0')
            tblPr.append(tblW)
            
            # Set autofit behavior
            tblLayout = OxmlElement('w:tblLayout')
            tblLayout.set(qn('w:type'), 'autofit')
            tblPr.append(tblLayout)
        except Exception:
            pass
    
    # Style header row
    header_row = table.rows[0]
    for cell in header_row.cells:
        # Set background color
        shading_elm = parse_xml(
            f'<w:shd {nsdecls("w")} w:fill="{header_hex}" w:val="clear"/>'
        )
        cell._tc.get_or_add_tcPr().append(shading_elm)
        
        # Set cell padding
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            tcMar = OxmlElement('w:tcMar')
            for margin_name in ['top', 'left', 'bottom', 'right']:
                margin = OxmlElement(f'w:{margin_name}')
                margin.set(qn('w:w'), '60')
                margin.set(qn('w:type'), 'dxa')
                tcMar.append(margin)
            tcPr.append(tcMar)
        except Exception:
            pass
        
        # Style text
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.name = "Calibri"
                run.font.color.rgb = MODERN_COLORS["white"]
                run.font.size = Pt(10)
    
    # Style data rows with zebra striping
    light_gray = "F8F9FA"
    white = "FFFFFF"
    
    for i, row in enumerate(table.rows[1:], start=1):
        bg_color = light_gray if zebra and i % 2 == 0 else white
        for cell in row.cells:
            shading_elm = parse_xml(
                f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>'
            )
            cell._tc.get_or_add_tcPr().append(shading_elm)
            
            # Style text
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(9)
                    run.font.color.rgb = MODERN_COLORS["dark"]


def _add_cover_page(doc, title="Bibliometric Report", subtitle=None, date_str=None, 
                    summary_stats=None, color=None):
    """
    Add a modern cover page to the document.
    
    Parameters
    ----------
    doc : Document
        The Word document.
    title : str
        Main title.
    subtitle : str, optional
        Subtitle or dataset description.
    date_str : str, optional
        Date string (defaults to today).
    summary_stats : dict, optional
        Key statistics to display (e.g., {"Documents": 500, "Sources": 50}).
    color : RGBColor, optional
        Accent color.
    """
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from datetime import datetime
    
    if color is None:
        color = MODERN_COLORS["primary"]
    
    if date_str is None:
        date_str = datetime.now().strftime("%B %d, %Y")
    
    # Add spacing at top
    for _ in range(3):
        doc.add_paragraph()
    
    # Main title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(title)
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = color
    
    # Subtitle
    if subtitle:
        sub_para = doc.add_paragraph()
        sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_para.add_run(subtitle)
        sub_run.font.size = Pt(14)
        sub_run.font.color.rgb = MODERN_COLORS["muted"]
    
    # Decorative line
    doc.add_paragraph()
    line_para = doc.add_paragraph()
    line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    line_run = line_para.add_run("─" * 40)
    line_run.font.color.rgb = MODERN_COLORS["light"]
    
    # Date
    doc.add_paragraph()
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_run = date_para.add_run(date_str)
    date_run.font.size = Pt(12)
    date_run.font.color.rgb = MODERN_COLORS["muted"]
    
    # Summary statistics box
    if summary_stats:
        doc.add_paragraph()
        doc.add_paragraph()
        
        stats_para = doc.add_paragraph()
        stats_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        stats_text = "  │  ".join([f"{k}: {v:,}" if isinstance(v, int) else f"{k}: {v}" 
                                   for k, v in summary_stats.items()])
        stats_run = stats_para.add_run(stats_text)
        stats_run.font.size = Pt(11)
        stats_run.font.color.rgb = MODERN_COLORS["secondary"]
    
    # Page break after cover
    doc.add_page_break()


def _add_table_of_contents(doc, color=None):
    """
    Add a Table of Contents placeholder to the document.
    
    Note: The TOC needs to be updated in Word after opening.
    
    Parameters
    ----------
    doc : Document
        The Word document.
    color : RGBColor, optional
        Heading color.
    """
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    
    if color is None:
        color = MODERN_COLORS["primary"]
    
    # TOC heading
    toc_heading = doc.add_heading("Table of Contents", level=1)
    toc_heading.runs[0].font.color.rgb = color
    
    # Add TOC field
    paragraph = doc.add_paragraph()
    run = paragraph.add_run()
    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')
    
    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
    
    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'separate')
    
    fldChar3 = OxmlElement('w:fldChar')
    fldChar3.set(qn('w:fldCharType'), 'end')
    
    run._r.append(fldChar1)
    run._r.append(instrText)
    run._r.append(fldChar2)
    run._r.append(fldChar3)
    
    # Instruction text
    note_para = doc.add_paragraph()
    note_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    note_run = note_para.add_run("(Right-click and select 'Update Field' to populate)")
    note_run.font.size = Pt(9)
    note_run.font.italic = True
    note_run.font.color.rgb = MODERN_COLORS["muted"]
    
    doc.add_page_break()


def _get_summary_stats(self):
    """
    Extract summary statistics from the BiblioStats object for the cover page.
    
    Parameters
    ----------
    self : Any
        BiblioStats or similar object.
        
    Returns
    -------
    dict
        Dictionary of key statistics.
    """
    stats = {}
    
    # Number of documents
    if hasattr(self, 'n'):
        stats["Documents"] = self.n
    elif hasattr(self, 'df'):
        stats["Documents"] = len(self.df)
    
    # Date range - handle mixed types safely
    if hasattr(self, 'df') and 'Year' in self.df.columns:
        try:
            years = pd.to_numeric(self.df['Year'], errors='coerce').dropna()
            if len(years) > 0:
                min_year = int(years.min())
                max_year = int(years.max())
                stats["Period"] = f"{min_year}-{max_year}"
        except Exception:
            pass
    
    # Number of sources
    if hasattr(self, 'sources_counts_df'):
        stats["Sources"] = len(self.sources_counts_df)
    
    # Number of authors
    if hasattr(self, 'authors_counts_df'):
        stats["Authors"] = len(self.authors_counts_df)
    
    # Total citations - handle mixed types safely
    if hasattr(self, 'df') and 'Cited by' in self.df.columns:
        try:
            citations = pd.to_numeric(self.df['Cited by'], errors='coerce').sum()
            if pd.notna(citations):
                stats["Citations"] = int(citations)
        except Exception:
            pass
    
    return stats


def _save_word_report_simple(
    self: Any,
    output_path: str = "bibliometric_report.docx",
    results_base: Union[str, Path] = "results",
    top_n: int = 10,
    max_plots: Optional[int] = None,
    plot_filter: Optional[List[str]] = None,
    report_title: str = "Bibliometric Analysis Report",
) -> Path:
    """
    Simple Word report generator that directly adds found tables and plots.
    Bypasses the template system for maximum reliability.
    
    Parameters
    ----------
    self : Any
        BiblioAnalysis object
    output_path : str
        Output file path
    results_base : str or Path
        Base folder for output
    top_n : int
        Max rows per table
    max_plots : int or None
        Maximum plots to include (None = all)
    plot_filter : list or None
        Keywords to filter plots (e.g., ["bradford", "lotka"])
    report_title : str
        Title for the report
        
    Returns
    -------
    Path
        Path to saved document
    """
    from datetime import datetime
    
    doc = Document()
    
    # Title
    title_para = doc.add_heading(report_title, 0)
    title_para.alignment = 1  # Center
    
    # Subtitle with date
    subtitle = doc.add_paragraph()
    subtitle.alignment = 1
    run = subtitle.add_run(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    run.font.size = DocxPt(11)
    run.font.color.rgb = DocxColor(128, 128, 128)
    
    # Summary stats
    if hasattr(self, 'df'):
        n_docs = len(self.df)
        summary = doc.add_paragraph()
        summary.alignment = 1
        run = summary.add_run(f"Documents: {n_docs:,}")
        run.font.size = DocxPt(12)
    
    doc.add_paragraph()
    
    # Track what we add
    tables_added = 0
    plots_added = 0
    
    # Tables to look for
    TABLE_MAP = {
        "descriptives_df": "Dataset Descriptives",
        "performances_df": "Global Performance Metrics", 
        "production_df": "Scientific Production",
        "time_series_stats_df": "Time Series Statistics",
        "missings_df": "Missing Values",
        "sources_counts_df": "Top Sources",
        "authors_counts_df": "Top Authors",
        "author_keywords_counts_df": "Author Keywords",
        "index_keywords_counts_df": "Index Keywords",
        "document_types_counts_df": "Document Types",
        "ca_country_counts_df": "Corresponding Author Countries",
        "all_countries_counts_df": "All Countries",
        "all_countries_stats_df": "Countries Statistics",
        "affiliations_counts_df": "Affiliations",
        "references_counts_df": "References",
        "bradford_df": "Bradford's Law",
        "lotka_df": "Lotka's Law",
        "zipf_df": "Zipf's Law",
        "top_cited_docs_global_df": "Top Cited Documents (Global)",
        "top_cited_docs_local_df": "Top Cited Documents (Local)",
    }
    
    # Add tables
    doc.add_heading("Tables", level=1)
    
    for attr, title_text in TABLE_MAP.items():
        if hasattr(self, attr):
            df = getattr(self, attr)
            if isinstance(df, pd.DataFrame) and not df.empty:
                doc.add_heading(title_text, level=2)
                
                # Limit rows and columns
                df_display = df.head(top_n)
                if len(df_display.columns) > 7:
                    df_display = df_display.iloc[:, :7]
                
                # Create table
                rows, cols = df_display.shape
                table = doc.add_table(rows=rows + 1, cols=cols)
                table.style = "Table Grid"
                
                # Header with formatting
                for j, col_name in enumerate(df_display.columns):
                    cell = table.cell(0, j)
                    cell.text = str(col_name)
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                
                # Data
                for i in range(rows):
                    for j in range(cols):
                        val = df_display.iloc[i, j]
                        if isinstance(val, float):
                            cell_text = f"{val:.3f}" if abs(val) < 1000 else f"{val:,.0f}"
                        elif isinstance(val, int):
                            cell_text = f"{val:,}"
                        else:
                            cell_text = str(val)[:50]
                        table.cell(i + 1, j).text = cell_text
                
                doc.add_paragraph()
                tables_added += 1
                print(f"  ✓ Added table: {title_text} ({len(df)} rows)")
    
    # Add plots
    doc.add_page_break()
    doc.add_heading("Plots", level=1)
    
    # Find plots folder
    res_folder = getattr(self, "res_folder", None)
    plots_folder = Path(res_folder) / "plots" if res_folder else Path(results_base) / "plots"
    
    if plots_folder.exists():
        png_files = sorted(plots_folder.glob("*.png"))
        
        # Apply filter
        if plot_filter:
            plot_filter_lower = [f.lower() for f in plot_filter]
            png_files = [
                f for f in png_files 
                if any(keyword in f.name.lower() for keyword in plot_filter_lower)
            ]
            print(f"  Filtered to {len(png_files)} plots matching: {plot_filter}")
        
        # Apply max limit
        if max_plots and len(png_files) > max_plots:
            png_files = png_files[:max_plots]
            print(f"  Limited to first {max_plots} plots")
        
        for png_file in png_files:
            try:
                title_text = png_file.stem.replace("_", " ").replace("-", " ").title()
                doc.add_heading(title_text, level=2)
                doc.add_picture(str(png_file), width=DocxInches(5.5))
                doc.add_paragraph()
                plots_added += 1
                print(f"  ✓ Added plot: {png_file.name}")
            except Exception as e:
                print(f"  ✗ Failed to add {png_file.name}: {e}")
    else:
        print(f"  ⚠ Plots folder not found: {plots_folder}")
    
    # Save
    out_path = _resolve_path(self, output_path, results_base)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    
    print(f"\n{'='*50}")
    print(f"Report saved to: {out_path}")
    print(f"Tables added: {tables_added}")
    print(f"Plots added: {plots_added}")
    print(f"{'='*50}")
    
    return out_path


def save_word_report_from_template(
    self: Any,
    output_path: str = "bibliometric_report.docx",
    template_path: Union[str, Path] = FD
    / "additional files"
    / "template for report output.xlsx",
    template_sheet: str = "all",
    top_n: Optional[int] = 5,
    default_table_style: str = "Table Grid",
    enumerate_figures: bool = True,
    enumerate_tables: bool = True,
    enumerate_sections: bool = True,
    include_logo: bool = True,
    heading_color: str = "primary",
    results_base: Union[str, Path] = "results",
    # New modern styling options
    modern_style: bool = True,
    include_cover_page: bool = True,
    include_toc: bool = True,
    report_title: str = "Bibliometric Analysis Report",
    report_subtitle: str = None,
    zebra_tables: bool = True,
    # Simple mode - bypass template system
    use_simple_mode: bool = False,
    simple_max_plots: Optional[int] = None,
    simple_plot_filter: Optional[List[str]] = None,
) -> Path:
    """
    Generate a Word document from a template definition sheet.

    Parameters
    ----------
    self : Any
        Object holding DataFrame attributes and optionally `res_folder`.
    output_path : str, default "bibliometric_report.docx"
        Output DOCX file name or path.
    template_path : str or Path
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    top_n : int or None, default 5
        Default max rows for tables when Head column is empty.
    default_table_style : str, default "Table Grid"
        Fallback Word table style.
    enumerate_figures : bool, default True
        Prefix figure captions with "Figure N:".
    enumerate_tables : bool, default True
        Prefix table captions with "Table N:".
    enumerate_sections : bool, default True
        Prefix headings with "1.", "1.1", etc.
    include_logo : bool, default True
        Place logo image in the footer if available.
    heading_color : str, default "primary"
        Color name for headings (primary, secondary, accent, or legacy colors).
    results_base : str or Path, default "results"
        Base folder for relative paths and output.
    modern_style : bool, default True
        Apply modern table styling with colored headers and zebra stripes.
    include_cover_page : bool, default True
        Add a cover page with title, date, and summary statistics.
    include_toc : bool, default True
        Add a Table of Contents.
    report_title : str, default "Bibliometric Analysis Report"
        Title for the cover page.
    report_subtitle : str, optional
        Subtitle for the cover page.
    zebra_tables : bool, default True
        Apply zebra striping to tables.

    Returns
    -------
    Path
        Path to the saved DOCX file.
    """
    # Simple mode - bypass template entirely
    if use_simple_mode:
        return _save_word_report_simple(
            self,
            output_path=output_path,
            results_base=results_base,
            top_n=top_n or 10,
            max_plots=simple_max_plots,
            plot_filter=simple_plot_filter,
            report_title=report_title,
        )
    
    template_path = _resolve_path(self, template_path, results_base)
    df_template = _read_template(template_path, template_sheet)

    doc = Document()
    
    # Extended color map with modern colors
    color_map_docx = {
        # Modern palette
        "primary": MODERN_COLORS["primary"],
        "secondary": MODERN_COLORS["secondary"],
        "accent": MODERN_COLORS["accent"],
        "dark": MODERN_COLORS["dark"],
        "muted": MODERN_COLORS["muted"],
        # Legacy colors
        "orange": DocxColor(255, 165, 0),
        "blue": DocxColor(0, 112, 192),
        "black": DocxColor(0, 0, 0),
        "gray": DocxColor(128, 128, 128),
        "green": DocxColor(0, 176, 80),
        "red": DocxColor(255, 0, 0),
        "purple": DocxColor(112, 48, 160),
    }
    
    heading_rgb = color_map_docx.get(heading_color.lower(), MODERN_COLORS["primary"])
    
    # Add cover page
    if include_cover_page:
        summary_stats = _get_summary_stats(self)
        _add_cover_page(
            doc, 
            title=report_title, 
            subtitle=report_subtitle,
            summary_stats=summary_stats,
            color=heading_rgb
        )
    
    # Add Table of Contents
    if include_toc:
        _add_table_of_contents(doc, color=heading_rgb)
    
    # If no cover page, add simple title
    if not include_cover_page:
        title_para = doc.add_heading(report_title, level=0)
        title_para.runs[0].font.color.rgb = heading_rgb
        doc.add_paragraph()

    figure_counter = 1
    table_counter = 1

    current_lvl1: Optional[str] = None
    sec_idx = 0
    sub_idx = 0

    for _, row in df_template.iterrows():
        # Skip rows marked as xlsx_only
        if _is_xlsx_only(row):
            continue
            
        lvl1 = _clean_str(row.get("Level 1", ""))
        lvl2 = _clean_str(row.get("Level 2", ""))

        # New Level 1 section
        if lvl1 and lvl1 != current_lvl1:
            sec_idx += 1
            sub_idx = 0
            heading1 = f"{sec_idx}. {lvl1}" if enumerate_sections else lvl1
            para1 = doc.add_heading(heading1, level=1)
            para1.runs[0].font.color.rgb = heading_rgb
            current_lvl1 = lvl1

        # Level 2 heading for every row that has it
        if lvl2:
            sub_idx += 1
            if enumerate_sections and sec_idx > 0:
                heading2 = f"{sec_idx}.{sub_idx} {lvl2}"
            else:
                heading2 = lvl2
            para2 = doc.add_heading(heading2, level=2)
            para2.runs[0].font.color.rgb = heading_rgb

        # Section-level description / narrative (attributes)
        desc_attr = _clean_str(row.get("Description", ""))
        if desc_attr and hasattr(self, desc_attr):
            desc_text = getattr(self, desc_attr)
            if isinstance(desc_text, str) and desc_text:
                doc.add_paragraph(desc_text)

        narrative_attr = _clean_str(row.get("Narrative", ""))
        if narrative_attr and hasattr(self, narrative_attr):
            narrative = getattr(self, narrative_attr)
            if isinstance(narrative, str) and narrative:
                for line in narrative.splitlines():
                    doc.add_paragraph(line)

        item_type = _infer_item_type(row)
        if not item_type:
            if _flag_true(row.get("Page Break After", "")):
                doc.add_page_break()
            continue

        it = item_type.lower()
        
        # DEBUG: Show what item type we're processing
        lvl2_debug = _clean_str(row.get("Level 2", ""))

        if it == "table":
            data_attr = _clean_str(row.get("Data Attr", ""))
            df_obj = _get_df_attr(self, data_attr)
            if isinstance(df_obj, pd.DataFrame) and not df_obj.empty:
                dfp_raw = _apply_head(df_obj, row.get("Head", None), top_n)
                # Filter columns for docx/pptx/tex output
                columns_spec = _clean_str(row.get("Columns (docx/pptx/tex)", ""))
                dfp_filtered = _filter_columns(dfp_raw, columns_spec)
                dfp = _df_display(dfp_filtered)
                if dfp is not None and not dfp.empty:
                    rows, cols = dfp.shape
                    table = doc.add_table(rows=rows + 1, cols=cols)

                    # Apply base style first
                    style_name = _clean_str(row.get("Table Style", ""))
                    try:
                        table.style = style_name or default_table_style
                    except KeyError:
                        try:
                            table.style = default_table_style
                        except KeyError:
                            try:
                                table.style = "Table Grid"
                            except KeyError:
                                pass

                    # Header row
                    for j, name in enumerate(dfp.columns):
                        table.cell(0, j).text = str(name)

                    # Data rows
                    for i in range(rows):
                        for j in range(cols):
                            table.cell(i + 1, j).text = str(dfp.iloc[i, j])

                    # Apply modern styling if enabled
                    if modern_style:
                        _apply_modern_table_style(
                            table, 
                            header_color=heading_rgb,
                            zebra=zebra_tables
                        )

                    caption = _clean_str(row.get("Caption", ""))
                    if caption:
                        prefix = (
                            f"Table {table_counter}: "
                            if enumerate_tables
                            else ""
                        )
                        table_counter += 1
                        cap = doc.add_paragraph(prefix + caption, style="Caption")
                        cap.runs[0].font.color.rgb = heading_rgb

        elif it == "plot":
            tpath = _clean_str(row.get("Path", ""))
            filename = _clean_str(row.get("Plot Filename", ""))
            if filename:
                img = _resolve_image_path(self, tpath, filename, results_base)
                if img is not None:
                    doc.add_picture(str(img), width=DocxInches(5.5))
                    caption = _clean_str(row.get("Caption", ""))
                    if caption:
                        prefix = (
                            f"Figure {figure_counter}: "
                            if enumerate_figures
                            else ""
                        )
                        figure_counter += 1
                        cap = doc.add_paragraph(
                            prefix + caption,
                            style="Caption",
                        )
                        cap.runs[0].font.color.rgb = heading_rgb

        elif it == "text":
            text_value = _clean_str(row.get("Text", ""))
            if text_value:
                for line in text_value.splitlines():
                    doc.add_paragraph(line)

        elif it == "description":
            data_attr = _clean_str(row.get("Data Attr", ""))
            if data_attr and hasattr(self, data_attr):
                desc_text2 = getattr(self, data_attr)
                if isinstance(desc_text2, str) and desc_text2:
                    for line in desc_text2.splitlines():
                        doc.add_paragraph(line)

        if _flag_true(row.get("Page Break After", "")):
            doc.add_page_break()

    # Footer logo
    if include_logo:
        logo = ICONS_DIR / "logo.png"
        if logo.exists():
            for section in doc.sections:
                footer = section.footer
                para = (
                    footer.add_paragraph()
                    if not footer.paragraphs
                    else footer.paragraphs[0]
                )
                run = para.add_run()
                run.add_picture(str(logo), width=DocxInches(1.0))

    # DEBUG: Print summary

    out_path = _resolve_path(self, output_path, results_base)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    return out_path


# -----------------------------
# LaTeX
# -----------------------------

def save_tex_report_from_template(
    self: Any,
    output_path: str = "bibliometric_report.tex",
    template_path: Union[str, Path] = FD
    / "additional files"
    / "template for report output.xlsx",
    template_sheet: str = "all",
    top_n: int = 50,
    enumerate_figures: bool = True,
    enumerate_tables: bool = True,
    enumerate_sections: bool = True,
    results_base: Union[str, Path] = "results",
    # New modern options
    modern_style: bool = True,
    report_title: str = "Bibliometric Analysis Report",
    report_author: str = "",
    include_toc: bool = True,
) -> Path:
    """
    Generate a LaTeX (.tex) report from a template definition sheet.

    Parameters
    ----------
    self : Any
        Object holding DataFrame attributes and optionally `res_folder`.
    output_path : str, default "bibliometric_report.tex"
        Output TeX file name or path (relative to `results_base` / `self.res_folder`).
    template_path : str or Path, default "<FD>/additional files/template for report output.xlsx"
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    top_n : int, default 50
        Default max rows for tables when Head column is empty.
    enumerate_figures : bool, default True
        Prefix figure captions with "Figure N:".
    enumerate_tables : bool, default True
        Prefix table captions with "Table N:".
    enumerate_sections : bool, default True
        Prefix headings with "1.", "1.1", etc.
    results_base : str or Path, default "results"
        Base folder for relative paths and output.
    modern_style : bool, default True
        Use modern LaTeX styling with colors and improved typography.
    report_title : str, default "Bibliometric Analysis Report"
        Report title for title page.
    report_author : str, default ""
        Author name for title page.
    include_toc : bool, default True
        Include table of contents.

    Returns
    -------
    Path
        Path to the saved TeX file.
    """
    from datetime import datetime
    
    template_path = _resolve_path(self, template_path, results_base)
    df_template = _read_template(template_path, template_sheet)

    figure_counter = 1
    table_counter = 1
    
    # Escape title and author for LaTeX
    def _tex_escape(text):
        chars = {"&": r"\&", "%": r"\%", "$": r"\$", "#": r"\#", "_": r"\_", 
                 "{": r"\{", "}": r"\}", "~": r"\textasciitilde{}", "^": r"\textasciicircum{}"}
        for char, replacement in chars.items():
            text = text.replace(char, replacement)
        return text

    lines: List[str] = []
    
    if modern_style:
        # Modern LaTeX preamble with colors and styling
        lines.append(r"\documentclass[11pt,a4paper]{article}")
        lines.append(r"")
        lines.append(r"% Encoding and fonts")
        lines.append(r"\usepackage[utf8]{inputenc}")
        lines.append(r"\usepackage[T1]{fontenc}")
        lines.append(r"\usepackage{lmodern}")
        lines.append(r"")
        lines.append(r"% Layout")
        lines.append(r"\usepackage[margin=1in]{geometry}")
        lines.append(r"\usepackage{parskip}")
        lines.append(r"")
        lines.append(r"% Graphics and tables")
        lines.append(r"\usepackage{graphicx}")
        lines.append(r"\usepackage{booktabs}")
        lines.append(r"\usepackage{longtable}")
        lines.append(r"\usepackage{array}")
        lines.append(r"\usepackage{colortbl}")
        lines.append(r"")
        lines.append(r"% Colors")
        lines.append(r"\usepackage{xcolor}")
        lines.append(r"\definecolor{primary}{HTML}{2C3E50}")
        lines.append(r"\definecolor{secondary}{HTML}{3498DB}")
        lines.append(r"\definecolor{accent}{HTML}{E74C3C}")
        lines.append(r"\definecolor{tableheader}{HTML}{2C3E50}")
        lines.append(r"\definecolor{tablealt}{HTML}{F8F9FA}")
        lines.append(r"\definecolor{tableborder}{HTML}{DEE2E6}")
        lines.append(r"")
        lines.append(r"% Section styling")
        lines.append(r"\usepackage{titlesec}")
        lines.append(r"\titleformat{\section}{\normalfont\Large\bfseries\color{primary}}{\thesection}{1em}{}")
        lines.append(r"\titleformat{\subsection}{\normalfont\large\bfseries\color{secondary}}{\thesubsection}{1em}{}")
        lines.append(r"")
        lines.append(r"% Hyperlinks")
        lines.append(r"\usepackage{hyperref}")
        lines.append(r"\hypersetup{colorlinks=true,linkcolor=secondary,urlcolor=secondary,citecolor=secondary}")
        lines.append(r"")
        lines.append(r"% Headers and footers")
        lines.append(r"\usepackage{fancyhdr}")
        lines.append(r"\pagestyle{fancy}")
        lines.append(r"\fancyhf{}")
        lines.append(r"\fancyhead[L]{\small\textcolor{primary}{" + _tex_escape(report_title) + r"}}")
        lines.append(r"\fancyhead[R]{\small\textcolor{primary}{\thepage}}")
        lines.append(r"\renewcommand{\headrulewidth}{0.4pt}")
        lines.append(r"\renewcommand{\headrule}{\hbox to\headwidth{\color{secondary}\leaders\hrule height \headrulewidth\hfill}}")
        lines.append(r"")
        lines.append(r"% Table styling commands")
        lines.append(r"\newcommand{\tableheaderrow}{\rowcolor{tableheader}}")
        lines.append(r"\newcommand{\tablealtrow}{\rowcolor{tablealt}}")
        lines.append(r"")
        lines.append(r"% Document info")
        lines.append(r"\title{\textcolor{primary}{\textbf{" + _tex_escape(report_title) + r"}}}")
        lines.append(r"\author{" + _tex_escape(report_author) + r"}")
        lines.append(r"\date{" + datetime.now().strftime("%B %d, %Y") + r"}")
        lines.append(r"")
        lines.append(r"\begin{document}")
        lines.append(r"\maketitle")
        if include_toc:
            lines.append(r"\tableofcontents")
            lines.append(r"\newpage")
    else:
        # Simple preamble
        lines.append(r"\documentclass{article}")
        lines.append(r"\usepackage[utf8]{inputenc}")
        lines.append(r"\usepackage{graphicx}")
        lines.append(r"\usepackage{booktabs}")
        lines.append(r"\usepackage{geometry}")
        lines.append(r"\geometry{margin=1in}")
        lines.append(r"\title{" + _tex_escape(report_title) + r"}")
        lines.append(r"\date{}")
        lines.append(r"\begin{document}")
        lines.append(r"\maketitle")

    current_lvl1: Optional[str] = None
    sec_idx = 0
    sub_idx = 0

    for _, row in df_template.iterrows():
        # Skip rows marked as xlsx_only
        if _is_xlsx_only(row):
            continue
            
        lvl1 = _clean_str(row.get("Level 1", ""))
        lvl2 = _clean_str(row.get("Level 2", ""))

        if lvl1 and lvl1 != current_lvl1:
            sec_idx += 1
            sub_idx = 0
            heading1 = f"{sec_idx}. {lvl1}" if enumerate_sections else lvl1
            lines.append(f"\\section{{{heading1}}}")
            current_lvl1 = lvl1

        if lvl2:
            sub_idx += 1
            heading2 = (
                f"{sec_idx}.{sub_idx} {lvl2}"
                if enumerate_sections and sec_idx > 0
                else lvl2
            )
            lines.append(f"\\subsection{{{heading2}}}")

        desc_attr = _clean_str(row.get("Description", ""))
        if desc_attr and hasattr(self, desc_attr):
            desc_text = getattr(self, desc_attr)
            if isinstance(desc_text, str) and desc_text:
                lines.append(desc_text + r"\\")

        narrative_attr = _clean_str(row.get("Narrative", ""))
        if narrative_attr and hasattr(self, narrative_attr):
            narrative = getattr(self, narrative_attr)
            if isinstance(narrative, str) and narrative:
                for ln in narrative.splitlines():
                    lines.append(ln + r"\\")

        item_type = _infer_item_type(row)
        if not item_type:
            continue
        it = item_type.lower()

        if it == "table":
            data_attr = _clean_str(row.get("Data Attr", ""))
            df_obj = _get_df_attr(self, data_attr)
            if isinstance(df_obj, pd.DataFrame) and not df_obj.empty:
                dfp_raw = _apply_head(df_obj, row.get("Head", None), top_n)
                # Filter columns for tex output
                columns_spec = _clean_str(row.get("Columns (docx/pptx/tex)", ""))
                dfp_filtered = _filter_columns(dfp_raw, columns_spec)
                dfp = _df_display(dfp_filtered)
                if dfp is not None and not dfp.empty:
                    if modern_style:
                        # Generate modern styled table
                        n_cols = len(dfp.columns)
                        col_spec = "l" + "r" * (n_cols - 1)
                        
                        lines.append(r"\begin{table}[htbp]")
                        lines.append(r"\centering")
                        lines.append(r"\small")
                        lines.append(r"\begin{tabular}{" + col_spec + r"}")
                        lines.append(r"\toprule")
                        
                        # Header row with colored background
                        lines.append(r"\rowcolor{tableheader}")
                        header_cells = []
                        for col in dfp.columns:
                            col_escaped = _tex_escape(str(col))
                            header_cells.append(r"\textcolor{white}{\textbf{" + col_escaped + r"}}")
                        lines.append(" & ".join(header_cells) + r" \\")
                        lines.append(r"\midrule")
                        
                        # Data rows with alternating colors
                        for i, (_, row_data) in enumerate(dfp.iterrows()):
                            if i % 2 == 1:
                                lines.append(r"\rowcolor{tablealt}")
                            row_cells = [_tex_escape(str(val)) for val in row_data.values]
                            lines.append(" & ".join(row_cells) + r" \\")
                        
                        lines.append(r"\bottomrule")
                        lines.append(r"\end{tabular}")
                        
                        caption = _clean_str(row.get("Caption", ""))
                        if caption:
                            prefix = f"Table {table_counter}: " if enumerate_tables else ""
                            caption_escaped = _tex_escape(prefix + caption)
                            lines.append(r"\caption{" + caption_escaped + r"}")
                            table_counter += 1
                        
                        lines.append(r"\end{table}")
                        lines.append("")
                    else:
                        # Simple table
                        latex_table = dfp.to_latex(index=False, escape=True)
                        lines.append(latex_table)
                        caption = _clean_str(row.get("Caption", ""))
                        if caption:
                            prefix = (
                                f"Table {table_counter}: "
                                if enumerate_tables
                                else ""
                            )
                            lines.append(rf"\textbf{{{prefix}{caption}}}\\")
                            table_counter += 1

        elif it == "plot":
            tpath = _clean_str(row.get("Path", ""))
            filename = _clean_str(row.get("Plot Filename", ""))
            if filename:
                img = _resolve_image_path(self, tpath, filename, results_base)
                if img is not None:
                    tex_img_path = str(img).replace("\\", "/")
                    lines.append(
                        rf"\begin{{center}}\includegraphics[width=0.9\textwidth]{{{tex_img_path}}}\end{{center}}"
                    )
                    caption = _clean_str(row.get("Caption", ""))
                    if caption:
                        prefix = (
                            f"Figure {figure_counter}: "
                            if enumerate_figures
                            else ""
                        )
                        lines.append(rf"\textit{{{prefix}{caption}}}\\")
                        figure_counter += 1

        elif it == "text":
            text_value = _clean_str(row.get("Text", ""))
            if text_value:
                for ln in text_value.splitlines():
                    lines.append(ln + r"\\")

        elif it == "description":
            data_attr = _clean_str(row.get("Data Attr", ""))
            if data_attr and hasattr(self, data_attr):
                desc_text2 = getattr(self, data_attr)
                if isinstance(desc_text2, str) and desc_text2:
                    for ln in desc_text2.splitlines():
                        lines.append(ln + r"\\")

    lines.append(r"\end{document}")

    out_path = _resolve_path(self, output_path, results_base)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    return out_path


# =============================================================================
# REPORT DATA AVAILABILITY CHECK
# =============================================================================

def check_report_data_availability(
    self: Any,
    template_path: Union[str, Path] = FD / "additional files" / "template for report output.xlsx",
    template_sheet: str = "all",
    results_base: Union[str, Path] = "results",
    verbose: bool = True,
) -> Dict[str, Any]:
    """
    Check which template items have data available and which are missing.
    
    This function helps diagnose why reports might have limited content by
    identifying which analyses have been run and which data attributes exist.
    
    Parameters
    ----------
    self : Any
        BiblioStats or BiblioAnalysis object to check.
    template_path : str or Path
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    results_base : str or Path, default "results"
        Base folder for relative paths.
    verbose : bool, default True
        Print detailed information about available/missing data.
        
    Returns
    -------
    dict
        Dictionary with 'available', 'missing_data', 'missing_plots', 
        'summary' keys.
    
    Examples
    --------
    >>> ba = BiblioAnalysis("data.csv", db="scopus")
    >>> ba.count_sources()  # Run some analysis
    >>> report_status = check_report_data_availability(ba)
    >>> print(report_status['summary'])
    """
    template_path = _resolve_path(self, template_path, results_base)
    df_template = _read_template(template_path, template_sheet)
    
    available = []
    missing_data = []
    missing_plots = []
    
    for _, row in df_template.iterrows():
        level1 = _clean_str(row.get("Level 1", ""))
        level2 = _clean_str(row.get("Level 2", ""))
        item_type = _infer_item_type(row)
        data_attr = _clean_str(row.get("Data Attr", ""))
        plot_filename = _clean_str(row.get("Plot Filename", ""))
        
        item_name = f"{level1} > {level2}" if level2 else level1
        
        if item_type == "table":
            if data_attr:
                df_obj = _get_df_attr(self, data_attr)
                if df_obj is not None and not df_obj.empty:
                    available.append({
                        "name": item_name,
                        "type": "table",
                        "attr": data_attr,
                        "rows": len(df_obj),
                    })
                else:
                    missing_data.append({
                        "name": item_name,
                        "type": "table",
                        "attr": data_attr,
                        "reason": "attribute not found" if df_obj is None else "empty DataFrame",
                    })
        
        elif item_type == "plot":
            if plot_filename:
                tpath = _clean_str(row.get("Path", ""))
                img = _resolve_image_path(self, tpath, plot_filename, results_base)
                if img is not None:
                    available.append({
                        "name": item_name,
                        "type": "plot",
                        "filename": plot_filename,
                    })
                else:
                    missing_plots.append({
                        "name": item_name,
                        "type": "plot",
                        "filename": plot_filename,
                        "reason": "plot file not found",
                    })
    
    summary = {
        "total_items": len(df_template),
        "available_items": len(available),
        "missing_tables": len(missing_data),
        "missing_plots": len(missing_plots),
        "coverage_pct": round(100 * len(available) / max(1, len(df_template)), 1),
    }
    
    result = {
        "available": available,
        "missing_data": missing_data,
        "missing_plots": missing_plots,
        "summary": summary,
    }
    
    if verbose:
        print("=" * 60)
        print("REPORT DATA AVAILABILITY CHECK")
        print("=" * 60)
        print(f"\nTemplate: {template_path.name} (sheet: {template_sheet})")
        print(f"Total items in template: {summary['total_items']}")
        print(f"Available items: {summary['available_items']} ({summary['coverage_pct']}%)")
        print(f"Missing tables: {summary['missing_tables']}")
        print(f"Missing plots: {summary['missing_plots']}")
        
        if missing_data:
            print("\n" + "-" * 40)
            print("MISSING TABLE DATA:")
            print("-" * 40)
            for item in missing_data:
                print(f"  • {item['name']}")
                print(f"    Attribute: {item['attr']}")
                print(f"    Reason: {item['reason']}")
        
        if missing_plots:
            print("\n" + "-" * 40)
            print("MISSING PLOT FILES:")
            print("-" * 40)
            for item in missing_plots:
                print(f"  • {item['name']}")
                print(f"    File: {item['filename']}")
        
        if available and verbose:
            print("\n" + "-" * 40)
            print("AVAILABLE ITEMS:")
            print("-" * 40)
            for item in available[:20]:  # Show first 20
                if item['type'] == 'table':
                    print(f"  ✓ {item['name']} ({item['rows']} rows)")
                else:
                    print(f"  ✓ {item['name']} (plot)")
            if len(available) > 20:
                print(f"  ... and {len(available) - 20} more")
        
        print("\n" + "=" * 60)
        
        # Suggest what to run
        if missing_data:
            print("\nHINT: To generate missing data, run the corresponding analysis methods:")
            method_hints = {
                "descriptives_df": "count_all() or individual count methods",
                "sources_counts_df": "count_sources()",
                "authors_counts_df": "count_authors()",
                "author_keywords_counts_df": "count_author_keywords()",
                "references_counts_df": "count_references()",
                "affiliations_counts_df": "count_affiliations()",
                "production_df": "get_scientific_production()",
                "lotka_df": "compute_lotka() or plot_lotka()",
                "bradford_df": "compute_bradford() or plot_bradford()",
                "zipf_df": "compute_zipf() or plot_zipf()",
            }
            suggested = set()
            for item in missing_data:
                attr = item['attr']
                for key, hint in method_hints.items():
                    if key in attr:
                        suggested.add(f"  - {hint}")
            if suggested:
                for hint in list(suggested)[:5]:
                    print(hint)
        
        print("")
    
    return result


def debug_report_generation(
    self: Any,
    template_path: Union[str, Path] = FD / "additional files" / "template for report output.xlsx",
    template_sheet: str = "all",
    results_base: Union[str, Path] = "results",
) -> None:
    """
    Debug helper to diagnose why reports might be empty.
    
    Prints detailed information about:
    - Template parsing
    - Available data attributes
    - Found/missing plot files
    """
    print("=" * 60)
    print("REPORT GENERATION DEBUG")
    print("=" * 60)
    
    # Resolve template
    template_path = _resolve_path(self, template_path, results_base)
    print(f"\n1. Template path: {template_path}")
    print(f"   Exists: {template_path.exists()}")
    
    # Read template
    try:
        df_template = _read_template(template_path, template_sheet)
        print(f"   Sheet: {template_sheet}")
        print(f"   Rows in template: {len(df_template)}")
    except Exception as e:
        print(f"   ERROR reading template: {e}")
        return
    
    # Check res_folder
    res_folder = getattr(self, "res_folder", None)
    print(f"\n2. self.res_folder: {res_folder}")
    print(f"   results_base: {results_base}")
    
    # Check plots folder
    plots_folder = Path(res_folder or results_base) / "plots"
    print(f"\n3. Plots folder: {plots_folder}")
    print(f"   Exists: {plots_folder.exists()}")
    if plots_folder.exists():
        plot_files = list(plots_folder.glob("*.png"))
        print(f"   PNG files found: {len(plot_files)}")
        for pf in plot_files[:10]:
            print(f"     - {pf.name}")
        if len(plot_files) > 10:
            print(f"     ... and {len(plot_files) - 10} more")
    
    # Process template rows
    print(f"\n4. Template items analysis:")
    tables_found = 0
    tables_missing = 0
    plots_found = 0
    plots_missing = 0
    
    for idx, row in df_template.iterrows():
        item_type = _infer_item_type(row)
        level1 = _clean_str(row.get("Level 1", ""))
        level2 = _clean_str(row.get("Level 2", ""))
        item_name = f"{level1} > {level2}" if level2 else level1
        
        if item_type == "table":
            data_attr = _clean_str(row.get("Data Attr", ""))
            df_obj = _get_df_attr(self, data_attr)
            if df_obj is not None and not df_obj.empty:
                tables_found += 1
            else:
                tables_missing += 1
                has_attr = hasattr(self, data_attr) if data_attr else False
                print(f"   ✗ TABLE missing: {item_name}")
                print(f"      attr='{data_attr}', hasattr={has_attr}")
                
        elif item_type == "plot":
            plot_filename = _clean_str(row.get("Plot Filename", ""))
            tpath = _clean_str(row.get("Path", ""))
            img_path = _resolve_image_path(self, tpath, plot_filename, results_base)
            if img_path is not None:
                plots_found += 1
            else:
                plots_missing += 1
                print(f"   ✗ PLOT missing: {item_name}")
                print(f"      filename='{plot_filename}'")
                # Show where we looked
                search_paths = []
                if res_folder:
                    search_paths.append(Path(res_folder) / plot_filename)
                    search_paths.append(Path(res_folder) / "plots" / plot_filename)
                search_paths.append(Path(results_base) / plot_filename)
                search_paths.append(Path(results_base) / "plots" / plot_filename)
                print(f"      Searched:")
                for sp in search_paths:
                    exists = sp.exists()
                    print(f"        {sp} {'✓' if exists else '✗'}")
    
    print(f"\n5. Summary:")
    print(f"   Tables found: {tables_found}")
    print(f"   Tables missing: {tables_missing}")
    print(f"   Plots found: {plots_found}")
    print(f"   Plots missing: {plots_missing}")
    print(f"   Total content items: {tables_found + plots_found}")
    
    if tables_found + plots_found == 0:
        print("\n   ⚠️  NO CONTENT FOUND - Report will be empty!")
        print("   Possible causes:")
        print("   - No analysis methods have been run (call get_main_info(), count_sources(), etc.)")
        print("   - res_folder path is incorrect")
        print("   - Template sheet name is wrong")
    
    print("=" * 60)


def generate_plots_from_template(
    self: Any,
    template_path: Union[str, Path] = FD / "additional files" / "template for report output.xlsx",
    template_sheet: str = "all",
    output_folder: Optional[Union[str, Path]] = None,
    skip_existing: bool = True,
    verbose: bool = True,
) -> Dict[str, Any]:
    """
    Generate all plots specified in the report template.
    
    This function reads the template, extracts Plot Method entries,
    and calls the corresponding plotting methods on the BiblioAnalysis object.
    
    Parameters
    ----------
    self : Any
        BiblioStats or BiblioAnalysis object.
    template_path : str or Path
        Path to the Excel template.
    template_sheet : str, default "all"
        Name of the sheet in the template workbook.
    output_folder : str or Path, optional
        Folder to save plots. Defaults to self.res_folder/plots.
    skip_existing : bool, default True
        Skip plots that already exist.
    verbose : bool, default True
        Print progress information.
        
    Returns
    -------
    dict
        Summary with 'generated', 'skipped', 'failed' lists.
        
    Examples
    --------
    >>> ba = BiblioAnalysis("data.csv", db="scopus")
    >>> ba.get_main_info()  # Prepare data
    >>> generate_plots_from_template(ba)  # Generate all plots
    >>> ba.generate_report()  # Now create report with plots
    """
    # Resolve paths
    template_path = Path(template_path)
    if not template_path.is_absolute():
        template_path = FD / template_path
    
    if output_folder is None:
        base = getattr(self, "res_folder", None) or "results"
        output_folder = Path(base) / "plots"
    else:
        output_folder = Path(output_folder)
    
    output_folder.mkdir(parents=True, exist_ok=True)
    
    # Read template
    df_template = _read_template(template_path, template_sheet)
    
    # Map of plot method names to actual methods
    PLOT_METHOD_MAP = {
        # Production plots
        "production_line": "plot_production_line",
        "production_area": "plot_production_area",
        "production_bar": "plot_production_bar",
        # Item bar charts
        "sources_bar": "plot_sources_bar",
        "authors_bar": "plot_authors_bar",
        "keywords_bar": "plot_keywords_bar",
        "countries_bar": "plot_countries_bar",
        "affiliations_bar": "plot_affiliations_bar",
        "references_bar": "plot_references_bar",
        # Treemaps
        "sources_treemap": "plot_sources_treemap",
        "authors_treemap": "plot_authors_treemap",
        "keywords_treemap": "plot_keywords_treemap",
        "countries_treemap": "plot_countries_treemap",
        # Lollipop charts
        "sources_lollipop": "plot_sources_lollipop",
        "authors_lollipop": "plot_authors_lollipop",
        "keywords_lollipop": "plot_keywords_lollipop",
        # Donuts and pies
        "doctypes_donut": "plot_document_types_donut",
        "doctypes_pie": "plot_document_types_pie",
        "doctypes_bar": "plot_document_types_bar",
        # Wordclouds
        "keywords_wordcloud": "plot_keywords_wordcloud",
        "title_wordcloud": "plot_title_wordcloud",
        "abstract_wordcloud": "plot_abstract_wordcloud",
        # Histograms
        "citations_histogram": "plot_citations_histogram",
        "citations_distribution": "plot_citations_histogram",
        # Maps
        "countries_map": "plot_countries_map",
        "world_map": "plot_world_map",
        # Networks
        "keyword_network": "plot_keyword_coocurrence_network",
        "author_network": "plot_coauthorship_network",
        "citation_network": "plot_cocitation_network",
        # Laws
        "bradford_law": "plot_bradford_law",
        "lotka_law": "plot_lotka_law",
        "zipf_law": "plot_zipf_law",
        # Time series
        "keywords_evolution": "plot_keywords_evolution",
        "citations_evolution": "plot_citations_evolution",
    }
    
    generated = []
    skipped = []
    failed = []
    
    # Reverse map: filename pattern -> method
    FILENAME_TO_METHOD = {
        "scientific_production": "plot_production_line",
        "production_line": "plot_production_line",
        "production_area": "plot_production_area",
        "production_bar": "plot_production_bar",
        "sources_bar": "plot_sources_bar",
        "sources_treemap": "plot_sources_treemap",
        "sources_lollipop": "plot_sources_lollipop",
        "sources_bubble": "plot_sources_bubble",
        "authors_bar": "plot_authors_bar",
        "authors_treemap": "plot_authors_treemap",
        "authors_lollipop": "plot_authors_lollipop",
        "authors_bubble": "plot_authors_bubble",
        "keywords_bar": "plot_keywords_bar",
        "keywords_wordcloud": "plot_keywords_wordcloud",
        "keywords_treemap": "plot_keywords_treemap",
        "doctypes_donut": "plot_document_types_donut",
        "doctypes_bar": "plot_document_types_bar",
        "doctypes_pie": "plot_document_types_pie",
        "countries_bar": "plot_countries_bar",
        "countries_treemap": "plot_countries_treemap",
        "countries_map": "plot_countries_map",
        "citations_histogram": "plot_citations_histogram",
        "citations_boxplot": "plot_citations_boxplot",
        "keyword_network": "plot_keyword_coocurrence_network",
        "keyword_cooccurrence_heatmap": "plot_keyword_cooccurrence_heatmap",
        "author_collaboration_network": "plot_coauthorship_network",
        "country_collaboration_network": "plot_country_collaboration_network",
        "bradford_zones": "plot_bradford_law",
        "lotka_distribution": "plot_lotka_law",
        "zipf_distribution": "plot_zipf_law",
        "thematic_map": "plot_thematic_map",
        "reference_spectrogram": "plot_reference_spectrogram",
    }
    
    for _, row in df_template.iterrows():
        item_type = _infer_item_type(row)
        if item_type != "plot":
            continue
            
        plot_method_name = _clean_str(row.get("Plot Method", ""))
        plot_filename = _clean_str(row.get("Plot Filename", ""))
        level2 = _clean_str(row.get("Level 2", ""))
        
        if not plot_filename:
            continue
        
        # If Plot Method is missing, try to infer from filename
        if not plot_method_name:
            # Extract base name without extension
            base_name = plot_filename.rsplit(".", 1)[0] if "." in plot_filename else plot_filename
            # Look for matching pattern
            inferred_method = FILENAME_TO_METHOD.get(base_name)
            if inferred_method:
                plot_method_name = inferred_method
                if verbose:
                    print(f"  ℹ Inferred method '{plot_method_name}' from filename '{plot_filename}'")
            else:
                if verbose:
                    print(f"  ⚠ No Plot Method for '{plot_filename}' and couldn't infer from filename")
                failed.append({
                    "method": "unknown",
                    "filename": plot_filename,
                    "reason": "no Plot Method specified and couldn't infer from filename"
                })
                continue
        
        output_path = output_folder / plot_filename
        
        # Skip if exists
        if skip_existing and output_path.exists():
            if verbose:
                print(f"  ⏭ Skipping (exists): {plot_filename}")
            skipped.append({"method": plot_method_name, "filename": plot_filename})
            continue
        
        # Resolve actual method name
        actual_method = PLOT_METHOD_MAP.get(plot_method_name, plot_method_name)
        
        # Try to find and call the method
        method = getattr(self, actual_method, None)
        if method is None:
            # Try with plot_ prefix
            method = getattr(self, f"plot_{plot_method_name}", None)
        if method is None:
            # Try without plot_ prefix
            method = getattr(self, plot_method_name.replace("plot_", ""), None)
        
        if method is None or not callable(method):
            if verbose:
                print(f"  ✗ Method not found: {plot_method_name} ({actual_method})")
            failed.append({
                "method": plot_method_name, 
                "filename": plot_filename,
                "reason": "method not found"
            })
            continue
        
        # Call the method
        try:
            if verbose:
                print(f"  → Generating: {plot_filename}...", end=" ", flush=True)
            
            # Try to call with filename parameter
            try:
                method(filename=str(output_path))
            except TypeError:
                # Method doesn't accept filename, call and save manually
                import matplotlib.pyplot as plt
                fig = method()
                if fig is not None and hasattr(fig, 'savefig'):
                    fig.savefig(output_path, dpi=300, bbox_inches='tight')
                    plt.close(fig)
                elif plt.get_fignums():
                    plt.savefig(output_path, dpi=300, bbox_inches='tight')
                    plt.close()
            
            if verbose:
                print("✓")
            generated.append({"method": plot_method_name, "filename": plot_filename})
            
        except Exception as e:
            if verbose:
                print(f"✗ ({type(e).__name__}: {str(e)[:50]})")
            failed.append({
                "method": plot_method_name,
                "filename": plot_filename,
                "reason": str(e)
            })
    
    summary = {
        "generated": generated,
        "skipped": skipped,
        "failed": failed,
        "total_generated": len(generated),
        "total_skipped": len(skipped),
        "total_failed": len(failed),
    }
    
    if verbose:
        print(f"\nPlot generation summary:")
        print(f"  Generated: {len(generated)}")
        print(f"  Skipped: {len(skipped)}")
        print(f"  Failed: {len(failed)}")
    
    return summary


# =============================================================================
# GROUP ANALYSIS REPORTS
# =============================================================================

GROUP_TEMPLATE_PATH: Path = FD / "additional files" / "template for group report output.xlsx"


def save_group_excel_report(
    self: Any,
    output_path: str = "group_report.xlsx",
    template_path: Union[str, Path] = GROUP_TEMPLATE_PATH,
    template_sheet: str = "all",
    modern_style: bool = True,
    include_cover_sheet: bool = True,
    report_title: str = "Group Analysis Report",
    results_base: Union[str, Path] = "results",
    **kwargs
) -> Path:
    """
    Generate Excel report for group analysis.
    
    Parameters
    ----------
    self : BiblioGroupAnalysis
        Group analysis object.
    output_path : str
        Output file path.
    template_path : str or Path
        Path to template Excel file.
    template_sheet : str
        Sheet name in template.
    modern_style : bool
        Apply modern styling.
    include_cover_sheet : bool
        Include cover sheet.
    report_title : str
        Report title.
    results_base : str or Path
        Base folder for output.
    
    Returns
    -------
    Path
        Path to saved file.
    """
    # Use the main Excel report function with group template
    return save_excel_report_from_template(
        self,
        output_path=output_path,
        template_path=template_path,
        template_sheet=template_sheet,
        modern_style=modern_style,
        include_cover_sheet=include_cover_sheet,
        report_title=report_title,
        results_base=results_base,
        **kwargs
    )


def save_group_word_report(
    self: Any,
    output_path: str = "group_report.docx",
    template_path: Union[str, Path] = GROUP_TEMPLATE_PATH,
    template_sheet: str = "all",
    modern_style: bool = True,
    include_cover_page: bool = True,
    include_toc: bool = True,
    report_title: str = "Group Analysis Report",
    results_base: Union[str, Path] = "results",
    **kwargs
) -> Path:
    """
    Generate Word report for group analysis.
    
    Parameters
    ----------
    self : BiblioGroupAnalysis
        Group analysis object.
    output_path : str
        Output file path.
    template_path : str or Path
        Path to template Excel file.
    template_sheet : str
        Sheet name in template.
    modern_style : bool
        Apply modern styling.
    include_cover_page : bool
        Include cover page.
    include_toc : bool
        Include table of contents.
    report_title : str
        Report title.
    results_base : str or Path
        Base folder for output.
    
    Returns
    -------
    Path
        Path to saved file.
    """
    return save_word_report_from_template(
        self,
        output_path=output_path,
        template_path=template_path,
        template_sheet=template_sheet,
        modern_style=modern_style,
        include_cover_page=include_cover_page,
        include_toc=include_toc,
        report_title=report_title,
        results_base=results_base,
        **kwargs
    )


def save_group_pptx_report(
    self: Any,
    output_path: str = "group_report.pptx",
    template_path: Union[str, Path] = GROUP_TEMPLATE_PATH,
    template_sheet: str = "all",
    modern_style: bool = True,
    presentation_title: str = "Group Analysis Report",
    results_base: Union[str, Path] = "results",
    **kwargs
) -> Path:
    """
    Generate PowerPoint report for group analysis.
    
    Parameters
    ----------
    self : BiblioGroupAnalysis
        Group analysis object.
    output_path : str
        Output file path.
    template_path : str or Path
        Path to template Excel file.
    template_sheet : str
        Sheet name in template.
    modern_style : bool
        Apply modern styling.
    presentation_title : str
        Presentation title.
    results_base : str or Path
        Base folder for output.
    
    Returns
    -------
    Path
        Path to saved file.
    """
    return save_powerpoint_report_from_template(
        self,
        output_path=output_path,
        template_path=template_path,
        template_sheet=template_sheet,
        modern_style=modern_style,
        presentation_title=presentation_title,
        results_base=results_base,
        **kwargs
    )


def save_group_tex_report(
    self: Any,
    output_path: str = "group_report.tex",
    template_path: Union[str, Path] = GROUP_TEMPLATE_PATH,
    template_sheet: str = "all",
    modern_style: bool = True,
    report_title: str = "Group Analysis Report",
    results_base: Union[str, Path] = "results",
    **kwargs
) -> Path:
    """
    Generate LaTeX report for group analysis.
    
    Parameters
    ----------
    self : BiblioGroupAnalysis
        Group analysis object.
    output_path : str
        Output file path.
    template_path : str or Path
        Path to template Excel file.
    template_sheet : str
        Sheet name in template.
    modern_style : bool
        Apply modern styling.
    report_title : str
        Report title.
    results_base : str or Path
        Base folder for output.
    
    Returns
    -------
    Path
        Path to saved file.
    """
    return save_tex_report_from_template(
        self,
        output_path=output_path,
        template_path=template_path,
        template_sheet=template_sheet,
        modern_style=modern_style,
        report_title=report_title,
        results_base=results_base,
        **kwargs
    )


__all__ = [
    "save_excel_report_from_template",
    "save_powerpoint_report_from_template",
    "save_word_report_from_template",
    "save_tex_report_from_template",
    # Group reports
    "save_group_excel_report",
    "save_group_word_report",
    "save_group_pptx_report",
    "save_group_tex_report",
    # Utilities
    "check_report_data_availability",
    "validate_template",
]
