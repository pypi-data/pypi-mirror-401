"""PowerPoint document reader for AgenticX Knowledge Management System"""

import asyncio
import logging
from pathlib import Path
from typing import List, Optional, Union, Any

from ..base import BaseReader
from ..document import Document, DocumentMetadata

logger = logging.getLogger(__name__)


class PowerPointReader(BaseReader):
    """Reader for PowerPoint documents (.ppt, .pptx)"""
    
    def __init__(
        self,
        extract_images: bool = False,
        extract_notes: bool = True,
        include_slide_numbers: bool = True,
        **kwargs
    ):
        """
        Initialize PowerPointReader
        
        Args:
            extract_images: Whether to extract images from slides
            extract_notes: Whether to extract speaker notes
            include_slide_numbers: Whether to include slide numbers in output
            **kwargs: Additional configuration
        """
        self.extract_images = extract_images
        self.extract_notes = extract_notes
        self.include_slide_numbers = include_slide_numbers
        self.config = kwargs
        
        # Try to import required libraries
        self._ppt_library = self._get_ppt_library()
    
    def _get_ppt_library(self) -> Optional[str]:
        """Detect available PowerPoint processing library"""
        try:
            from pptx import Presentation
            return 'python-pptx'
        except ImportError:
            pass
        
        return None
    
    async def read(self, source: Union[str, Path]) -> List[Document]:
        """Read PowerPoint document and return documents
        
        Args:
            source: File path to read
            
        Returns:
            List containing single document with extracted text
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ImportError: If no PowerPoint library is available
            Exception: If PowerPoint reading fails
        """
        
        if not self._ppt_library:
            raise ImportError(
                "No PowerPoint library available. Install: "
                "python-pptx (pip install python-pptx)"
            )
        
        file_path = Path(source)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not file_path.is_file():
            raise ValueError(f"Path is not a file: {file_path}")
        
        # Extract text based on available library
        if self._ppt_library == 'python-pptx':
            content, metadata_dict = await self._read_with_python_pptx(file_path)
        else:
            raise ImportError(f"Unsupported PowerPoint library: {self._ppt_library}")
        
        # Create document metadata
        metadata = DocumentMetadata(
            name=file_path.name,
            source=str(file_path),
            source_type="file",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" if file_path.suffix.lower() == '.pptx' else "application/vnd.ms-powerpoint",
            size=len(content),
            reader_name=self.__class__.__name__,
            custom=metadata_dict
        )
        
        # Create document
        document = Document(
            content=content,
            metadata=metadata
        )
        
        return [document]
    
    async def read_async(self, source: Union[str, Path], **kwargs):
        """Asynchronously read PowerPoint file and yield documents
        
        Args:
            source: File path to read
            **kwargs: Additional arguments
            
        Yields:
            Documents one by one
        """
        documents = await self.read(source)
        for document in documents:
            yield document
    
    async def _read_with_python_pptx(self, file_path: Path) -> tuple[str, dict]:
        """Read PowerPoint document using python-pptx"""
        
        def _extract():
            from pptx import Presentation
            
            # Only works with .pptx files
            if file_path.suffix.lower() != '.pptx':
                raise ValueError("python-pptx only supports .pptx files")
            
            prs = Presentation(str(file_path))
            
            text_parts = []
            slide_count = 0
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_text = []
                
                if self.include_slide_numbers:
                    slide_text.append(f"=== 幻灯片 {slide_idx} ===")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract speaker notes if requested
                if self.extract_notes and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"备注: {notes_text}")
                
                if slide_text:
                    text_parts.append('\n'.join(slide_text))
            
            # Get presentation properties
            core_props = prs.core_properties
            metadata = {
                'ppt_reader_library': 'python-pptx',
                'file_extension': file_path.suffix.lower(),
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'slides_count': slide_count,
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
            
            return '\n\n'.join(text_parts), metadata
        
        # Run in thread pool
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _extract)
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions
        
        Returns:
            List of supported extensions
        """
        
        return ['.ppt', '.pptx']
    
    def __str__(self) -> str:
        return f"PowerPointReader(library={self._ppt_library}, extract_notes={self.extract_notes})"