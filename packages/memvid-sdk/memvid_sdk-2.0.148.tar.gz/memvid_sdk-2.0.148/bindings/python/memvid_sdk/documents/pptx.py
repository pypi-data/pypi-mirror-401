"""PowerPoint Parser using python-pptx."""

from __future__ import annotations

import os
from typing import Optional

from . import ParseOptions, ParseResult, DocumentItem


def parse_pptx(file_path: str, options: Optional[ParseOptions] = None) -> ParseResult:
    """
    Parse a PowerPoint file, extracting text per slide.

    Args:
        file_path: Path to the PowerPoint file
        options: Parsing options (max_items limits slides)

    Returns:
        ParseResult with per-slide items
    """
    filename = os.path.basename(file_path)

    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for PowerPoint parsing. "
            "Install with: pip install python-pptx"
        )

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        max_items = (options or {}).get("max_items") or total_slides

        items: list[DocumentItem] = []
        for i, slide in enumerate(prs.slides):
            if i >= max_items:
                break

            text_parts: list[str] = []
            title: Optional[str] = None

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                    # First text shape with content is likely the title
                    if title is None:
                        title = shape.text.strip().split("\n")[0]

            if text_parts:
                items.append({
                    "number": i + 1,
                    "title": title,
                    "text": "\n".join(text_parts),
                })

        return {
            "type": "pptx",
            "filename": filename,
            "total_items": total_slides,
            "items": items,
        }
    except Exception as e:
        raise RuntimeError(
            f'Failed to parse PowerPoint file "{filename}": {e}. '
            f"Ensure the file is a valid .pptx/.ppt file."
        )
