# 使用 Spire.Presentation 合并PPT文件，保留原始设计
# 安装: pip install spire.presentation

import os
from spire.presentation import *
from spire.presentation.common import *
from pptx import Presentation as PptxPresentation

def combine_ppt_files_with_spire(source_folder, output_file, png_names=None):
    """
    使用 Spire.Presentation 合并PPT文件，每个PPT只保留第一页，并保留原始设计
    
    Args:
        source_folder: 源PPT文件所在的文件夹路径
        output_file: 输出的合并PPT文件路径
    """
    # 获取所有pptx文件并按字典序排序
    
    ppt_files = sorted([f for f in os.listdir(source_folder) if f.endswith('.pptx')])

    ppt_names = [png_name.replace(".png", ".pptx") for png_name in png_names] if png_names else None
    
    if not ppt_files:
        print("未找到任何PPT文件")
        return
    
    print(f"找到 {len(ppt_files)} 个PPT文件:")
    valid_ppt_files = []
    for idx, file in enumerate(ppt_files, 1):
        print(f"  {idx}. {file}")
        if ppt_names and file in ppt_names:
            valid_ppt_files.append(file)
    if ppt_names:
        ppt_files = valid_ppt_files
        print(f"\n根据提供的PNG名称过滤后，剩余 {len(ppt_files)} 个PPT文件:")
    
    # 创建主演示文稿对象，使用第一个PPT作为基础
    first_ppt_path = os.path.join(source_folder, ppt_files[0])
    main_pres = Presentation()
    main_pres.LoadFromFile(first_ppt_path)
    
    # 删除第一个PPT的多余页面，只保留第一页
    while main_pres.Slides.Count > 1:
        main_pres.Slides.RemoveAt(1)
    
    print(f"  已添加: {ppt_files[0]} (第1页)")
    
    # 遍历剩余的PPT文件
    for ppt_file in ppt_files[1:]:
        file_path = os.path.join(source_folder, ppt_file)
        
        # 加载当前PPT文件
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        
        if temp_pres.Slides.Count > 0:
            # 使用 AppendBySlide 方法追加第一页，保留原始设计
            main_pres.Slides.AppendBySlide(temp_pres.Slides[0])
            print(f"  已添加: {ppt_file} (第1页)")
        else:
            print(f"  跳过: {ppt_file} (无幻灯片)")
        
        # 释放临时演示文稿资源
        temp_pres.Dispose()
    
    # 保存合并后的PPT
    main_pres.SaveToFile(output_file, FileFormat.Pptx2016)
    print(f"\n合并完成！输出文件: {output_file}")
    print(f"总共合并了 {main_pres.Slides.Count} 页幻灯片")
    
    # 释放资源
    main_pres.Dispose()

    valid_png_names = [ppt_name.replace(".pptx", ".png") for ppt_name in ppt_files] # 返回真正存在的png names

    return valid_png_names



def clean_ppt(in_ppt_file, out_ppt_file):
    """
    使用 python-pptx 删除PPT中的多余形状
    
    Args:
        ppt_file: 需要清理的PPT文件路径
    """
    ppt = PptxPresentation(in_ppt_file)
    for slide in ppt.slides:
        for shape in list(slide.shapes):
            # 如果shape name 叫做"New shape",删除它
            if shape.name == "New shape":
                sp = slide.shapes._spTree.remove(shape._element)

    ppt.save(out_ppt_file)
    

def combine_ppt(source_folder, out_ppt_file, png_names = None):
    # 确保是字符串路径，因为后面用到了 .replace
    source_folder = str(source_folder)
    out_ppt_file = str(out_ppt_file)
    
    # 方法1: 保留原始设计（推荐）
    output_file1 = out_ppt_file.replace(".pptx", "_combined_original_design.pptx")
    valid_png_names = combine_ppt_files_with_spire(source_folder, output_file1, png_names=png_names)

    clean_ppt(output_file1, out_ppt_file)
    print(f"\n已生成合并的PPT文件: {out_ppt_file}")
    os.remove(output_file1)
    
    return valid_png_names
